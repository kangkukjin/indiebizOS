"""doc_formats.py — 문서 IR(blocks) → 출력 포맷 (handler.py 에서 분리, 2026-08-06 1500줄 규칙).

한 방향 변환기만 산다: blocks → html/css/typst/docx/pptx/markdown + 이미지 해소.
입력을 blocks 로 *만드는* 쪽(마크다운 파서·structure/render 오케스트레이터)은
doc_build.py. 무거운 의존성(docx·pptx·typst)은 여기서도 함수 안 지연 import —
모듈레벨 stdlib 만(폰 import-safe 불변식).
"""
import json
import re
import os

#   · table{columns,rows}(=데이터 통화 재사용) · quote{text,cite?} · code{text,lang?} · divider
# 포맷 중립 IR을 여러 emitter가 렌더(현재 html). slide/newspaper와 달리 단일 IR이 단일 진실 소스.
def _doc_blocks_to_html(blocks: list) -> str:
    import html as _html

    def esc(s):
        return _html.escape(str(s if s is not None else ""))

    parts = []
    for b in blocks:
        if not isinstance(b, dict):
            continue
        t = (b.get("type") or "paragraph").lower()
        if t == "heading":
            lvl = max(1, min(6, int(b.get("level") or 2)))
            txt = b.get("text")
            anchor = b.get("anchor") or (str(txt) if txt else "")  # 목차 점프용 id
            id_attr = f' id="{esc(anchor)}"' if anchor else ""
            parts.append(f"<h{lvl}{id_attr}>{esc(txt)}</h{lvl}>")
        elif t == "list":
            tag = "ol" if b.get("ordered") else "ul"
            li = []
            for i in (b.get("items") or []):
                # 항목은 문자열 또는 {text, url}(링크 — 목차·북마크 등)
                if isinstance(i, dict):
                    # note = 링크 뒤에 붙는 설명(평문). 링크는 제목에만 건다.
                    note = f' <span class="li-note">{esc(i.get("note"))}</span>' if i.get("note") else ""
                    if i.get("url"):
                        li.append(f'<li><a href="{esc(i.get("url"))}">{esc(i.get("text"))}</a>{note}</li>')
                    else:
                        li.append(f"<li>{esc(i.get('text'))}{note}</li>")
                else:
                    li.append(f"<li>{esc(i)}</li>")
            parts.append(f"<{tag}>{''.join(li)}</{tag}>")
        elif t == "image":
            src = b.get("src") or b.get("path") or ""
            cap = b.get("caption")
            cap_html = f"<figcaption>{esc(cap)}</figcaption>" if cap else ""
            parts.append(f'<figure><img src="{esc(src)}" alt="{esc(cap)}">{cap_html}</figure>')
        elif t == "table":
            # 데이터 통화 재사용: {columns, rows}
            cols = b.get("columns") or []
            rows = b.get("rows") or []
            thead = "".join(f"<th>{esc(c)}</th>" for c in cols)
            tbody = "".join("<tr>" + "".join(f"<td>{esc(c)}</td>" for c in r) + "</tr>"
                            for r in rows if isinstance(r, (list, tuple)))
            parts.append(f"<table><thead><tr>{thead}</tr></thead><tbody>{tbody}</tbody></table>")
        elif t == "quote":
            cite = b.get("cite")
            cite_html = f"<cite>— {esc(cite)}</cite>" if cite else ""
            parts.append(f"<blockquote>{esc(b.get('text'))}{cite_html}</blockquote>")
        elif t == "code":
            parts.append(f"<pre><code>{esc(b.get('text'))}</code></pre>")
        elif t == "cards":
            # 링크 달린 카드 그리드 — 뉴스/검색결과/북마크/상품목록 공용 문서 원시.
            # 각 item: {title, meta?, summary?, url?, link_label?}. columns(기본 2).
            ncol = max(1, min(4, int(b.get("columns") or 2)))
            cell = []
            for it in (b.get("items") or []):
                if not isinstance(it, dict):
                    continue
                img = it.get("image")
                img_h = f'<img class="card-img" src="{esc(img)}" alt="" loading="lazy">' if img else ""
                meta_h = f'<p class="card-meta">{esc(it.get("meta"))}</p>' if it.get("meta") else ""
                sum_h = f'<p class="card-sum">{esc(it.get("summary"))}</p>' if it.get("summary") else ""
                url = it.get("url")
                link_h = (f'<a href="{esc(url)}" target="_blank" rel="noopener">'
                          f'{esc(it.get("link_label") or "열기")}</a>') if url else ""
                cell.append(f'<div class="card">{img_h}<h3>{esc(it.get("title"))}</h3>{meta_h}{sum_h}{link_h}</div>')
            parts.append(f'<div class="cards" style="--cols:{ncol}">{"".join(cell)}</div>')
        elif t == "divider":
            parts.append("<hr>")
        else:  # paragraph (기본)
            parts.append(f"<p>{esc(b.get('text'))}</p>")
    return "\n".join(parts)


def _item_line(it) -> str:
    """목록 항목 → 한 줄 평문(note 포함). 링크를 못 다는 emitter(pptx 등)에서 유실 방지."""
    if not isinstance(it, dict):
        return str(it)
    return f"{it.get('text') or ''} {it.get('note') or ''}".strip()


def _doc_css(theme: str) -> str:
    """문서 emitter의 <style> 본문을 테마별로. default(기사형) / newspaper(제호+카드 그리드)."""
    # 공통: 카드 그리드 골격(테마 무관 동일 구조, 색·여백만 테마가 덧칠)
    base_cards = """
.cards{display:grid;grid-template-columns:repeat(var(--cols,2),1fr);gap:18px;margin:1.2em 0}
@media(max-width:680px){.cards{grid-template-columns:1fr}}
.card{border:1px solid var(--line,#e1e4e8);border-radius:10px;padding:16px 18px;display:flex;flex-direction:column;background:var(--card,#fff)}
.card img.card-img{width:100%;max-height:200px;object-fit:contain;border-radius:6px;margin-bottom:10px;background:#f5f5f5}
.card h3{margin:0 0 8px;font-size:1.08em;line-height:1.4;color:var(--ink,#22223b)}
.card .card-meta{color:#888;font-size:0.82em;margin:0 0 8px}
.card .card-sum{color:var(--ink,#444);font-size:0.92em;margin:0 0 10px;flex:1}
.card a{margin-top:auto;color:#3d5a80;font-weight:bold;font-size:0.9em;text-decoration:none}
.card a:hover{text-decoration:underline}
.li-note{color:var(--dim,#4a5563)}
"""
    if theme == "newspaper":
        return """
body{max-width:1100px;margin:30px auto;padding:0 16px;font-family:'Pretendard',-apple-system,'Apple SD Gothic Neo','Noto Sans KR','Malgun Gothic',sans-serif;line-height:1.6;color:#333;background:#f0f2f5}
.docwrap{background:#fff;padding:40px;border-radius:12px;box-shadow:0 4px 20px rgba(0,0,0,0.08)}
h1{color:#1a1a2e;font-size:2.5em;margin:0 0 15px;border-bottom:4px solid #1a1a2e;padding-bottom:15px;text-align:center}
.doc-meta{text-align:center;color:#666;font-size:0.95em;margin-bottom:30px;background:#f8f9fa;padding:15px;border-radius:8px}
h2{color:#1a1a2e;font-size:1.8em;margin:40px 0 20px;padding-bottom:8px;border-bottom:2px solid #eee}
h3{color:#22223b}
img{max-width:100%;border-radius:8px} figure{margin:1.2em 0} figcaption{color:#666;font-size:0.9em;text-align:center}
table{border-collapse:collapse;width:100%;margin:1.2em 0} th,td{border:1px solid #ddd;padding:8px 12px;text-align:left} th{background:#f5f5f5}
blockquote{border-left:4px solid #ccc;margin:1.2em 0;padding:0.5em 1em;color:#555}
hr{border:none;border-top:1px solid #eee;margin:30px 0}
""" + base_cards
    # default(기사형 단일단) — 읽는 문서용. html/pdf/png 가 같은 CSS 를 쓰므로
    # 화면(카드+배경틴트)·인쇄(흰 바탕)·다크모드를 전부 여기서 갈라둔다.
    return """
:root{--bg:#f6f7f9;--card:#fff;--ink:#1b2330;--dim:#6b7683;--line:#e2e6ec;--accent:#2f5fd0;--code:#eef1f6}
@media(prefers-color-scheme:dark){:root{--bg:#12151a;--card:#191d24;--ink:#e4e8ee;--dim:#98a3b0;--line:#2a313b;--accent:#7aa2f7;--code:#212831}}
*{box-sizing:border-box}
body{margin:0;background:var(--bg);color:var(--ink);line-height:1.75;-webkit-text-size-adjust:100%;
  font-family:'Pretendard',-apple-system,'Apple SD Gothic Neo','Noto Sans KR','Malgun Gothic',sans-serif}
.docwrap{max-width:820px;margin:32px auto;background:var(--card);border:1px solid var(--line);
  border-radius:14px;padding:34px 38px 46px}
@media(max-width:680px){.docwrap{margin:0;border:0;border-radius:0;padding:22px 16px 40px}}
.docwrap>h1:first-child{margin-top:0}
h1,h2,h3,h4{line-height:1.32;font-weight:700}
h1{font-size:clamp(1.5rem,4.5vw,2.05rem);margin:0 0 .5em;padding-bottom:.34em;border-bottom:2px solid var(--line);letter-spacing:-.01em}
h2{font-size:1.24rem;margin:1.9em 0 .6em;padding-left:11px;border-left:5px solid var(--accent)}
h3{font-size:1.06rem;margin:1.5em 0 .45em;color:var(--dim)}
p{margin:.75em 0;word-break:break-word}
.doc-meta{color:var(--dim);font-size:.9rem;margin:-.2em 0 1.6em}
ul,ol{margin:.7em 0 1em;padding-left:1.35em} li{margin:.34em 0}
.li-note{color:var(--dim)}
a{color:var(--accent);text-decoration:none;word-break:break-word} a:hover{text-decoration:underline}
img{max-width:100%;border-radius:8px} figure{margin:1.5em 0}
figcaption{color:var(--dim);font-size:.88rem;text-align:center;margin-top:.5em}
table{border-collapse:collapse;width:100%;margin:1.4em 0;font-size:.93rem;display:block;overflow-x:auto}
th,td{border:1px solid var(--line);padding:8px 11px;text-align:left} th{background:var(--code)}
blockquote{margin:1.2em 0;padding:.6em 15px;border-left:4px solid var(--line);background:var(--code);
  color:var(--dim);border-radius:0 8px 8px 0}
blockquote cite{display:block;margin-top:.5em;font-size:.9em;font-style:normal}
code{background:var(--code);border-radius:5px;padding:.1em .38em;font-size:.9em;
  font-family:'SF Mono',ui-monospace,Menlo,Consolas,monospace}
pre{background:var(--code);padding:1em;border-radius:10px;overflow:auto} pre code{background:none;padding:0}
hr{border:0;border-top:1px solid var(--line);margin:2em 0}
@media print{
  body{background:#fff} 
  .docwrap{margin:0;max-width:none;border:0;border-radius:0;padding:0}
  a{color:inherit;text-decoration:underline}
}
""" + base_cards


def _resolve_image_bytes(src: str):
    """이미지 src(로컬 경로/data URI/http URL)를 file-like(BytesIO)로 해소. 실패 시 None.
    docx·pptx emitter 공용 — 둘 다 파일/스트림만 받음(HTML <img>와 달리)."""
    import io
    import os
    import base64
    if not src:
        return None
    try:
        s = str(src).strip()
        if s.startswith("data:"):  # data:image/png;base64,....
            b64 = s.split(",", 1)[1]
            return io.BytesIO(base64.b64decode(b64))
        if s.startswith("file://"):
            s = s[7:]
        if s.startswith(("http://", "https://")):
            import urllib.request
            with urllib.request.urlopen(s, timeout=15) as r:
                return io.BytesIO(r.read())
        if os.path.isfile(s):
            with open(s, "rb") as f:
                return io.BytesIO(f.read())
    except Exception:
        return None
    return None


def _typ_esc(s) -> str:
    """typst 마크업 특수문자 이스케이프."""
    s = str(s if s is not None else "")
    for ch in ("\\", "#", "$", "*", "_", "`", "<", ">", "@", "[", "]"):
        s = s.replace(ch, "\\" + ch)
    return s


def _doc_blocks_to_typst(blocks: list, title: str, meta: str, out_path: str):
    """문서 IR → typst 컴파일 PDF (책 품질 조판). 산문·보고서에 최적 — 정렬·페이지·타이포가 강점.
    HTML theme/cards 그리드는 무시(조판 모델이 다름). 한글 = Apple SD Gothic Neo."""
    import os
    import subprocess
    import tempfile

    lines = [
        '#set text(font: "Apple SD Gothic Neo", size: 11pt, lang: "ko")',
        '#set page(paper: "a4", margin: (x: 2.2cm, y: 2.4cm), numbering: "1")',
        '#set par(justify: true, leading: 0.8em)',
        '#show heading: set block(above: 1.2em, below: 0.6em)',
        '#set heading(numbering: none)',
        "",
    ]
    if title:
        lines.append(f'#align(center)[#text(size: 22pt, weight: "bold")[{_typ_esc(title)}]]')
    if meta:
        lines.append(f'#align(center)[#text(size: 9pt, fill: gray)[{_typ_esc(meta)}]]')
    if title or meta:
        lines.append("#v(0.5em)")
        lines.append("")

    for b in blocks:
        if not isinstance(b, dict):
            continue
        t = (b.get("type") or "paragraph").lower()
        if t == "heading":
            lvl = max(1, min(5, int(b.get("level") or 2)))
            lines.append("=" * lvl + " " + _typ_esc(b.get("text")))
        elif t == "list":
            for it in (b.get("items") or []):
                txt = it.get("text") if isinstance(it, dict) else it
                if isinstance(it, dict) and it.get("note"):
                    txt = f"{txt} {it['note']}"
                lines.append("- " + _typ_esc(txt))
        elif t == "table":
            cols = b.get("columns") or []
            rows = [r for r in (b.get("rows") or []) if isinstance(r, (list, tuple))]
            ncol = max([len(cols)] + [len(r) for r in rows] or [0])
            if ncol:
                cells = []
                if cols:
                    cells += [f"[*{_typ_esc(c)}*]" for c in cols[:ncol]] + ["[]"] * (ncol - len(cols[:ncol]))
                for r in rows:
                    cells += [f"[{_typ_esc(v)}]" for v in r[:ncol]] + ["[]"] * (ncol - len(r[:ncol]))
                lines.append(f"#table(columns: {ncol}, " + ", ".join(cells) + ")")
        elif t == "cards":
            for it in (b.get("items") or []):
                if not isinstance(it, dict):
                    continue
                lines.append("=== " + _typ_esc(it.get("title")))
                if it.get("meta"):
                    lines.append(f'#text(size: 9pt, fill: gray)[{_typ_esc(it.get("meta"))}]')
                if it.get("summary"):
                    lines.append(_typ_esc(it.get("summary")))
                if it.get("url"):
                    lines.append(f'#link("{it.get("url")}")[{_typ_esc(it.get("link_label") or "열기")}]')
        elif t == "quote":
            cite = b.get("cite")
            q = f'#quote(block: true)[{_typ_esc(b.get("text"))}]'
            lines.append(q + (f" #text(size: 9pt, fill: gray)[— {_typ_esc(cite)}]" if cite else ""))
        elif t == "code":
            lines.append("```\n" + str(b.get("text") or "") + "\n```")
        elif t == "image":
            src = b.get("src") or b.get("path") or ""
            if src and os.path.isfile(str(src)):
                lines.append(f'#figure(image("{src}", width: 80%)' +
                             (f', caption: [{_typ_esc(b.get("caption"))}]' if b.get("caption") else "") + ")")
        elif t == "divider":
            lines.append("#line(length: 100%, stroke: 0.5pt + gray)")
        else:
            lines.append(_typ_esc(b.get("text")))
        lines.append("")

    typ_src = "\n".join(lines)
    with tempfile.NamedTemporaryFile("w", suffix=".typ", delete=False, encoding="utf-8") as tf:
        tf.write(typ_src)
        typ_file = tf.name
    try:
        proc = subprocess.run(["typst", "compile", typ_file, out_path],
                              capture_output=True, text=True, timeout=60)
        if proc.returncode != 0:
            raise RuntimeError(f"typst compile 실패: {proc.stderr[:300]}")
    finally:
        try:
            os.unlink(typ_file)
        except Exception:
            pass


def _add_hyperlink(paragraph, url: str, text: str):
    """python-docx 문단에 클릭 가능한 하이퍼링크 추가(네이티브 미지원이라 관계+XML 수작업)."""
    from docx.oxml.ns import qn
    from docx.oxml import OxmlElement
    try:
        part = paragraph.part
        r_id = part.relate_to(url, "http://schemas.openxmlformats.org/officeDocument/2006/relationships/hyperlink",
                              is_external=True)
        hyperlink = OxmlElement("w:hyperlink")
        hyperlink.set(qn("r:id"), r_id)
        new_run = OxmlElement("w:r")
        rPr = OxmlElement("w:rPr")
        color = OxmlElement("w:color"); color.set(qn("w:val"), "3D5A80"); rPr.append(color)
        u = OxmlElement("w:u"); u.set(qn("w:val"), "single"); rPr.append(u)
        new_run.append(rPr)
        t = OxmlElement("w:t"); t.text = text or url; new_run.append(t)
        hyperlink.append(new_run)
        paragraph._p.append(hyperlink)
    except Exception:
        paragraph.add_run(f"{text}: {url}" if text else url)


def _doc_blocks_to_docx(blocks: list, title: str, out_path: str):
    """문서 IR → .docx (python-docx). html emitter와 같은 IR을 소비.
    table 블록 = 데이터 통화 {columns,rows} 그대로 재사용."""
    from docx import Document
    from docx.shared import Pt, Inches
    from docx.oxml.ns import qn
    from docx.oxml import OxmlElement

    doc = Document()
    if title:
        doc.add_heading(str(title), level=0)

    for b in blocks:
        if not isinstance(b, dict):
            continue
        t = (b.get("type") or "paragraph").lower()
        if t == "heading":
            lvl = max(1, min(6, int(b.get("level") or 2)))
            doc.add_heading(str(b.get("text") or ""), level=lvl)
        elif t == "list":
            ordered = bool(b.get("ordered"))
            style = "List Number" if ordered else "List Bullet"
            for it in (b.get("items") or []):
                txt = it.get("text") if isinstance(it, dict) else it
                doc.add_paragraph(str(txt), style=style)
        elif t == "image":
            stream = _resolve_image_bytes(b.get("src") or b.get("path") or "")
            if stream is not None:
                try:
                    doc.add_picture(stream, width=Inches(6.0))
                except Exception:
                    pass
            cap = b.get("caption")
            if cap:
                p = doc.add_paragraph(str(cap))
                p.italic = True
        elif t == "table":
            cols = b.get("columns") or []
            rows = [r for r in (b.get("rows") or []) if isinstance(r, (list, tuple))]
            ncol = max([len(cols)] + [len(r) for r in rows] or [0])
            if ncol:
                tbl = doc.add_table(rows=0, cols=ncol)
                try:
                    tbl.style = "Light Grid Accent 1"
                except Exception:
                    pass
                if cols:
                    hdr = tbl.add_row().cells
                    for i, c in enumerate(cols[:ncol]):
                        hdr[i].text = str(c)
                    for cell in hdr:
                        for para in cell.paragraphs:
                            for run in para.runs:
                                run.bold = True
                for r in rows:
                    cells = tbl.add_row().cells
                    for i, v in enumerate(r[:ncol]):
                        cells[i].text = str(v)
        elif t == "quote":
            p = doc.add_paragraph(str(b.get("text") or ""))
            try:
                p.style = "Intense Quote"
            except Exception:
                pass
            cite = b.get("cite")
            if cite:
                doc.add_paragraph(f"— {cite}")
        elif t == "code":
            p = doc.add_paragraph()
            run = p.add_run(str(b.get("text") or ""))
            run.font.name = "Courier New"
            run.font.size = Pt(9.5)
        elif t == "cards":
            for it in (b.get("items") or []):
                if not isinstance(it, dict):
                    continue
                if it.get("image"):  # 썸네일(표지 등) — 다운로드 실패는 graceful
                    stream = _resolve_image_bytes(it.get("image"))
                    if stream is not None:
                        try:
                            doc.add_picture(stream, width=Inches(1.6))
                        except Exception:
                            pass
                doc.add_heading(str(it.get("title") or ""), level=3)
                if it.get("meta"):
                    mp = doc.add_paragraph()
                    mr = mp.add_run(str(it.get("meta")))
                    mr.italic = True
                    mr.font.size = Pt(9)
                if it.get("summary"):
                    doc.add_paragraph(str(it.get("summary")))
                if it.get("url"):
                    _add_hyperlink(doc.add_paragraph(), str(it.get("url")),
                                   str(it.get("link_label") or "열기"))
        elif t == "divider":
            p = doc.add_paragraph()
            pPr = p._p.get_or_add_pPr()
            pbdr = OxmlElement("w:pBdr")
            bottom = OxmlElement("w:bottom")
            bottom.set(qn("w:val"), "single")
            bottom.set(qn("w:sz"), "6")
            bottom.set(qn("w:space"), "1")
            bottom.set(qn("w:color"), "auto")
            pbdr.append(bottom)
            pPr.append(pbdr)
        else:  # paragraph
            doc.add_paragraph(str(b.get("text") or ""))

    doc.save(out_path)


def _doc_blocks_to_pptx(blocks: list, title: str, out_path: str):
    """문서 IR → .pptx (python-pptx). ★종류 경계 주의: 슬라이드 IR이 아니라 *문서 IR을 슬라이드로 투영*.
    문서 IR이 정본, pptx는 emitter일 뿐 — heading(level≤2)이 새 슬라이드, 그 아래 내용이 글머리표.
    슬라이드 전용 시각 레이아웃이 필요하면 self:slide{op:create} 를 써야지 이걸 쓰면 안 됨."""
    from pptx import Presentation
    from pptx.util import Inches, Pt

    prs = Presentation()
    SW, SH = prs.slide_width, prs.slide_height
    blank = prs.slide_layouts[6]
    title_layout = prs.slide_layouts[0]

    # 표지 슬라이드(title 있으면)
    if title:
        s = prs.slides.add_slide(title_layout)
        s.shapes.title.text = str(title)

    state = {"slide": None, "body": None}

    def new_content_slide(heading_text=""):
        # "제목+내용" 레이아웃(1) — 제목 placeholder + 본문 placeholder
        s = prs.slides.add_slide(prs.slide_layouts[1])
        s.shapes.title.text = str(heading_text or "")
        body_tf = None
        for ph in s.placeholders:
            if ph.placeholder_format.idx == 1:
                body_tf = ph.text_frame
                break
        if body_tf is not None:
            body_tf.clear()
            body_tf.word_wrap = True
        state["slide"], state["body"] = s, body_tf

    def add_bullet(text, level=0, italic=False, mono=False):
        if state["body"] is None:
            new_content_slide("")
        tf = state["body"]
        # clear()가 남긴 빈 첫 문단 재사용, 이후엔 add
        if len(tf.paragraphs) == 1 and not tf.paragraphs[0].runs:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = str(text)
        p.level = min(4, max(0, level))
        for run in p.runs:
            if italic:
                run.font.italic = True
            if mono:
                run.font.name = "Courier New"
                run.font.size = Pt(14)

    for b in blocks:
        if not isinstance(b, dict):
            continue
        t = (b.get("type") or "paragraph").lower()
        if t == "heading":
            lvl = int(b.get("level") or 2)
            if lvl <= 2 or state["body"] is None:
                # 큰 섹션, 또는 표·구분선 뒤 첫 제목이면 새 슬라이드 제목으로
                new_content_slide(b.get("text") or "")
            else:  # 진행 중 슬라이드의 하위 섹션 = 글머리표
                add_bullet(b.get("text") or "", level=0)
        elif t == "paragraph":
            add_bullet(b.get("text") or "", level=0)
        elif t == "list":
            for it in (b.get("items") or []):
                add_bullet(_item_line(it), level=1)
        elif t == "quote":
            txt = str(b.get("text") or "")
            cite = b.get("cite")
            add_bullet(f"“{txt}”" + (f" — {cite}" if cite else ""), level=1, italic=True)
        elif t == "code":
            add_bullet(b.get("text") or "", level=1, mono=True)
        elif t == "cards":
            for it in (b.get("items") or []):
                if not isinstance(it, dict):
                    continue
                add_bullet(it.get("title") or "", level=0)
                sub = " / ".join(x for x in [it.get("meta"), it.get("summary")] if x)
                if sub:
                    add_bullet(sub, level=1)
        elif t == "image":
            stream = _resolve_image_bytes(b.get("src") or b.get("path") or "")
            if stream is not None:
                s = prs.slides.add_slide(blank)
                try:
                    s.shapes.add_picture(stream, Inches(0.6), Inches(0.6), width=SW - Inches(1.2))
                except Exception:
                    pass
                cap = b.get("caption")
                if cap:
                    tb = s.shapes.add_textbox(Inches(0.6), SH - Inches(0.9), SW - Inches(1.2), Inches(0.7))
                    tb.text_frame.text = str(cap)
                state["slide"], state["body"] = None, None  # 이미지 후 새 내용 슬라이드 강제
        elif t == "table":
            cols = b.get("columns") or []
            rows = [r for r in (b.get("rows") or []) if isinstance(r, (list, tuple))]
            ncol = max([len(cols)] + [len(r) for r in rows] or [0])
            nrow = len(rows) + (1 if cols else 0)
            if ncol and nrow:
                s = prs.slides.add_slide(blank)
                gt = s.shapes.add_table(nrow, ncol, Inches(0.5), Inches(0.6),
                                        SW - Inches(1.0), Inches(0.4) * nrow).table
                ri = 0
                if cols:
                    for ci, c in enumerate(cols[:ncol]):
                        gt.cell(0, ci).text = str(c)
                    ri = 1
                for r in rows:
                    for ci, v in enumerate(r[:ncol]):
                        gt.cell(ri, ci).text = str(v)
                    ri += 1
                state["slide"], state["body"] = None, None
        elif t == "divider":
            state["slide"], state["body"] = None, None  # 새 슬라이드 경계

    if len(prs.slides) == 0:  # 표지도 내용도 없으면 빈 슬라이드 하나
        prs.slides.add_slide(blank)
    prs.save(out_path)



def _doc_blocks_to_markdown(blocks: list, title: str = "", meta: str = "") -> str:
    """문서 IR → 마크다운. cards = 링크 목록(뉴스/검색결과 공용). NIP-23 발행 등 범용."""
    out = []
    if title:
        out.append(f"# {title}")
    if meta:
        out.append(f"_{meta}_")
    if out:
        out.append("")
    for b in blocks:
        if not isinstance(b, dict):
            continue
        t = (b.get("type") or "paragraph").lower()
        if t == "heading":
            lvl = max(1, min(6, int(b.get("level") or 2)))
            out += [f"{'#' * lvl} {b.get('text') or ''}", ""]
        elif t == "list":
            ordered = bool(b.get("ordered"))
            for idx, i in enumerate(b.get("items") or [], 1):
                mark = f"{idx}." if ordered else "-"
                if isinstance(i, dict):
                    txt = i.get("text") or ""
                    line = f"{mark} [{txt}]({i.get('url')})" if i.get("url") else f"{mark} {txt}"
                    if i.get("note"):
                        line += f" {i['note']}"
                    out.append(line)
                else:
                    out.append(f"{mark} {i}")
            out.append("")
        elif t == "cards":
            for it in (b.get("items") or []):
                if not isinstance(it, dict):
                    continue
                url = it.get("url")
                out.append(f"- [{it.get('title') or ''}]({url})" if url else f"- {it.get('title') or ''}")
                sub = " — ".join(x for x in [it.get("meta"), it.get("summary")] if x)
                if sub:
                    out.append(f"  {sub}")
            out.append("")
        elif t == "table":
            cols, rows = b.get("columns") or [], b.get("rows") or []
            if cols:
                out.append("| " + " | ".join(str(c) for c in cols) + " |")
                out.append("| " + " | ".join("---" for _ in cols) + " |")
            for r in rows:
                if isinstance(r, (list, tuple)):
                    out.append("| " + " | ".join(str(c) for c in r) + " |")
            out.append("")
        elif t == "quote":
            out.append(f"> {b.get('text') or ''}")
            if b.get("cite"):
                out.append(f"> — {b.get('cite')}")
            out.append("")
        elif t == "code":
            out += ["```", str(b.get("text") or ""), "```", ""]
        elif t == "image":
            out += [f"![{b.get('caption') or ''}]({b.get('src') or b.get('path') or ''})", ""]
        elif t == "divider":
            out += ["---", ""]
        else:  # paragraph
            out += [str(b.get("text") or ""), ""]
    return "\n".join(out).strip() + "\n"

