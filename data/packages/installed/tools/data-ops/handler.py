"""data-ops — 통화→통화 변환자 (currency algebra).

IBL의 깊이(depth)를 만드는 부품. 생산자(sense:* 등)가 내는 공유 통화를 입력으로,
같은 통화를 출력하는 닫힌(closure) 동사들. >> 파이프로 임의 깊이 조합:

    [sense:realty]{...} >> [table:filter]{where:"전세"} >> [table:sort]{by:price} >> [table:take]{n:5}

각 동사는 도메인에 무관 — 한 번 짜서 34개 생산자 × 임의 깊이에 곱해진다. 이 파일이
없던 시절엔 sort/top-N/dedup이 ~19개 핸들러에 사적으로 중복 구현돼 있었다(부채 회수).

통화 (단일):
  - items = {"items": [ {…열린 필드…} ]}  — 유일한 컬렉션 통화(2026-06-27 컷오버 완료).
    목록형·table 모두 items로 수렴. table(columns/rows)은 items의 파생 뷰(_get_table 가 items→table 재구성).
    변환자는 내부에서 행 dict(=items)로 정규화해 도메인·통화종류 무관하게 한 코드로 처리한다.

단항(filter/sort/take/select/dedup/groupby)·이항(join/union/merge, & 병렬 두 입력) 모두 구현됨.

+ 표준 코어 문서 emitter(2026-07-03 media_producer에서 이관 — 표준 코어 table 어휘가 개인
패키지에 살던 경계 이상 해소): [table:structure]=콘텐츠→문서 IR(경량 LLM 편집자),
[table:document]=문서 IR→html/pdf/png/docx/pptx/typst. 무거운 의존성(playwright·docx·pptx·typst)은
전부 함수 안 지연 import — 모듈레벨은 stdlib(json/re)만이라 폰 import-safe 불변식 유지.
"""

import json
import re


# ───────────────────────── 통화 추출/주입 (공유) ─────────────────────────

def _parse_prev(prev):
    """_prev_result(JSON 문자열 또는 dict/list) → 파이썬 객체."""
    if prev is None:
        return None
    if isinstance(prev, (dict, list)):
        return prev
    if isinstance(prev, str):
        try:
            return json.loads(prev)
        except Exception:
            return None
    return None


def _get_items(obj):
    """객체에서 항목 리스트를 꺼낸다.

    단일 통화 `items`가 유일한 키다(2026-06-27 단일통화 컷오버 완료). 모든 생산자가
    native 풍부 dict를 items로 방출하므로 옛 카드 통화가 버리던 필드까지 in-pipe로 흐른다
    (예: zigbang lat/lng, business level). dict/list 가 아니면 통화 아님으로 (None, None) 반환.

    반환: (items_list, envelope_dict)  — envelope에 변환 결과를 다시 끼워 비파괴 반환용.
    """
    if isinstance(obj, dict):
        r = obj.get("items")            # 단일 통화 — 모든 생산자가 items 방출
        if isinstance(r, list):
            return r, obj
    if isinstance(obj, list):
        return obj, {"items": obj}
    return None, None


def _get_table(obj):
    """객체에서 표준 table 통화를 꺼낸다. ({table:{columns,rows}} 또는 최상위 columns/rows)

    반환: (table_dict, envelope_dict).
    """
    if isinstance(obj, dict):
        t = obj.get("table")
        if isinstance(t, dict) and isinstance(t.get("rows"), list):
            return {"columns": t.get("columns") or [], "rows": t["rows"]}, obj
        if isinstance(obj.get("rows"), list) and isinstance(obj.get("columns"), list):
            return {"columns": obj["columns"], "rows": obj["rows"]}, obj
        # 단일 통화 items(행 dict) → table 재구성: 첫 dict의 키 순서=열, 값=행(§3 table 흡수).
        items = obj.get("items")
        if isinstance(items, list) and items and all(isinstance(x, dict) for x in items):
            cols = list(items[0].keys())
            return {"columns": cols, "rows": [[d.get(c) for c in cols] for d in items]}, obj
    return None, None


def _emit_items(envelope, new_items):
    """변환된 항목들을 원 envelope에 비파괴로 끼워 반환.

    단일 통화 키 `items`로 내보낸다(2026-06-27 단일통화 컷오버 완료 — 옛 이중방출 은퇴).
    """
    out = dict(envelope) if isinstance(envelope, dict) else {}
    out.pop("message", None)            # 변환 후 stale·O(items) 산문 제거 (파이프 블로업·정합성)
    out["items"] = new_items          # 단일 통화
    out["count"] = len(new_items)
    out.setdefault("success", True)
    return out


def _emit_table(envelope, new_table):
    out = dict(envelope) if isinstance(envelope, dict) else {}
    out.pop("message", None)            # 변환 후 stale 산문 제거
    if "table" in out:
        out["table"] = new_table
    else:
        out["columns"] = new_table.get("columns", [])
        out["rows"] = new_table.get("rows", [])
    out.setdefault("success", True)
    return out


def _row_dicts(table):
    """table rows → [{col: val}] (where/sort/dedup이 items와 같은 코드 쓰도록)."""
    cols = table.get("columns") or []
    out = []
    for r in table.get("rows") or []:
        d = {}
        for i, c in enumerate(cols):
            d[str(c)] = r[i] if i < len(r) else None
        out.append(d)
    return out


# ───────────────────────── where 미니 DSL ─────────────────────────

_OPS = {
    "==": lambda a, b: _num_eq(a, b),
    "eq": lambda a, b: _num_eq(a, b),
    "!=": lambda a, b: not _num_eq(a, b),
    "ne": lambda a, b: not _num_eq(a, b),
    "<": lambda a, b: _num_cmp(a, b) < 0,
    "lt": lambda a, b: _num_cmp(a, b) < 0,
    "<=": lambda a, b: _num_cmp(a, b) <= 0,
    "le": lambda a, b: _num_cmp(a, b) <= 0,
    ">": lambda a, b: _num_cmp(a, b) > 0,
    "gt": lambda a, b: _num_cmp(a, b) > 0,
    ">=": lambda a, b: _num_cmp(a, b) >= 0,
    "ge": lambda a, b: _num_cmp(a, b) >= 0,
    "contains": lambda a, b: str(b).lower() in str(a).lower(),
    "in": lambda a, b: (a in b) if isinstance(b, (list, tuple, set)) else (str(a).lower() in str(b).lower()),
}


def _as_num(v):
    try:
        return float(str(v).replace(",", "").strip())
    except Exception:
        return None


def _num_eq(a, b):
    na, nb = _as_num(a), _as_num(b)
    if na is not None and nb is not None:
        return na == nb
    return str(a).strip().lower() == str(b).strip().lower()


def _num_cmp(a, b):
    na, nb = _as_num(a), _as_num(b)
    if na is not None and nb is not None:
        return (na > nb) - (na < nb)
    sa, sb = str(a), str(b)
    return (sa > sb) - (sa < sb)


_CMP_RE = re.compile(r"^\s*(.+?)\s*(>=|<=|==|!=|>|<|=)\s*(.+?)\s*$")


def _match(item, where):
    """item(dict) 이 where 조건을 만족하나.

    where 형태:
      - str "필드 op 값"  : 비교 연산자(>= <= > < == != =)가 있으면 단일 비교로 파싱
                            (예 "연도 >= 2000" → {field:연도, op:>=, value:2000}).
                            모델이 자연스럽게 쓰는 SQL식 문자열을 침묵 부분일치로 삼키지 않는다.
      - str S            : 연산자 없으면 아무 필드 값에 S가 부분일치 (전 필드 substring)
      - {field, op, value}: SQL식 단일 조건 (op 기본 ==; field=col/column 별칭)
      - {col: value, ...}: 각 열=값 동등(AND) 단축형
      - [cond, cond, ...]: AND 결합
    """
    if where is None or where == "":
        return True
    if isinstance(where, str):
        m = _CMP_RE.match(where)
        if m:  # 비교 연산자가 든 문자열 → 단일 비교로 해석 (침묵 부분일치 함정 제거)
            field, op, val = m.group(1).strip(), m.group(2), m.group(3).strip()
            if op == "=":
                op = "=="
            if len(val) >= 2 and val[0] in "\"'" and val[-1] == val[0]:
                val = val[1:-1]  # 따옴표 제거
            fn = _OPS.get(op, _OPS["=="])
            return fn(item.get(field), val)
        s = where.lower()
        return any(s in str(v).lower() for v in item.values())
    if isinstance(where, list):
        return all(_match(item, w) for w in where)
    if isinstance(where, dict):
        field = where.get("field") or where.get("col") or where.get("column")
        if field is not None:  # 구조형 {field, op, value}
            op = str(where.get("op", "==")).lower()
            val = where.get("value")
            fn = _OPS.get(op, _OPS["=="])
            return fn(item.get(str(field)), val)
        # 단축형 {col: value, ...} — 모두 동등(AND)
        return all(_num_eq(item.get(str(k)), v) for k, v in where.items())
    return True


# ───────────────────────── sort 키 (수치 인식) ─────────────────────────

def _sort_key(field):
    def key(item):
        v = item.get(str(field)) if isinstance(item, dict) else None
        n = _as_num(v)
        # 숫자 먼저(0) 안정 정렬, 그다음 문자열(1). None은 맨 뒤.
        if v is None:
            return (2, 0.0, "")
        if n is not None:
            return (0, n, "")
        return (1, 0.0, str(v).lower())
    return key


# ───────────────────────── 단항 동사 ─────────────────────────

def _op_filter(prev, params):
    """items|table → 부분집합. params.where (미니 DSL). condition 별칭 수용."""
    where = params.get("where") or params.get("condition")
    recs, env = _get_items(prev)
    if recs is not None:
        return _emit_items(env, [r for r in recs if isinstance(r, dict) and _match(r, where)])
    table, env = _get_table(prev)
    if table is not None:
        kept = [d for d in _row_dicts(table) if _match(d, where)]
        cols = table.get("columns") or []
        rows = [[d.get(str(c)) for c in cols] for d in kept]
        return _emit_table(env, {"columns": cols, "rows": rows})
    return {"error": "filter: 입력에서 items 통화를 찾지 못했습니다."}


def _op_sort(prev, params):
    """items|table → 정렬. params.by(필드/열명), params.desc(bool)."""
    by = params.get("by")
    desc = bool(params.get("desc", False))
    if not by:
        return {"error": "sort: by(정렬 기준 필드/열명)가 필요합니다."}
    recs, env = _get_items(prev)
    if recs is not None:
        srt = sorted([r for r in recs if isinstance(r, dict)], key=_sort_key(by), reverse=desc)
        return _emit_items(env, srt)
    table, env = _get_table(prev)
    if table is not None:
        dicts = _row_dicts(table)
        dicts.sort(key=_sort_key(by), reverse=desc)
        cols = table.get("columns") or []
        rows = [[d.get(str(c)) for c in cols] for d in dicts]
        return _emit_table(env, {"columns": cols, "rows": rows})
    return {"error": "sort: 입력에서 items 통화를 찾지 못했습니다."}


def _op_take(prev, params):
    """items|table → 상위 n. params.n (기본 10). 음수면 뒤에서 n개."""
    n = params.get("n", params.get("limit", 10))
    try:
        n = int(n)
    except Exception:
        n = 10
    recs, env = _get_items(prev)
    if recs is not None:
        sliced = recs[n:] if n < 0 else recs[:n]
        return _emit_items(env, sliced)
    table, env = _get_table(prev)
    if table is not None:
        rows = table.get("rows") or []
        sliced = rows[n:] if n < 0 else rows[:n]
        return _emit_table(env, {"columns": table.get("columns") or [], "rows": sliced})
    return {"error": "take: 입력에서 items 통화를 찾지 못했습니다."}


def _op_select(prev, params):
    """table → 열 투영. params.columns(남길 열 이름 배열). items는 필드 추림."""
    cols_keep = params.get("columns") or params.get("cols") or params.get("fields")
    if not cols_keep:
        return {"error": "select: columns(남길 열/필드 이름 배열)가 필요합니다."}
    cols_keep = [str(c) for c in cols_keep]
    table, env = _get_table(prev)
    if table is not None:
        src_cols = [str(c) for c in (table.get("columns") or [])]
        idx = [src_cols.index(c) for c in cols_keep if c in src_cols]
        new_cols = [src_cols[i] for i in idx]
        new_rows = [[(r[i] if i < len(r) else None) for i in idx] for r in (table.get("rows") or [])]
        return _emit_table(env, {"columns": new_cols, "rows": new_rows})
    recs, env = _get_items(prev)
    if recs is not None:
        out = [{k: r.get(k) for k in cols_keep if k in r} for r in recs if isinstance(r, dict)]
        return _emit_items(env, out)
    return {"error": "select: 입력에서 items 통화를 찾지 못했습니다."}


def _norm(s):
    import re
    return re.sub(r"\s+", " ", str(s or "").strip().lower())


def _op_dedup(prev, params):
    """items|table → 중복 제거(첫 항목 유지). params.by(키 필드/열, 기본 title).

    by 미지정 시 items는 title, table은 첫 열을 키로. 정규화(공백/대소문자) 후 동등 비교.
    (newspaper 내부에 묻혀있던 _dedup_rank를 통화 동사로 끌어올림 — 전 생산자 공용화.)
    """
    by = params.get("by")
    recs, env = _get_items(prev)
    if recs is not None:
        key = by or "title"
        seen, out = set(), []
        for r in recs:
            if not isinstance(r, dict):
                continue
            k = _norm(r.get(key))
            if k and k in seen:
                continue
            seen.add(k)
            out.append(r)
        return _emit_items(env, out)
    table, env = _get_table(prev)
    if table is not None:
        cols = [str(c) for c in (table.get("columns") or [])]
        ki = cols.index(str(by)) if (by and str(by) in cols) else 0
        seen, rows = set(), []
        for r in table.get("rows") or []:
            k = _norm(r[ki] if ki < len(r) else None)
            if k and k in seen:
                continue
            seen.add(k)
            rows.append(r)
        return _emit_table(env, {"columns": cols, "rows": rows})
    return {"error": "dedup: 입력에서 items 통화를 찾지 못했습니다."}


_AGG = {
    "count": lambda vs: len(vs),
    "sum": lambda vs: round(sum(_as_num(v) or 0 for v in vs), 6),
    "avg": lambda vs: round(sum(_as_num(v) or 0 for v in vs) / len(vs), 6) if vs else 0,
    "min": lambda vs: min((_as_num(v) for v in vs if _as_num(v) is not None), default=None),
    "max": lambda vs: max((_as_num(v) for v in vs if _as_num(v) is not None), default=None),
}


def _rows_for_field(obj, field):
    """그룹/집계 키 field 를 실제로 담은 표현을 골라 list-of-dicts 로 반환.

    공유 통화 items(및 table) 외에도 도메인 생산자가 원천 행을 다른 키
    (data/results)에 담는 경우까지 후보로 본다. 옛 손실적 카드 투영이 도메인 필드를
    meta 문자열로 접어 groupby 가 그 필드를 못 찾던 문제(예: realty 가 '법정동'을 meta 로
    접음 — 원천은 data:[{...법정동...}] 에 그대로 있음)를 푼다.

    후보 중 field 를 가진 첫 리스트를 우선(없으면 첫 비어있지 않은 리스트)으로 고른다.
    """
    cands = []  # [(키, list-of-dicts)]
    if isinstance(obj, list):
        cands.append(("(root)", [x for x in obj if isinstance(x, dict)]))
    elif isinstance(obj, dict):
        t, _ = _get_table(obj)
        if t is not None:
            cands.append(("table", _row_dicts(t)))
        for k in ("data", "items", "results"):
            v = obj.get(k)
            if isinstance(v, list) and any(isinstance(x, dict) for x in v):
                cands.append((k, [x for x in v if isinstance(x, dict)]))
    if not cands:
        return None
    if field:
        for _, rows in cands:
            if any(field in r for r in rows):
                return rows
    for _, rows in cands:
        if rows:
            return rows
    return cands[0][1]


def _op_groupby(prev, params):
    """items|table|도메인 행목록 → 그룹 집계. params.by(그룹 키), params.agg.

    agg 형태: {새열명: [op, 원본열]} 또는 {원본열: op}. op = count/sum/avg/min/max.
    기본: agg 없으면 그룹별 count.  반환 table = [by열, 집계열들].

    형제 동사들처럼 items 를 받는다. 나아가 옛 카드 투영이 그룹 키를 접은 경우
    원천 data/items 행까지 거슬러 키를 찾는다(_rows_for_field).
    """
    by = params.get("by") or params.get("key") or params.get("group_by")
    if not by:
        return {"error": "groupby: by(그룹 키 열)가 필요합니다."}
    by = str(by)
    dicts = _rows_for_field(prev, by)
    if not dicts:
        return {"error": "groupby: 입력에서 items 통화(또는 data/items 행 목록)를 찾지 못했습니다."}
    _, env = _get_table(prev)
    env = env or {}
    agg = params.get("agg")
    # agg 정규화 → [(out_col, op, src_col)]
    specs = []
    if isinstance(agg, dict):
        for k, v in agg.items():
            if isinstance(v, (list, tuple)) and len(v) == 2:  # {새열: [op, 원본열]}
                specs.append((str(k), str(v[0]).lower(), str(v[1])))
            else:  # {원본열: op}
                specs.append((f"{v}_{k}", str(v).lower(), str(k)))
    if not specs:
        specs = [("count", "count", by)]
    # 그룹핑 (입력 순서 보존)
    groups, order = {}, []
    for d in dicts:
        gk = d.get(by)
        if gk not in groups:
            groups[gk] = []
            order.append(gk)
        groups[gk].append(d)
    out_cols = [by] + [s[0] for s in specs]
    out_rows = []
    for gk in order:
        members = groups[gk]
        row = [gk]
        for out_col, op, src in specs:
            vals = [m.get(src) for m in members]
            fn = _AGG.get(op, _AGG["count"])
            row.append(fn(vals))
        out_rows.append(row)
    return _emit_table(env, {"columns": out_cols, "rows": out_rows})


# ───────────────────────── 이항 동사 (& 병렬 두 입력) ─────────────────────────

def _extract_two(prev):
    """& 병렬 결과에서 두 객체 추출.

    prev = [elem0, elem1] (각 elem은 dict 또는 JSON 문자열 — 병렬 분기 결과).
    반환: (obj0, obj1) 둘 다 dict/list로 파싱. 부족하면 (None, None).
    """
    if not isinstance(prev, list) or len(prev) < 2:
        return None, None

    def _parse_elem(e):
        if isinstance(e, str):
            try:
                return json.loads(e)
            except Exception:
                return None
        return e

    return _parse_elem(prev[0]), _parse_elem(prev[1])


def _op_union(prev, params):
    """두 table(또는 두 items)을 행 결합. 같은 통화끼리. params 없음.

    table: 열 이름으로 통합(순서 보존, 한쪽에만 있는 열은 다른쪽 None). items: 단순 concat.
    중복 제거가 필요하면 뒤에 >> dedup.
    """
    a, b = _extract_two(prev)
    if a is None or b is None:
        return {"error": "union: & 병렬로 두 입력이 필요합니다. 예: [A] & [B] >> [table:union]"}
    ta, _ = _get_table(a)
    tb, _ = _get_table(b)
    if ta is not None and tb is not None:
        cols = [str(c) for c in (ta.get("columns") or [])]
        for c in (tb.get("columns") or []):
            if str(c) not in cols:
                cols.append(str(c))

        def remap(t):
            tcols = [str(c) for c in (t.get("columns") or [])]
            out = []
            for r in t.get("rows") or []:
                d = {tcols[i]: (r[i] if i < len(r) else None) for i in range(len(tcols))}
                out.append([d.get(c) for c in cols])
            return out

        return _emit_table({"table": {}}, {"columns": cols, "rows": remap(ta) + remap(tb)})
    ra, _ = _get_items(a)
    rb, _ = _get_items(b)
    if ra is not None and rb is not None:
        return _emit_items({}, list(ra) + list(rb))
    return {"error": "union: 두 입력의 통화 종류가 같아야 합니다(둘 다 table 또는 둘 다 items)."}


def _op_merge(prev, params):
    """두 items를 합친다(concat). params.by 지정 시 그 키로 중복 제거.

    여러 검색 결과를 한 목록으로 모을 때. (table 결합은 union.)
    """
    a, b = _extract_two(prev)
    if a is None or b is None:
        return {"error": "merge: & 병렬로 두 items 입력이 필요합니다. 예: [A] & [B] >> [table:merge]"}
    ra, _ = _get_items(a)
    rb, _ = _get_items(b)
    if ra is None or rb is None:
        return {"error": "merge: 두 입력 모두 items 통화여야 합니다(표형 결합은 table:union)."}
    out = list(ra) + list(rb)
    by = params.get("by")
    if by or params.get("dedup"):
        key = by or "title"
        seen, dd = set(), []
        for r in out:
            if not isinstance(r, dict):
                continue
            k = _norm(r.get(key))
            if k and k in seen:
                continue
            seen.add(k)
            dd.append(r)
        out = dd
    return _emit_items({}, out)


def _suffix_collisions(base_cols, add_cols):
    """add_cols 이름이 base_cols(또는 서로)와 충돌하면 _2,_3.. 접미사 — *표시 이름만*.
    (읽기는 원본 이름으로 따로 한다. 동명 열이 겹치면 다운스트림 select/sort가 이름으로
     둘째 열을 못 집어 조용히 첫째를 오선택하는 침묵 함정을 막는다.)"""
    used = set(base_cols)
    out = []
    for c in add_cols:
        name = c
        if name in used:
            i = 2
            while f"{c}_{i}" in used:
                i += 1
            name = f"{c}_{i}"
        used.add(name)
        out.append(name)
    return out


def _op_join(prev, params):
    """두 table을 키 열로 inner join. params.on(양쪽 공통 키 열명, 필수).

    결과 열 = 좌측 전체 + 우측(키 제외). 서로 다른 소스를 한 키로 묶어 분석.
    예: [sense:stock]{op:history} & [sense:world_bank]{...} >> [table:join]{on: "연도"}.
    """
    on = params.get("on") or params.get("key")
    if not on:
        return {"error": "join: on(조인 키 열 이름)이 필요합니다."}
    on = str(on)
    a, b = _extract_two(prev)
    if a is None or b is None:
        return {"error": "join: & 병렬로 두 table 입력이 필요합니다. 예: [A] & [B] >> [table:join]{on: \"연도\"}"}
    ta, _ = _get_table(a)
    tb, _ = _get_table(b)
    if ta is None or tb is None:
        # 두 입력이 items 통화면 items inner join (table 분기와 대칭).
        # items 행도 dict 라 키 필드로 조인 가능 — merge/union 이 items 를 받는 것과 일관.
        ra, _ = _get_items(a)
        rb, _ = _get_items(b)
        if ra is None or rb is None:
            return {"error": "join: 두 입력이 같은 통화여야 합니다(둘 다 table 또는 둘 다 items)."}
        index = {}
        for r in rb:
            if isinstance(r, dict) and r.get(on) is not None:
                index.setdefault(_norm(r.get(on)), []).append(r)
        out = []
        for l in ra:
            if not isinstance(l, dict) or l.get(on) is None:
                continue
            lkeys = list(l.keys())
            for r in index.get(_norm(l.get(on)), []):
                add = [k for k in r.keys() if k != on]
                disp = _suffix_collisions(lkeys, add)  # 동명 필드 _2 (침묵 오선택 방지)
                merged = dict(l)
                for orig, name in zip(add, disp):
                    merged[name] = r[orig]
                out.append(merged)
        return _emit_items({}, out)
    ca = [str(c) for c in (ta.get("columns") or [])]
    cb = [str(c) for c in (tb.get("columns") or [])]
    if on not in ca or on not in cb:
        return {"error": f"join: 키 '{on}'이 양쪽 table 열에 모두 있어야 합니다(좌:{ca} 우:{cb})."}
    lki, rki = ca.index(on), cb.index(on)
    # 우측을 키로 인덱싱
    index = {}
    for r in tb.get("rows") or []:
        k = _norm(r[rki] if rki < len(r) else None)
        index.setdefault(k, []).append(r)
    extra = [c for c in cb if c != on]  # 우측에서 가져올 열(키 제외, 읽기는 원본 이름)
    out_cols = ca + _suffix_collisions(ca, extra)  # 표시 이름만 충돌 회피
    out_rows = []
    for r in ta.get("rows") or []:
        k = _norm(r[lki] if lki < len(r) else None)
        for rb_row in index.get(k, []):
            rbd = {cb[i]: (rb_row[i] if i < len(rb_row) else None) for i in range(len(cb))}
            out_rows.append(list(r) + [rbd.get(c) for c in extra])
    return _emit_table({"table": {}}, {"columns": out_cols, "rows": out_rows})


# ── 문서 IR(공유 문서 모델) → 산출물 emitter ───────────────────────────
# 문서 IR: {title?, blocks:[{type, ...}]}. 블록 타입:
#   heading{level,text} · paragraph{text} · list{ordered?,items[]} · image{src,caption?}
#   · table{columns,rows}(=데이터 통화 재사용) · quote{text,cite?} · code{text,lang?} · divider
# 포맷 중립 IR을 여러 emitter가 렌더(현재 html). slide/newspaper와 달리 단일 IR이 단일 진실 소스.
def _doc_blocks_to_html(blocks: list) -> str:
    import html as _html

    def esc(s):
        return _html.escape(str(s if s is not None else ""))

    parts = []
    for b in blocks:
        if not isinstance(b, dict):
            continue
        t = (b.get("type") or "paragraph").lower()
        if t == "heading":
            lvl = max(1, min(6, int(b.get("level") or 2)))
            txt = b.get("text")
            anchor = b.get("anchor") or (str(txt) if txt else "")  # 목차 점프용 id
            id_attr = f' id="{esc(anchor)}"' if anchor else ""
            parts.append(f"<h{lvl}{id_attr}>{esc(txt)}</h{lvl}>")
        elif t == "list":
            tag = "ol" if b.get("ordered") else "ul"
            li = []
            for i in (b.get("items") or []):
                # 항목은 문자열 또는 {text, url}(링크 — 목차·북마크 등)
                if isinstance(i, dict):
                    if i.get("url"):
                        li.append(f'<li><a href="{esc(i.get("url"))}">{esc(i.get("text"))}</a></li>')
                    else:
                        li.append(f"<li>{esc(i.get('text'))}</li>")
                else:
                    li.append(f"<li>{esc(i)}</li>")
            parts.append(f"<{tag}>{''.join(li)}</{tag}>")
        elif t == "image":
            src = b.get("src") or b.get("path") or ""
            cap = b.get("caption")
            cap_html = f"<figcaption>{esc(cap)}</figcaption>" if cap else ""
            parts.append(f'<figure><img src="{esc(src)}" alt="{esc(cap)}">{cap_html}</figure>')
        elif t == "table":
            # 데이터 통화 재사용: {columns, rows}
            cols = b.get("columns") or []
            rows = b.get("rows") or []
            thead = "".join(f"<th>{esc(c)}</th>" for c in cols)
            tbody = "".join("<tr>" + "".join(f"<td>{esc(c)}</td>" for c in r) + "</tr>"
                            for r in rows if isinstance(r, (list, tuple)))
            parts.append(f"<table><thead><tr>{thead}</tr></thead><tbody>{tbody}</tbody></table>")
        elif t == "quote":
            cite = b.get("cite")
            cite_html = f"<cite>— {esc(cite)}</cite>" if cite else ""
            parts.append(f"<blockquote>{esc(b.get('text'))}{cite_html}</blockquote>")
        elif t == "code":
            parts.append(f"<pre><code>{esc(b.get('text'))}</code></pre>")
        elif t == "cards":
            # 링크 달린 카드 그리드 — 뉴스/검색결과/북마크/상품목록 공용 문서 원시.
            # 각 item: {title, meta?, summary?, url?, link_label?}. columns(기본 2).
            ncol = max(1, min(4, int(b.get("columns") or 2)))
            cell = []
            for it in (b.get("items") or []):
                if not isinstance(it, dict):
                    continue
                img = it.get("image")
                img_h = f'<img class="card-img" src="{esc(img)}" alt="" loading="lazy">' if img else ""
                meta_h = f'<p class="card-meta">{esc(it.get("meta"))}</p>' if it.get("meta") else ""
                sum_h = f'<p class="card-sum">{esc(it.get("summary"))}</p>' if it.get("summary") else ""
                url = it.get("url")
                link_h = (f'<a href="{esc(url)}" target="_blank" rel="noopener">'
                          f'{esc(it.get("link_label") or "열기")}</a>') if url else ""
                cell.append(f'<div class="card">{img_h}<h3>{esc(it.get("title"))}</h3>{meta_h}{sum_h}{link_h}</div>')
            parts.append(f'<div class="cards" style="--cols:{ncol}">{"".join(cell)}</div>')
        elif t == "divider":
            parts.append("<hr>")
        else:  # paragraph (기본)
            parts.append(f"<p>{esc(b.get('text'))}</p>")
    return "\n".join(parts)


def _doc_css(theme: str) -> str:
    """문서 emitter의 <style> 본문을 테마별로. default(기사형) / newspaper(제호+카드 그리드)."""
    # 공통: 카드 그리드 골격(테마 무관 동일 구조, 색·여백만 테마가 덧칠)
    base_cards = """
.cards{display:grid;grid-template-columns:repeat(var(--cols,2),1fr);gap:18px;margin:1.2em 0}
@media(max-width:680px){.cards{grid-template-columns:1fr}}
.card{border:1px solid #e1e4e8;border-radius:10px;padding:16px 18px;display:flex;flex-direction:column;background:#fff}
.card img.card-img{width:100%;max-height:200px;object-fit:contain;border-radius:6px;margin-bottom:10px;background:#f5f5f5}
.card h3{margin:0 0 8px;font-size:1.08em;line-height:1.4;color:#22223b}
.card .card-meta{color:#888;font-size:0.82em;margin:0 0 8px}
.card .card-sum{color:#444;font-size:0.92em;margin:0 0 10px;flex:1}
.card a{margin-top:auto;color:#3d5a80;font-weight:bold;font-size:0.9em;text-decoration:none}
.card a:hover{text-decoration:underline}
"""
    if theme == "newspaper":
        return """
body{max-width:1100px;margin:30px auto;padding:0 16px;font-family:'Noto Sans KR','Pretendard',-apple-system,sans-serif;line-height:1.6;color:#333;background:#f0f2f5}
.docwrap{background:#fff;padding:40px;border-radius:12px;box-shadow:0 4px 20px rgba(0,0,0,0.08)}
h1{color:#1a1a2e;font-size:2.5em;margin:0 0 15px;border-bottom:4px solid #1a1a2e;padding-bottom:15px;text-align:center}
.doc-meta{text-align:center;color:#666;font-size:0.95em;margin-bottom:30px;background:#f8f9fa;padding:15px;border-radius:8px}
h2{color:#1a1a2e;font-size:1.8em;margin:40px 0 20px;padding-bottom:8px;border-bottom:2px solid #eee}
h3{color:#22223b}
img{max-width:100%;border-radius:8px} figure{margin:1.2em 0} figcaption{color:#666;font-size:0.9em;text-align:center}
table{border-collapse:collapse;width:100%;margin:1.2em 0} th,td{border:1px solid #ddd;padding:8px 12px;text-align:left} th{background:#f5f5f5}
blockquote{border-left:4px solid #ccc;margin:1.2em 0;padding:0.5em 1em;color:#555}
hr{border:none;border-top:1px solid #eee;margin:30px 0}
""" + base_cards
    # default(기사형 단일단)
    return """
body{max-width:760px;margin:40px auto;padding:0 20px;font-family:'Pretendard','Noto Sans KR',sans-serif;line-height:1.7;color:#1a1a1a}
.docwrap{}
h1,h2,h3,h4{line-height:1.3;margin:1.4em 0 0.5em} h1{font-size:2em}
.doc-meta{color:#666;font-size:0.95em;margin:-0.5em 0 1.5em}
img{max-width:100%;border-radius:8px} figure{margin:1.5em 0} figcaption{color:#666;font-size:0.9em;text-align:center;margin-top:0.5em}
table{border-collapse:collapse;width:100%;margin:1.5em 0} th,td{border:1px solid #ddd;padding:8px 12px;text-align:left} th{background:#f5f5f5}
blockquote{border-left:4px solid #ccc;margin:1.5em 0;padding:0.5em 1em;color:#555} blockquote cite{display:block;margin-top:0.5em;font-size:0.9em}
pre{background:#f6f8fa;padding:1em;border-radius:6px;overflow:auto} hr{border:none;border-top:1px solid #e0e0e0;margin:2em 0}
""" + base_cards


def _resolve_image_bytes(src: str):
    """이미지 src(로컬 경로/data URI/http URL)를 file-like(BytesIO)로 해소. 실패 시 None.
    docx·pptx emitter 공용 — 둘 다 파일/스트림만 받음(HTML <img>와 달리)."""
    import io
    import os
    import base64
    if not src:
        return None
    try:
        s = str(src).strip()
        if s.startswith("data:"):  # data:image/png;base64,....
            b64 = s.split(",", 1)[1]
            return io.BytesIO(base64.b64decode(b64))
        if s.startswith("file://"):
            s = s[7:]
        if s.startswith(("http://", "https://")):
            import urllib.request
            with urllib.request.urlopen(s, timeout=15) as r:
                return io.BytesIO(r.read())
        if os.path.isfile(s):
            with open(s, "rb") as f:
                return io.BytesIO(f.read())
    except Exception:
        return None
    return None


def _typ_esc(s) -> str:
    """typst 마크업 특수문자 이스케이프."""
    s = str(s if s is not None else "")
    for ch in ("\\", "#", "$", "*", "_", "`", "<", ">", "@", "[", "]"):
        s = s.replace(ch, "\\" + ch)
    return s


def _doc_blocks_to_typst(blocks: list, title: str, meta: str, out_path: str):
    """문서 IR → typst 컴파일 PDF (책 품질 조판). 산문·보고서에 최적 — 정렬·페이지·타이포가 강점.
    HTML theme/cards 그리드는 무시(조판 모델이 다름). 한글 = Apple SD Gothic Neo."""
    import os
    import subprocess
    import tempfile

    lines = [
        '#set text(font: "Apple SD Gothic Neo", size: 11pt, lang: "ko")',
        '#set page(paper: "a4", margin: (x: 2.2cm, y: 2.4cm), numbering: "1")',
        '#set par(justify: true, leading: 0.8em)',
        '#show heading: set block(above: 1.2em, below: 0.6em)',
        '#set heading(numbering: none)',
        "",
    ]
    if title:
        lines.append(f'#align(center)[#text(size: 22pt, weight: "bold")[{_typ_esc(title)}]]')
    if meta:
        lines.append(f'#align(center)[#text(size: 9pt, fill: gray)[{_typ_esc(meta)}]]')
    if title or meta:
        lines.append("#v(0.5em)")
        lines.append("")

    for b in blocks:
        if not isinstance(b, dict):
            continue
        t = (b.get("type") or "paragraph").lower()
        if t == "heading":
            lvl = max(1, min(5, int(b.get("level") or 2)))
            lines.append("=" * lvl + " " + _typ_esc(b.get("text")))
        elif t == "list":
            for it in (b.get("items") or []):
                txt = it.get("text") if isinstance(it, dict) else it
                lines.append("- " + _typ_esc(txt))
        elif t == "table":
            cols = b.get("columns") or []
            rows = [r for r in (b.get("rows") or []) if isinstance(r, (list, tuple))]
            ncol = max([len(cols)] + [len(r) for r in rows] or [0])
            if ncol:
                cells = []
                if cols:
                    cells += [f"[*{_typ_esc(c)}*]" for c in cols[:ncol]] + ["[]"] * (ncol - len(cols[:ncol]))
                for r in rows:
                    cells += [f"[{_typ_esc(v)}]" for v in r[:ncol]] + ["[]"] * (ncol - len(r[:ncol]))
                lines.append(f"#table(columns: {ncol}, " + ", ".join(cells) + ")")
        elif t == "cards":
            for it in (b.get("items") or []):
                if not isinstance(it, dict):
                    continue
                lines.append("=== " + _typ_esc(it.get("title")))
                if it.get("meta"):
                    lines.append(f'#text(size: 9pt, fill: gray)[{_typ_esc(it.get("meta"))}]')
                if it.get("summary"):
                    lines.append(_typ_esc(it.get("summary")))
                if it.get("url"):
                    lines.append(f'#link("{it.get("url")}")[{_typ_esc(it.get("link_label") or "열기")}]')
        elif t == "quote":
            cite = b.get("cite")
            q = f'#quote(block: true)[{_typ_esc(b.get("text"))}]'
            lines.append(q + (f" #text(size: 9pt, fill: gray)[— {_typ_esc(cite)}]" if cite else ""))
        elif t == "code":
            lines.append("```\n" + str(b.get("text") or "") + "\n```")
        elif t == "image":
            src = b.get("src") or b.get("path") or ""
            if src and os.path.isfile(str(src)):
                lines.append(f'#figure(image("{src}", width: 80%)' +
                             (f', caption: [{_typ_esc(b.get("caption"))}]' if b.get("caption") else "") + ")")
        elif t == "divider":
            lines.append("#line(length: 100%, stroke: 0.5pt + gray)")
        else:
            lines.append(_typ_esc(b.get("text")))
        lines.append("")

    typ_src = "\n".join(lines)
    with tempfile.NamedTemporaryFile("w", suffix=".typ", delete=False, encoding="utf-8") as tf:
        tf.write(typ_src)
        typ_file = tf.name
    try:
        proc = subprocess.run(["typst", "compile", typ_file, out_path],
                              capture_output=True, text=True, timeout=60)
        if proc.returncode != 0:
            raise RuntimeError(f"typst compile 실패: {proc.stderr[:300]}")
    finally:
        try:
            os.unlink(typ_file)
        except Exception:
            pass


def _add_hyperlink(paragraph, url: str, text: str):
    """python-docx 문단에 클릭 가능한 하이퍼링크 추가(네이티브 미지원이라 관계+XML 수작업)."""
    from docx.oxml.ns import qn
    from docx.oxml import OxmlElement
    try:
        part = paragraph.part
        r_id = part.relate_to(url, "http://schemas.openxmlformats.org/officeDocument/2006/relationships/hyperlink",
                              is_external=True)
        hyperlink = OxmlElement("w:hyperlink")
        hyperlink.set(qn("r:id"), r_id)
        new_run = OxmlElement("w:r")
        rPr = OxmlElement("w:rPr")
        color = OxmlElement("w:color"); color.set(qn("w:val"), "3D5A80"); rPr.append(color)
        u = OxmlElement("w:u"); u.set(qn("w:val"), "single"); rPr.append(u)
        new_run.append(rPr)
        t = OxmlElement("w:t"); t.text = text or url; new_run.append(t)
        hyperlink.append(new_run)
        paragraph._p.append(hyperlink)
    except Exception:
        paragraph.add_run(f"{text}: {url}" if text else url)


def _doc_blocks_to_docx(blocks: list, title: str, out_path: str):
    """문서 IR → .docx (python-docx). html emitter와 같은 IR을 소비.
    table 블록 = 데이터 통화 {columns,rows} 그대로 재사용."""
    from docx import Document
    from docx.shared import Pt, Inches
    from docx.oxml.ns import qn
    from docx.oxml import OxmlElement

    doc = Document()
    if title:
        doc.add_heading(str(title), level=0)

    for b in blocks:
        if not isinstance(b, dict):
            continue
        t = (b.get("type") or "paragraph").lower()
        if t == "heading":
            lvl = max(1, min(6, int(b.get("level") or 2)))
            doc.add_heading(str(b.get("text") or ""), level=lvl)
        elif t == "list":
            ordered = bool(b.get("ordered"))
            style = "List Number" if ordered else "List Bullet"
            for it in (b.get("items") or []):
                txt = it.get("text") if isinstance(it, dict) else it
                doc.add_paragraph(str(txt), style=style)
        elif t == "image":
            stream = _resolve_image_bytes(b.get("src") or b.get("path") or "")
            if stream is not None:
                try:
                    doc.add_picture(stream, width=Inches(6.0))
                except Exception:
                    pass
            cap = b.get("caption")
            if cap:
                p = doc.add_paragraph(str(cap))
                p.italic = True
        elif t == "table":
            cols = b.get("columns") or []
            rows = [r for r in (b.get("rows") or []) if isinstance(r, (list, tuple))]
            ncol = max([len(cols)] + [len(r) for r in rows] or [0])
            if ncol:
                tbl = doc.add_table(rows=0, cols=ncol)
                try:
                    tbl.style = "Light Grid Accent 1"
                except Exception:
                    pass
                if cols:
                    hdr = tbl.add_row().cells
                    for i, c in enumerate(cols[:ncol]):
                        hdr[i].text = str(c)
                    for cell in hdr:
                        for para in cell.paragraphs:
                            for run in para.runs:
                                run.bold = True
                for r in rows:
                    cells = tbl.add_row().cells
                    for i, v in enumerate(r[:ncol]):
                        cells[i].text = str(v)
        elif t == "quote":
            p = doc.add_paragraph(str(b.get("text") or ""))
            try:
                p.style = "Intense Quote"
            except Exception:
                pass
            cite = b.get("cite")
            if cite:
                doc.add_paragraph(f"— {cite}")
        elif t == "code":
            p = doc.add_paragraph()
            run = p.add_run(str(b.get("text") or ""))
            run.font.name = "Courier New"
            run.font.size = Pt(9.5)
        elif t == "cards":
            for it in (b.get("items") or []):
                if not isinstance(it, dict):
                    continue
                if it.get("image"):  # 썸네일(표지 등) — 다운로드 실패는 graceful
                    stream = _resolve_image_bytes(it.get("image"))
                    if stream is not None:
                        try:
                            doc.add_picture(stream, width=Inches(1.6))
                        except Exception:
                            pass
                doc.add_heading(str(it.get("title") or ""), level=3)
                if it.get("meta"):
                    mp = doc.add_paragraph()
                    mr = mp.add_run(str(it.get("meta")))
                    mr.italic = True
                    mr.font.size = Pt(9)
                if it.get("summary"):
                    doc.add_paragraph(str(it.get("summary")))
                if it.get("url"):
                    _add_hyperlink(doc.add_paragraph(), str(it.get("url")),
                                   str(it.get("link_label") or "열기"))
        elif t == "divider":
            p = doc.add_paragraph()
            pPr = p._p.get_or_add_pPr()
            pbdr = OxmlElement("w:pBdr")
            bottom = OxmlElement("w:bottom")
            bottom.set(qn("w:val"), "single")
            bottom.set(qn("w:sz"), "6")
            bottom.set(qn("w:space"), "1")
            bottom.set(qn("w:color"), "auto")
            pbdr.append(bottom)
            pPr.append(pbdr)
        else:  # paragraph
            doc.add_paragraph(str(b.get("text") or ""))

    doc.save(out_path)


def _doc_blocks_to_pptx(blocks: list, title: str, out_path: str):
    """문서 IR → .pptx (python-pptx). ★종류 경계 주의: 슬라이드 IR이 아니라 *문서 IR을 슬라이드로 투영*.
    문서 IR이 정본, pptx는 emitter일 뿐 — heading(level≤2)이 새 슬라이드, 그 아래 내용이 글머리표.
    슬라이드 전용 시각 레이아웃이 필요하면 engines:slide(슬라이드 IR)를 써야지 이걸 쓰면 안 됨."""
    from pptx import Presentation
    from pptx.util import Inches, Pt

    prs = Presentation()
    SW, SH = prs.slide_width, prs.slide_height
    blank = prs.slide_layouts[6]
    title_layout = prs.slide_layouts[0]

    # 표지 슬라이드(title 있으면)
    if title:
        s = prs.slides.add_slide(title_layout)
        s.shapes.title.text = str(title)

    state = {"slide": None, "body": None}

    def new_content_slide(heading_text=""):
        # "제목+내용" 레이아웃(1) — 제목 placeholder + 본문 placeholder
        s = prs.slides.add_slide(prs.slide_layouts[1])
        s.shapes.title.text = str(heading_text or "")
        body_tf = None
        for ph in s.placeholders:
            if ph.placeholder_format.idx == 1:
                body_tf = ph.text_frame
                break
        if body_tf is not None:
            body_tf.clear()
            body_tf.word_wrap = True
        state["slide"], state["body"] = s, body_tf

    def add_bullet(text, level=0, italic=False, mono=False):
        if state["body"] is None:
            new_content_slide("")
        tf = state["body"]
        # clear()가 남긴 빈 첫 문단 재사용, 이후엔 add
        if len(tf.paragraphs) == 1 and not tf.paragraphs[0].runs:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = str(text)
        p.level = min(4, max(0, level))
        for run in p.runs:
            if italic:
                run.font.italic = True
            if mono:
                run.font.name = "Courier New"
                run.font.size = Pt(14)

    for b in blocks:
        if not isinstance(b, dict):
            continue
        t = (b.get("type") or "paragraph").lower()
        if t == "heading":
            lvl = int(b.get("level") or 2)
            if lvl <= 2 or state["body"] is None:
                # 큰 섹션, 또는 표·구분선 뒤 첫 제목이면 새 슬라이드 제목으로
                new_content_slide(b.get("text") or "")
            else:  # 진행 중 슬라이드의 하위 섹션 = 글머리표
                add_bullet(b.get("text") or "", level=0)
        elif t == "paragraph":
            add_bullet(b.get("text") or "", level=0)
        elif t == "list":
            for it in (b.get("items") or []):
                add_bullet(it.get("text") if isinstance(it, dict) else it, level=1)
        elif t == "quote":
            txt = str(b.get("text") or "")
            cite = b.get("cite")
            add_bullet(f"“{txt}”" + (f" — {cite}" if cite else ""), level=1, italic=True)
        elif t == "code":
            add_bullet(b.get("text") or "", level=1, mono=True)
        elif t == "cards":
            for it in (b.get("items") or []):
                if not isinstance(it, dict):
                    continue
                add_bullet(it.get("title") or "", level=0)
                sub = " / ".join(x for x in [it.get("meta"), it.get("summary")] if x)
                if sub:
                    add_bullet(sub, level=1)
        elif t == "image":
            stream = _resolve_image_bytes(b.get("src") or b.get("path") or "")
            if stream is not None:
                s = prs.slides.add_slide(blank)
                try:
                    s.shapes.add_picture(stream, Inches(0.6), Inches(0.6), width=SW - Inches(1.2))
                except Exception:
                    pass
                cap = b.get("caption")
                if cap:
                    tb = s.shapes.add_textbox(Inches(0.6), SH - Inches(0.9), SW - Inches(1.2), Inches(0.7))
                    tb.text_frame.text = str(cap)
                state["slide"], state["body"] = None, None  # 이미지 후 새 내용 슬라이드 강제
        elif t == "table":
            cols = b.get("columns") or []
            rows = [r for r in (b.get("rows") or []) if isinstance(r, (list, tuple))]
            ncol = max([len(cols)] + [len(r) for r in rows] or [0])
            nrow = len(rows) + (1 if cols else 0)
            if ncol and nrow:
                s = prs.slides.add_slide(blank)
                gt = s.shapes.add_table(nrow, ncol, Inches(0.5), Inches(0.6),
                                        SW - Inches(1.0), Inches(0.4) * nrow).table
                ri = 0
                if cols:
                    for ci, c in enumerate(cols[:ncol]):
                        gt.cell(0, ci).text = str(c)
                    ri = 1
                for r in rows:
                    for ci, v in enumerate(r[:ncol]):
                        gt.cell(ri, ci).text = str(v)
                    ri += 1
                state["slide"], state["body"] = None, None
        elif t == "divider":
            state["slide"], state["body"] = None, None  # 새 슬라이드 경계

    if len(prs.slides) == 0:  # 표지도 내용도 없으면 빈 슬라이드 하나
        prs.slides.add_slide(blank)
    prs.save(out_path)


# ── B: 구조화 원자 — 콘텐츠 → 문서 IR (A획득→B구조화→IR→emit 파이프라인) ──
_STRUCTURE_PROMPT = """당신은 콘텐츠를 깔끔한 문서 구조로 정리하는 편집자입니다. 주어진 내용을 문서 IR(JSON)로 변환합니다.

출력은 JSON 한 객체만: {"title": "...", "blocks": [ ... ]}
블록 타입:
- {"type":"heading","level":2,"text":"..."}   (level 1~4)
- {"type":"paragraph","text":"..."}
- {"type":"list","ordered":false,"items":["...","..."]}
- {"type":"table","columns":["...","..."],"rows":[["...","..."]]}
- {"type":"quote","text":"...","cite":"..."}
- {"type":"code","text":"...","lang":"..."}
- {"type":"divider"}

원칙: 내용을 지어내지 말고 주어진 것에서만. title=핵심을 담은 명제. 긴 글은 heading으로 섹션화, 나열은 list, 비교·수치는 table. JSON 외 텍스트 금지."""


def structure_document(tool_input, output_base="."):
    """[table:structure] — 원본 콘텐츠를 문서 IR(blocks)로 구조화 (LLM 편집자).

    파라미터: content(필수, 원본 텍스트) · instruction(선택, 정리 방향).
    반환: {success, title, blocks, block_count}. render_document로 이어 렌더(>> 파이프 지원).
    """
    import json as _json

    content = (tool_input.get("content") or "").strip()
    # >> 파이프: 이전 액션의 텍스트 결과를 content로 받음
    if not content:
        pr = tool_input.get("_prev_result")
        if isinstance(pr, str):
            content = pr.strip()
        elif isinstance(pr, dict):
            content = str(pr.get("summary") or pr.get("content") or pr.get("text") or "").strip()
    if not content:
        return _json.dumps({"success": False, "message": "content(구조화할 원본 내용)가 필요합니다."},
                           ensure_ascii=False)

    instruction = (tool_input.get("instruction") or "").strip()
    user = f"# 정리할 내용\n{content[:16000]}"
    if instruction:
        user += f"\n\n# 정리 방향\n{instruction}"
    user += "\n\n위 내용을 문서 IR(JSON 한 객체)로 출력하라."

    try:
        from consciousness_agent import lightweight_ai_call
        resp = lightweight_ai_call(user, system_prompt=_STRUCTURE_PROMPT)
    except Exception as e:
        return _json.dumps({"success": False, "message": f"구조화 AI 호출 실패: {e}"}, ensure_ascii=False)
    if not resp or not resp.strip():
        return _json.dumps({"success": False, "message": "구조화 AI 응답 없음"}, ensure_ascii=False)

    txt = resp.strip()
    if txt.startswith("```"):
        txt = txt.strip("`")
        if txt.lower().startswith("json"):
            txt = txt[4:]
        txt = txt.split("```")[0].strip()
    try:
        a, b = txt.find("{"), txt.rfind("}")
        ir = _json.loads(txt[a:b + 1])
    except Exception as e:
        return _json.dumps({"success": False, "message": f"IR JSON 파싱 실패: {e}", "raw": resp[:300]},
                           ensure_ascii=False)
    blocks = ir.get("blocks")
    if not isinstance(blocks, list) or not blocks:
        return _json.dumps({"success": False, "message": "blocks가 없습니다.", "raw": resp[:300]},
                           ensure_ascii=False)
    return _json.dumps({"success": True, "title": ir.get("title", ""), "blocks": blocks,
                        "block_count": len(blocks),
                        "message": f"{len(blocks)}블록 문서 IR로 구조화."}, ensure_ascii=False)


def render_document(tool_input, output_base="."):
    """문서 IR → 산출물. 현재 emitter: html (단일 IR, 향후 pdf/docx/pptx emitter 추가).

    파라미터: blocks(필수, IR 블록 배열) · title(선택) · format(기본 html) · filename(선택).
    반환: {success, path, format, blocks}.
    """
    import os
    import html as _html
    import json as _json

    blocks = tool_input.get("blocks")
    if not blocks:
        # >> 파이프: 이전 생산자 결과(_prev_result)의 blocks·title·meta·theme 자동 수용
        pr = tool_input.get("_prev_result")
        if pr:
            try:
                po = _json.loads(pr) if isinstance(pr, str) else pr
                if isinstance(po, dict):
                    _rows = po.get("items")
                    if po.get("blocks"):
                        blocks = po["blocks"]
                    elif isinstance(_rows, list) and _rows:
                        if isinstance(_rows[0], dict) and "type" in _rows[0] and "text" in _rows[0]:
                            # 문서 IR items(type+text — crawl·read 등) = blocks 그 자체(산문).
                            blocks = _rows
                        else:
                            # 단일 통화 items([{title,meta,summary,url,image}]) → cards 블록으로 래핑.
                            blocks = [{"type": "cards", "columns": 2, "items": _rows}]
                    for k in ("title", "meta", "theme"):
                        if not tool_input.get(k) and po.get(k):
                            tool_input[k] = po[k]
            except Exception:
                pass
    if isinstance(blocks, str):
        try:
            blocks = _json.loads(blocks)
        except Exception:
            blocks = None
    if not isinstance(blocks, list) or not blocks:
        return _json.dumps({"success": False, "message": "blocks(문서 IR 블록 배열)가 필요합니다."},
                           ensure_ascii=False)

    title = tool_input.get("title") or ""
    meta = tool_input.get("meta") or ""
    theme = (tool_input.get("theme") or "default").strip().lower()
    fmt = (tool_input.get("format") or "html").strip().lower()
    note = ""
    if fmt not in ("html", "pdf", "png", "docx", "pptx", "typst"):
        note = f" (format '{fmt}' 미지원 — html로 산출)"
        fmt = "html"

    os.makedirs(output_base, exist_ok=True)
    base = tool_input.get("filename") or "document"
    base = os.path.splitext(os.path.basename(str(base)))[0] or "document"

    # typst emitter — 책 품질 조판 PDF(산문·보고서). HTML theme/cards 그리드는 무시(조판 모델 상이).
    if fmt == "typst":
        try:
            out_path = os.path.join(output_base, f"{base}.pdf")
            _doc_blocks_to_typst(blocks, title, meta, out_path)
            return _json.dumps({"success": True, "path": out_path, "file": out_path,
                                "title": title, "format": "typst_pdf", "blocks": len(blocks),
                                "message": f"문서 {len(blocks)}블록을 typst 책 품질 PDF로 조판했습니다."},
                               ensure_ascii=False)
        except Exception as e:
            note = f" (typst 조판 실패 → 브라우저 PDF 폴백: {e})"
            fmt = "pdf"

    # docx/pptx emitter — 같은 문서 IR을 사무 포맷으로. pptx는 문서 IR을 슬라이드로 *투영*(종류 경계 주의).
    if fmt in ("docx", "pptx"):
        try:
            out_path = os.path.join(output_base, f"{base}.{fmt}")
            if fmt == "docx":
                _doc_blocks_to_docx(blocks, title, out_path)
            else:
                _doc_blocks_to_pptx(blocks, title, out_path)
            return _json.dumps({"success": True, "path": out_path, "file": out_path,
                                "title": title, "format": fmt, "blocks": len(blocks),
                                "message": f"문서 {len(blocks)}블록을 {fmt.upper()}로 렌더했습니다."},
                               ensure_ascii=False)
        except Exception as e:
            note = f" ({fmt} 렌더 실패 → HTML 폴백: {e})"
            fmt = "html"

    body = _doc_blocks_to_html(blocks)
    title_h = f"<h1>{_html.escape(str(title))}</h1>" if title else ""
    meta_h = f'<div class="doc-meta">{_html.escape(str(meta))}</div>' if meta else ""
    doc = f"""<!DOCTYPE html><html lang="ko"><head><meta charset="utf-8">
<title>{_html.escape(str(title))}</title><style>{_doc_css(theme)}</style></head><body>
<div class="docwrap">
{title_h}
{meta_h}
{body}
</div>
</body></html>"""

    # 같은 문서 IR을 여러 emitter로 — html/pdf/png. pdf·png는 Playwright로 동일 HTML 렌더.
    if fmt in ("pdf", "png"):
        try:
            from playwright.sync_api import sync_playwright
            out_path = os.path.join(output_base, f"{base}.{fmt}")
            with sync_playwright() as pw:
                br = pw.chromium.launch()
                pg = br.new_page(viewport={"width": 900, "height": 1200})
                pg.set_content(doc, wait_until="networkidle")
                pg.wait_for_timeout(300)
                if fmt == "pdf":
                    pg.pdf(path=out_path, format="A4", print_background=True,
                           margin={"top": "20mm", "bottom": "20mm", "left": "16mm", "right": "16mm"})
                else:
                    pg.screenshot(path=out_path, full_page=True)
                br.close()
            return _json.dumps({"success": True, "path": out_path, "file": out_path,
                                "title": title, "format": fmt, "blocks": len(blocks),
                                "message": f"문서 {len(blocks)}블록을 {fmt.upper()}로 렌더했습니다."},
                               ensure_ascii=False)
        except Exception as e:
            # emitter 실패 시 HTML로 폴백(산출 보존)
            note = f" ({fmt} 렌더 실패 → HTML 폴백: {e})"

    out_path = os.path.join(output_base, f"{base}.html")
    with open(out_path, "w", encoding="utf-8") as f:
        f.write(doc)
    return _json.dumps({"success": True, "path": out_path, "file": out_path,
                        "title": title, "format": "html", "blocks": len(blocks),
                        # 렌더된 HTML을 결과에 동봉 — 액션이 다른 몸(맥)으로 포워드돼 파일이
                        # 거기 생겨도, 호출한 몸(폰)이 파일 위치 의존 없이 콘텐츠로 바로 띄운다.
                        "html": doc,
                        "message": f"문서 {len(blocks)}블록을 HTML로 렌더했습니다.{note}"},
                       ensure_ascii=False)


_DISPATCH = {
    "data_filter": _op_filter,
    "data_sort": _op_sort,
    "data_take": _op_take,
    "data_select": _op_select,
    "data_dedup": _op_dedup,
    "data_groupby": _op_groupby,
    "data_join": _op_join,
    "data_union": _op_union,
    "data_merge": _op_merge,
}


def execute(tool_input: dict, context):
    """표준 시그니처. context.tool_name 으로 동사 분기, _prev_result에서 통화 수용."""
    tool_name = getattr(context, "tool_name", None)
    # 표준 코어 emitter(table:structure/document) — 변환자와 달리 산출 경로(output_dir) 사용.
    if tool_name == "structure_document":
        return structure_document(tool_input, context.output_dir())
    if tool_name == "render_document":
        return render_document(tool_input, context.output_dir())
    fn = _DISPATCH.get(tool_name)
    if not fn:
        return {"error": f"data-ops: 알 수 없는 변환자 '{tool_name}'."}
    params = dict(tool_input or {})
    prev = _parse_prev(params.get("_prev_result"))
    if prev is None:
        # 파이프 입력(>>)이 없으면 params 에서 통화를 직접 수용 — 단독 호출/자가점검 지원.
        # (파이프 통화와 params 통화는 같은 모양이라 정합적이다.)
        # 이항(merge/join/union): (left,right)/(table1,table2)/(a,b) 쌍 또는 inputs 리스트 → [A, B].
        # 단항(filter/sort/take/select/dedup/groupby): items(단일 통화) 또는 table(표형).
        for k1, k2 in (("left", "right"), ("table1", "table2"), ("a", "b")):
            if params.get(k1) is not None and params.get(k2) is not None:
                prev = [params[k1], params[k2]]
                break
        if prev is None:
            if isinstance(params.get("inputs"), list):
                prev = params["inputs"]
            elif params.get("items") is not None:
                prev = {"items": params["items"]}
            elif params.get("table") is not None:
                prev = {"table": params["table"]}
    if prev is None:
        return {"error": (
            f"{tool_name}: 입력 통화가 없습니다. 변환자는 >> 파이프로 앞 액션의 "
            "items 통화(표형은 table)를 받습니다. 예: [sense:search]{...} >> [table:filter]{where:...}"
        )}
    return fn(prev, params)
