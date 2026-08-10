"""
lecture_export.py - 강의 데크를 PDF / PPTX로 내보내기

저장 위치: {lecture_dir}/exports/{lecture_id}_{timestamp}.{pdf,pptx}

설계 원칙:
- 데크의 slide_order 순서대로 PNG를 모아서 합본
- PDF: PIL의 다중 페이지 저장 (의존성 추가 없음)
- PPTX: python-pptx로 1280×720 슬라이드에 PNG를 전체 배경으로 삽입
- 슬라이드 PNG가 없는 항목은 건너뜀 (경고)
"""

from __future__ import annotations

import sys
from datetime import datetime
from pathlib import Path
from typing import Optional

# lecture_store import
_THIS_DIR = Path(__file__).resolve().parent
if str(_THIS_DIR) not in sys.path:
    sys.path.insert(0, str(_THIS_DIR))
import lecture_store  # noqa: E402


def _collect_slide_pngs(deck: dict, lecture_dir: Path) -> list[Path]:
    """slide_order 순서대로 존재하는 PNG 경로 모음."""
    paths = []
    for sid in deck.get("slide_order", []):
        meta = deck.get("slides", {}).get(sid, {})
        rel = meta.get("png_file")
        if not rel:
            continue
        p = lecture_dir / rel
        if p.exists():
            paths.append(p)
    return paths


def _exports_dir(lecture_id: str) -> Path:
    d = lecture_store.lecture_dir(lecture_id) / "exports"
    d.mkdir(parents=True, exist_ok=True)
    return d


def _timestamp() -> str:
    return datetime.now().strftime("%Y%m%d_%H%M%S")


# ─────────────────────────────────────────────────────────────────────
# PDF
# ─────────────────────────────────────────────────────────────────────

def export_pdf(lecture_id: str, output_path: Optional[Path] = None) -> dict:
    """슬라이드 데크를 PDF 한 파일로. PIL의 PDF 저장 기능 사용.

    Returns: {success, path, slide_count, skipped, format: "pdf"}
    """
    from PIL import Image

    deck = lecture_store.read_deck(lecture_id)
    lecture_dir = lecture_store.lecture_dir(lecture_id)
    pngs = _collect_slide_pngs(deck, lecture_dir)
    total_planned = len(deck.get("slide_order", []))
    skipped = total_planned - len(pngs)

    if not pngs:
        raise ValueError("내보낼 슬라이드가 없습니다 (PNG 없음)")

    if output_path is None:
        output_path = _exports_dir(lecture_id) / f"{lecture_id}_{_timestamp()}.pdf"

    # RGB로 변환 (PDF는 알파 채널 미지원)
    images = [Image.open(p).convert("RGB") for p in pngs]
    images[0].save(
        str(output_path),
        "PDF",
        save_all=True,
        append_images=images[1:] if len(images) > 1 else [],
        resolution=150.0,
    )

    return {
        "success": True,
        "format": "pdf",
        "path": str(output_path.resolve()),
        "slide_count": len(pngs),
        "skipped": skipped,
        "filename": output_path.name,
    }


# ─────────────────────────────────────────────────────────────────────
# PPTX (이미지 모드 — 디자인 완벽 보존, 편집 불가)
# ─────────────────────────────────────────────────────────────────────

def export_pptx(lecture_id: str, output_path: Optional[Path] = None) -> dict:
    """슬라이드 데크를 PPTX로 (이미지 모드).

    각 슬라이드 PNG를 통째 배경에 박아넣음. 디자인 완벽 보존하나 PPT에서 편집 불가.
    """
    try:
        from pptx import Presentation
        from pptx.util import Emu
    except ImportError:
        raise RuntimeError(
            "python-pptx 미설치. 'pip install python-pptx' 후 재시도하세요."
        )

    deck = lecture_store.read_deck(lecture_id)
    lecture_dir = lecture_store.lecture_dir(lecture_id)
    pngs = _collect_slide_pngs(deck, lecture_dir)
    total_planned = len(deck.get("slide_order", []))
    skipped = total_planned - len(pngs)

    if not pngs:
        raise ValueError("내보낼 슬라이드가 없습니다 (PNG 없음)")

    if output_path is None:
        output_path = _exports_dir(lecture_id) / f"{lecture_id}_{_timestamp()}.pptx"

    prs = Presentation()
    # 16:9 슬라이드 — 1280×720 픽셀에 가깝게.
    # PPT 표준 widescreen은 13.333" × 7.5" = 12,192,000 × 6,858,000 EMU
    prs.slide_width = Emu(12192000)
    prs.slide_height = Emu(6858000)
    blank_layout = prs.slide_layouts[6]  # blank

    for png_path in pngs:
        slide = prs.slides.add_slide(blank_layout)
        slide.shapes.add_picture(
            str(png_path), 0, 0,
            width=prs.slide_width,
            height=prs.slide_height,
        )

    prs.save(str(output_path))

    return {
        "success": True,
        "format": "pptx",
        "mode": "image",
        "path": str(output_path.resolve()),
        "slide_count": len(pngs),
        "skipped": skipped,
        "filename": output_path.name,
    }


# ─────────────────────────────────────────────────────────────────────
# PPTX (편집 가능 모드 — 텍스트박스로 분해)
# ─────────────────────────────────────────────────────────────────────
#
# 슬라이드 spec을 파싱해서 PPT 텍스트박스 + 이미지 객체로 분해.
# PPT에서 자유 위치 조정, 폰트 변경, 텍스트 편집 모두 가능.
# 디자인 시스템의 폰트·텍스처는 PPT가 재현 못 함 — 톤은 단순화됨.
#
# 좌표: 1280×720 px = 12192000×6858000 EMU. 1px = 9525 EMU.

# 1280×720 기준 위치 표 (px). 함수에서 EMU로 변환.
PX_TO_EMU = 9525
SLIDE_W_PX = 1280
SLIDE_H_PX = 720


def export_pptx_editable(lecture_id: str, output_path: Optional[Path] = None) -> dict:
    """슬라이드 spec → 편집 가능 PPTX. 텍스트박스 + 이미지로 분해.

    Layout별로 약속된 위치에 텍스트박스 배치. 일러스트 이미지는 별도 picture로 삽입.
    PPT에서 자유롭게 위치 조정·텍스트 편집 가능.
    """
    try:
        from pptx import Presentation
        from pptx.util import Emu, Pt
    except ImportError:
        raise RuntimeError(
            "python-pptx 미설치. 'pip install python-pptx' 후 재시도하세요."
        )

    deck = lecture_store.read_deck(lecture_id)
    lecture_dir = lecture_store.lecture_dir(lecture_id)

    slide_order = deck.get("slide_order", [])
    if not slide_order:
        raise ValueError("내보낼 슬라이드가 없습니다.")

    if output_path is None:
        output_path = _exports_dir(lecture_id) / f"{lecture_id}_{_timestamp()}_editable.pptx"

    prs = Presentation()
    prs.slide_width = Emu(SLIDE_W_PX * PX_TO_EMU)
    prs.slide_height = Emu(SLIDE_H_PX * PX_TO_EMU)
    blank_layout = prs.slide_layouts[6]

    placed = 0
    fallback_image = 0  # spec 파싱 실패 시 이미지로 fallback한 수
    overlay_boxes = 0  # '글자 얹기' 슬라이드 — 원본 그림 + 텍스트박스로 분해된 수

    for sid in slide_order:
        meta = deck.get("slides", {}).get(sid)
        if not meta:
            continue
        slide = prs.slides.add_slide(blank_layout)

        # 글자가 PNG에 구워진 슬라이드(native 통짜 / composite 합성 / image 업로드)는
        # 분해 불가 — 편집모드 내보내기에서도 이미지 그대로 얹어 비주얼을 보존한다.
        # 단 '글자 얹기'(text_overlays)가 있으면 원본 그림 + 진짜 텍스트박스로 분해:
        # 얹은 글자는 구조화된 데이터(문구·위치·크기·색)가 살아 있어 PPT에서 편집 가능하게 나간다.
        if meta.get("layout") in ("native", "composite", "image"):
            overlays = meta.get("text_overlays") or []
            base_png = lecture_dir / "slides" / f"{sid}.base.png"
            if overlays and base_png.exists():
                slide.shapes.add_picture(
                    str(base_png), 0, 0, width=prs.slide_width, height=prs.slide_height
                )
                for ov in overlays:
                    _add_overlay_textbox(slide, ov, Emu, Pt)
                overlay_boxes += 1
            else:
                png = lecture_dir / meta.get("png_file", "")
                if png.exists():
                    slide.shapes.add_picture(
                        str(png), 0, 0, width=prs.slide_width, height=prs.slide_height
                    )
                    fallback_image += 1
            continue

        spec_file = lecture_dir / meta.get("spec_file", "")
        try:
            if not spec_file.exists():
                raise FileNotFoundError(f"spec 없음: {spec_file}")
            import json
            with open(spec_file, "r", encoding="utf-8") as f:
                spec = json.load(f)
            _populate_editable_slide(slide, spec, lecture_dir, Emu, Pt)
            placed += 1
        except Exception as e:
            # spec 파싱 실패 → PNG로 fallback
            print(f"[pptx_editable] {sid} fallback to image: {e}")
            png = lecture_dir / meta.get("png_file", "")
            if png.exists():
                slide.shapes.add_picture(
                    str(png), 0, 0,
                    width=prs.slide_width,
                    height=prs.slide_height,
                )
                fallback_image += 1

    prs.save(str(output_path))

    return {
        "success": True,
        "format": "pptx",
        "mode": "editable",
        "path": str(output_path.resolve()),
        "slide_count": placed + fallback_image + overlay_boxes,
        "editable_count": placed,
        "fallback_image_count": fallback_image,
        # '글자 얹기' 슬라이드 — 원본 그림 + 편집 가능 텍스트박스로 나간 수
        "overlay_textbox_count": overlay_boxes,
        "filename": output_path.name,
    }


def _populate_editable_slide(slide, spec: dict, lecture_dir: Path, Emu, Pt):
    """단일 슬라이드에 layout별 텍스트박스·이미지 배치."""
    layout = spec.get("layout", "lecture_body")

    if layout in ("hero", "hero_illustration"):
        _add_image_if_any(slide, spec, "image_path", lecture_dir, 320, 100, 640, 360, Emu)
        _add_text(slide, spec.get("eyebrow"), 80, 60, 1120, 40, Emu, Pt, font_size=14, bold=False, gray=True)
        _add_text(slide, spec.get("title"), 80, 480, 1120, 100, Emu, Pt, font_size=48, bold=True, align="center")
        _add_text(slide, spec.get("subtitle"), 80, 590, 1120, 60, Emu, Pt, font_size=24, align="center", gray=True)

    elif layout == "quote":
        _add_text(slide, spec.get("quote"), 100, 180, 1080, 300, Emu, Pt, font_size=44, italic=True, align="center")
        _add_text(slide, spec.get("attribution"), 100, 500, 1080, 50, Emu, Pt, font_size=20, align="center", gray=True)
        _add_text(slide, spec.get("context"), 100, 560, 1080, 80, Emu, Pt, font_size=16, align="center", gray=True)

    elif layout == "split_concept":
        # 좌우 분할
        _add_image_if_any(slide, spec, "left_image_path", lecture_dir, 60, 100, 540, 360, Emu)
        _add_image_if_any(slide, spec, "right_image_path", lecture_dir, 680, 100, 540, 360, Emu)
        _add_text(slide, spec.get("eyebrow"), 80, 30, 1120, 30, Emu, Pt, font_size=14, gray=True)
        _add_text(slide, spec.get("left_title"), 60, 470, 540, 50, Emu, Pt, font_size=24, bold=True, align="center")
        _add_text(slide, spec.get("left_body"), 60, 525, 540, 100, Emu, Pt, font_size=16, align="center")
        _add_text(slide, spec.get("right_title"), 680, 470, 540, 50, Emu, Pt, font_size=24, bold=True, align="center")
        _add_text(slide, spec.get("right_body"), 680, 525, 540, 100, Emu, Pt, font_size=16, align="center")
        _add_text(slide, spec.get("conclusion"), 80, 640, 1120, 60, Emu, Pt, font_size=18, bold=True, align="center")

    elif layout == "comparison_table":
        _add_text(slide, spec.get("eyebrow"), 80, 40, 1120, 30, Emu, Pt, font_size=14, gray=True)
        _add_text(slide, spec.get("title"), 80, 75, 1120, 50, Emu, Pt, font_size=32, bold=True)
        # 표는 PPT 표 객체로
        headers = spec.get("headers") or []
        rows = spec.get("rows") or []
        if headers and rows:
            n_cols = len(headers)
            n_rows = len(rows) + 1  # +1 for header row
            try:
                table_shape = slide.shapes.add_table(
                    n_rows, n_cols,
                    Emu(80 * PX_TO_EMU), Emu(150 * PX_TO_EMU),
                    Emu(1120 * PX_TO_EMU), Emu(min(n_rows * 60, 500) * PX_TO_EMU),
                )
                tbl = table_shape.table
                # 헤더
                for ci, h in enumerate(headers):
                    cell = tbl.cell(0, ci)
                    cell.text = str(h)
                # 행
                for ri, row in enumerate(rows):
                    if not isinstance(row, list):
                        continue
                    for ci in range(min(n_cols, len(row))):
                        tbl.cell(ri + 1, ci).text = str(row[ci])
            except Exception as e:
                print(f"[pptx_editable] 표 생성 실패: {e}")

    elif layout == "factbox":
        _add_text(slide, spec.get("eyebrow"), 80, 40, 1120, 30, Emu, Pt, font_size=14, gray=True)
        _add_text(slide, spec.get("title"), 80, 75, 1120, 60, Emu, Pt, font_size=32, bold=True)
        _add_text(slide, spec.get("body"), 80, 145, 1120, 80, Emu, Pt, font_size=18)
        items = spec.get("items") or []
        if items:
            text = "\n".join(f"• {it}" for it in items if it)
            _add_text(slide, text, 80, 240, 1120, 400, Emu, Pt, font_size=20)
        _add_text(slide, spec.get("source"), 80, 660, 1120, 40, Emu, Pt, font_size=12, gray=True, italic=True)

    elif layout == "metaphor_story":
        _add_text(slide, spec.get("eyebrow"), 80, 40, 800, 30, Emu, Pt, font_size=14, gray=True)
        _add_text(slide, spec.get("title"), 80, 75, 1120, 60, Emu, Pt, font_size=32, bold=True)
        _add_text(slide, spec.get("label"), 950, 75, 250, 40, Emu, Pt, font_size=14, gray=True, align="right")
        _add_text(slide, spec.get("story"), 80, 160, 1120, 350, Emu, Pt, font_size=20)
        _add_text(slide, spec.get("takeaway"), 80, 540, 1120, 140, Emu, Pt, font_size=22, bold=True)

    elif layout in ("illustration_anchor", "illustration_background", "illustration_overlay"):
        # 상단 이미지 + 하단 텍스트
        _add_image_if_any(slide, spec, "image_path", lecture_dir, 80, 60, 1120, 380, Emu)
        _add_text(slide, spec.get("eyebrow"), 80, 460, 1120, 30, Emu, Pt, font_size=14, gray=True)
        _add_text(slide, spec.get("title"), 80, 495, 1120, 60, Emu, Pt, font_size=28, bold=True)
        _add_text(slide, spec.get("body"), 80, 565, 1120, 80, Emu, Pt, font_size=18)
        _add_text(slide, spec.get("takeaway") or spec.get("subtitle"),
                  80, 650, 1120, 50, Emu, Pt, font_size=16, bold=True, gray=True)

    else:
        # 기본: lecture_body 패턴 (comparison_iconic 등 미매핑 layout 포함)
        _add_text(slide, spec.get("eyebrow"), 80, 40, 1120, 30, Emu, Pt, font_size=14, gray=True)
        _add_text(slide, spec.get("title"), 80, 75, 1120, 70, Emu, Pt, font_size=36, bold=True)
        _add_text(slide, spec.get("body"), 80, 160, 1120, 100, Emu, Pt, font_size=18)
        bullets = spec.get("bullets") or []
        if bullets:
            text = "\n".join(f"• {b}" for b in bullets if b)
            _add_text(slide, text, 80, 280, 1120, 340, Emu, Pt, font_size=20)
        _add_text(slide, spec.get("quote"), 80, 630, 1120, 50, Emu, Pt, font_size=16, italic=True, gray=True)
        _add_text(slide, spec.get("footer"), 80, 680, 1120, 30, Emu, Pt, font_size=11, gray=True, align="right")


def _add_text(
    slide, text, x_px, y_px, w_px, h_px, Emu, Pt,
    font_size=18, bold=False, italic=False, align=None, gray=False,
):
    """텍스트박스 추가. 빈 텍스트는 무시."""
    if not text or not str(text).strip():
        return
    tb = slide.shapes.add_textbox(
        Emu(x_px * PX_TO_EMU), Emu(y_px * PX_TO_EMU),
        Emu(w_px * PX_TO_EMU), Emu(h_px * PX_TO_EMU),
    )
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = str(text)
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    if gray:
        from pptx.dml.color import RGBColor
        run.font.color.rgb = RGBColor(0x66, 0x66, 0x66)
    if align == "center":
        from pptx.enum.text import PP_ALIGN
        p.alignment = PP_ALIGN.CENTER
    elif align == "right":
        from pptx.enum.text import PP_ALIGN
        p.alignment = PP_ALIGN.RIGHT


def _add_image_if_any(slide, spec, key, lecture_dir, x_px, y_px, w_px, h_px, Emu):
    """spec[key]에 일러스트 경로가 있으면 그림 추가. 없으면 무시."""
    path_str = spec.get(key)
    if not path_str:
        return
    p = Path(path_str)
    if not p.is_absolute():
        p = lecture_dir / path_str
    if not p.exists():
        return
    try:
        slide.shapes.add_picture(
            str(p),
            Emu(x_px * PX_TO_EMU), Emu(y_px * PX_TO_EMU),
            width=Emu(w_px * PX_TO_EMU), height=Emu(h_px * PX_TO_EMU),
        )
    except Exception as e:
        print(f"[pptx_editable] 이미지 삽입 실패 {key}: {e}")


# '글자 얹기' 오버레이 → 텍스트박스 좌표 (slide_overlay.py CSS 어휘의 1280×720 사상)
_OVERLAY_POSITIONS = {
    "top-left": ("left", "top"), "top": ("center", "top"), "top-right": ("right", "top"),
    "left": ("left", "middle"), "center": ("center", "middle"), "right": ("right", "middle"),
    "bottom-left": ("left", "bottom"), "bottom": ("center", "bottom"),
    "bottom-right": ("right", "bottom"),
}
_OVERLAY_SIZE_VW = {"small": 2.0, "medium": 2.9, "large": 4.0}
# PPT 서체 사상 — sans/serif 는 Office 표준 한글 서체, 웹폰트 계열은 실제 이름
# (설치돼 있으면 그대로, 없으면 PPT 가 기본 서체로 폴백)
_OVERLAY_FONT_NAME = {
    "sans": "맑은 고딕", "serif": "바탕", "gowun": "Gowun Batang",
    "jua": "Jua", "black": "Black Han Sans",
    "pen": "Nanum Pen Script", "brush": "Nanum Brush Script",
}


def _overlay_rgb(color: str):
    """오버레이 color 어휘(white|black|#hex) → (RGBColor, 어두운 글자 여부)."""
    from pptx.dml.color import RGBColor
    c = (color or "white").strip()
    if c == "black":
        return RGBColor(0x11, 0x11, 0x11), True
    if c.startswith("#") and len(c) in (4, 7, 9):
        hexpart = c[1:7] if len(c) >= 7 else "".join(ch * 2 for ch in c[1:4])
        try:
            r, g, b = (int(hexpart[i:i + 2], 16) for i in (0, 2, 4))
            return RGBColor(r, g, b), (0.299 * r + 0.587 * g + 0.114 * b) < 128
        except ValueError:
            pass
    return RGBColor(0xFF, 0xFF, 0xFF), False


def _add_overlay_textbox(slide, ov: dict, Emu, Pt):
    """'글자 얹기' 1건을 편집 가능 텍스트박스로 — 합성판과 같은 위치·크기·색·서체.

    자유 좌표(x/y% — 박스 좌상단)가 있으면 9방(position)보다 우선(합성기와 같은 규칙).
    합성판의 그림자(text-shadow)는 PPTX 표현 밖이라 생략, 배경칩은 도형 채움으로 근사.
    """
    from pptx.enum.text import PP_ALIGN
    text = str(ov.get("text") or "").strip()
    if not text:
        return
    size_vw = _OVERLAY_SIZE_VW.get(ov.get("size") or "small", _OVERLAY_SIZE_VW["small"])
    try:
        v = float(ov.get("size_vw") or 0)
        if 0.5 <= v <= 12:
            size_vw = v
    except (TypeError, ValueError):
        pass
    font_px = int(SLIDE_W_PX * size_vw / 100)
    rgb, dark_text = _overlay_rgb(ov.get("color"))
    chip = bool(ov.get("chip"))
    font_name = _OVERLAY_FONT_NAME.get(ov.get("font") or "sans", _OVERLAY_FONT_NAME["sans"])

    mx, my = 41, 32          # 3.2% / 4.5% 여백
    w = 896                  # max-width 70%
    lines = max(1, text.count("\n") + 1)
    h = int(lines * font_px * 1.4) + (int(font_px * 0.9) if chip else 8)
    free_xy = None
    try:
        fx, fy = float(ov.get("x")), float(ov.get("y"))
        if -10 <= fx <= 110 and -10 <= fy <= 110:
            free_xy = (fx, fy)
    except (TypeError, ValueError):
        pass
    if free_xy is not None:
        # 자유 배치 — 좌상단 % 그대로, 좌측 정렬 (합성기와 동일 규칙)
        h_align = "left"
        x = int(SLIDE_W_PX * free_xy[0] / 100)
        y = int(SLIDE_H_PX * free_xy[1] / 100)
        w = min(w, max(80, SLIDE_W_PX - x - 8))
    else:
        h_align, v_align = _OVERLAY_POSITIONS.get(
            ov.get("position") or "bottom-right", _OVERLAY_POSITIONS["bottom-right"])
        x = mx if h_align == "left" else (SLIDE_W_PX - w) // 2 if h_align == "center" else SLIDE_W_PX - mx - w
        y = my if v_align == "top" else (SLIDE_H_PX - h) // 2 if v_align == "middle" else SLIDE_H_PX - my - h

    tb = slide.shapes.add_textbox(
        Emu(x * PX_TO_EMU), Emu(y * PX_TO_EMU), Emu(w * PX_TO_EMU), Emu(h * PX_TO_EMU))
    if chip:
        from pptx.dml.color import RGBColor
        tb.fill.solid()
        tb.fill.fore_color.rgb = RGBColor(0xF2, 0xF2, 0xF2) if dark_text else RGBColor(0x0C, 0x0C, 0x0E)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = p.add_run()
        run.text = line
        run.font.size = Pt(round(font_px * 0.75))
        run.font.bold = ov.get("weight") != "normal"  # 합성판 기본 font-weight 600 대응
        run.font.name = font_name
        run.font.color.rgb = rgb
        p.alignment = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER,
                       "right": PP_ALIGN.RIGHT}[h_align]


# ─────────────────────────────────────────────────────────────────────
# 통합 진입점
# ─────────────────────────────────────────────────────────────────────

# ─────────────────────────────────────────────────────────────────────
# 이미지 폴더 (슬라이드 각 장 = PNG 파일 하나)
# ─────────────────────────────────────────────────────────────────────

def export_images(lecture_id: str) -> dict:
    """슬라이드 각 장을 PNG 파일로 폴더 하나에 담아 내보내기.

    exports/{lecture_id}_{ts}_images/ 에 순번(+제목) 이름으로 복사하고,
    같은 이름 .zip 도 만든다 — 브라우저/원격 다운로드는 단일 파일만 전달되므로
    ZIP(풀면 그 폴더)이 전달체, 로컬(데스크탑)에선 폴더가 바로 산출물.

    Returns: {success, path(zip), folder, filename(zip), slide_count, skipped, format: "images"}
    """
    import re
    import shutil
    import zipfile

    deck = lecture_store.read_deck(lecture_id)
    lecture_dir = lecture_store.lecture_dir(lecture_id)
    order = deck.get("slide_order", [])
    folder = _exports_dir(lecture_id) / f"{lecture_id}_{_timestamp()}_images"
    folder.mkdir(parents=True, exist_ok=True)

    width = max(2, len(str(len(order) or 1)))  # 01, 02 … (100장 넘으면 3자리)
    count = skipped = 0
    for i, sid in enumerate(order, 1):
        meta = deck.get("slides", {}).get(sid, {})
        rel = meta.get("png_file")
        src = (lecture_dir / rel) if rel else None
        if src is None or not src.exists():
            skipped += 1
            continue
        # 파일명 = 순번 + 제목(파일시스템 금지문자 제거, 과장 방지 60자 캡)
        title = re.sub(r'[\\/:*?"<>|\r\n\t]+', " ", str(meta.get("title") or "")).strip()
        title = re.sub(r"\s+", " ", title)[:60].strip()
        # 확장자는 실제 매직바이트로 — 이미지 모델이 .png 이름에 JPEG 바이트를 담는 일이 실재한다
        with src.open("rb") as fh:
            head = fh.read(8)
        if head.startswith(b"\x89PNG"):
            suffix = ".png"
        elif head.startswith(b"\xff\xd8\xff"):
            suffix = ".jpg"
        else:
            suffix = src.suffix or ".png"
        name = f"{i:0{width}d}" + (f" {title}" if title else "") + suffix
        shutil.copy2(src, folder / name)
        count += 1

    if count == 0:
        shutil.rmtree(folder, ignore_errors=True)
        raise RuntimeError("내보낼 슬라이드 이미지가 없습니다.")

    zip_path = folder.with_suffix(".zip")
    # PNG는 이미 압축돼 있어 STORED(무압축)가 빠르고 크기도 거의 같다
    with zipfile.ZipFile(zip_path, "w", zipfile.ZIP_STORED) as zf:
        for p in sorted(folder.iterdir()):
            zf.write(p, arcname=f"{folder.name}/{p.name}")

    return {
        "success": True,
        "format": "images",
        "path": str(zip_path.resolve()),
        "folder": str(folder.resolve()),
        "filename": zip_path.name,
        "slide_count": count,
        "skipped": skipped,
    }


def export_deck(lecture_id: str, format: str) -> dict:
    """format에 따라 분기.

    format:
      - "pdf": 다중 페이지 PDF (PIL)
      - "pptx" 또는 "pptx_image": 통째 이미지 PPTX (디자인 완벽 보존, 편집 불가)
      - "pptx_editable": 텍스트박스로 분해된 PPTX (PPT에서 자유 편집 가능, 디자인 단순화)
      - "images": 슬라이드 각 장을 PNG 파일로 폴더 하나에 (+다운로드용 ZIP)
    """
    fmt = (format or "").lower().strip()
    if fmt == "pdf":
        return export_pdf(lecture_id)
    elif fmt in ("pptx", "pptx_image"):
        return export_pptx(lecture_id)
    elif fmt == "pptx_editable":
        return export_pptx_editable(lecture_id)
    elif fmt == "images":
        return export_images(lecture_id)
    else:
        raise ValueError(
            f"지원하지 않는 형식: {format} (pdf/pptx/pptx_editable/images)"
        )
