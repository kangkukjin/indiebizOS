import os
import uuid
import json
import asyncio
import base64
# ★무거운 미디어 의존성(edge_tts·moviepy)은 모듈레벨에서 import하지 않는다 — 폰엔 이 라이브러리가
# 없어 모듈 전체 import가 실패하고, 그러면 같은 파일의 *순수* 연산(document html 렌더·structure·
# image_critic·image_gemini = 문자열/HTTP뿐)까지 mac에 갇혔다. 무거운 건 쓰는 함수 안에서 지연 import.
# (moviepy는 이미 각 함수가 로컬 re-import 중이라 모듈레벨 심볼은 0회 사용 — 그냥 제거. edge_tts는 generate_tts로.)

# Helper functions
def get_image_base64(image_path):
    if not image_path or not os.path.exists(image_path):
        return None
    try:
        with open(image_path, "rb") as image_file:
            encoded_string = base64.b64encode(image_file.read()).decode('utf-8')
            ext = os.path.splitext(image_path)[1].lower()
            mime = "image/jpeg" if ext in [".jpg", ".jpeg"] else "image/png" if ext == ".png" else "image/webp" if ext == ".webp" else "image/jpeg"
            return f"data:{mime};base64,{encoded_string}"
    except Exception as e:
        print(f"Base64 인코딩 실패: {e}")
        return None


def execute(tool_input: dict, context) -> str:
    """도구 실행 함수 (ToolContext 기반 신규 시그니처).

    context는 tool_context.ToolContext. tool_name 분기와 출력 경로 결정에 사용.
    output_dir()은 항상 절대경로 + mkdir 자동.
    """
    tool_name = context.tool_name
    output_base = context.output_dir()

    # 통합 액션: tool_name=="html_video" + scene_dir 유무로 분기
    if tool_name == "html_video":
        if tool_input.get("scene_dir"):
            return render_html_video(tool_input, output_base)
        return create_html_video(tool_input, output_base)

    if tool_name == "render_html_to_image":
        return render_html_to_image(tool_input, output_base)
    elif tool_name == "render_document":
        return render_document(tool_input, output_base)
    elif tool_name == "structure_document":
        return structure_document(tool_input, output_base)
    elif tool_name == "generate_gemini_image":
        return generate_gemini_image(tool_input, output_base)
    elif tool_name == "critique_gemini_image":
        return critique_gemini_image(tool_input, output_base)
    elif tool_name == "read_gemini_image":
        return read_gemini_image(tool_input, output_base)
    elif tool_name == "create_tts":
        return create_tts(tool_input, output_base)
    elif tool_name == "create_shadcn_slides":
        # shadcn 스타일 슬라이드 생성 (별도 모듈)
        import importlib.util
        import sys
        module_path = os.path.join(os.path.dirname(__file__), "shadcn_slides.py")
        spec = importlib.util.spec_from_file_location("shadcn_slides", module_path)
        shadcn_slides = importlib.util.module_from_spec(spec)
        sys.modules["shadcn_slides"] = shadcn_slides
        spec.loader.exec_module(shadcn_slides)
        return shadcn_slides.create_shadcn_slides(tool_input, output_base)
    elif tool_name == "create_slide":
        # [engines:slide] 저작 기층 — 명령 → 슬라이드 1장 (별도 모듈, 호출마다 재로드)
        import importlib.util
        import sys
        module_path = os.path.join(os.path.dirname(__file__), "slide_author.py")
        spec = importlib.util.spec_from_file_location("slide_author", module_path)
        slide_author = importlib.util.module_from_spec(spec)
        sys.modules["slide_author"] = slide_author
        spec.loader.exec_module(slide_author)
        return slide_author.create_slide(tool_input, output_base)
    # (2026-06-04 은퇴) create_lecture_plan/write/illustrate/compose 4단계 파이프라인 제거.
    # 강의 = AI가 lecture_slide_principles.md 가이드로 계획 → [engines:slide] × N(같은 style, & 병렬).
    # 계획은 인지(가이드+추론)지 액션 아님 (ibl_design_philosophy Mode C).

    return f"알 수 없는 도구: {tool_name}"

async def generate_tts(text, output_path, voice="ko-KR-SunHiNeural", rate="+0%", pitch="+0Hz"):
    import edge_tts  # 지연 import — 폰엔 없음(tts는 mac_only라 폰선 호출 안 됨)
    communicate = edge_tts.Communicate(text, voice, rate=rate, pitch=pitch)
    await communicate.save(output_path)

def _prepare_scene_html(html, base_path, video_width, video_height):
    """씬 HTML을 렌더링 가능하도록 전처리 (이미지 Base64 변환, 뷰포트 설정)"""
    import re

    # 항상 이미지 경로를 Base64로 변환 (절대 경로 포함)
    resolved_base = os.path.abspath(base_path) if base_path else None

    def replace_img_src(match):
        src = match.group(1)
        if src.startswith(('data:', 'http://', 'https://')):
            return match.group(0)
        # file:// 프로토콜 처리
        if src.startswith('file://'):
            src = src[7:]  # file:// 제거
        # 절대 경로이면 바로 사용, 상대 경로이면 base_path 기준
        if os.path.isabs(src):
            abs_path = src
        elif resolved_base:
            abs_path = os.path.join(resolved_base, src)
        else:
            return match.group(0)  # base_path도 없고 상대 경로면 스킵
        base64_data = get_image_base64(abs_path)
        if base64_data:
            return f'src="{base64_data}"'
        return match.group(0)

    html = re.sub(r'src=["\']([^"\']+)["\']', replace_img_src, html)

    def replace_bg_url(match):
        url = match.group(1)
        if url.startswith(('data:', 'http://', 'https://')):
            return match.group(0)
        if url.startswith('file://'):
            url = url[7:]
        if os.path.isabs(url):
            abs_path = url
        elif resolved_base:
            abs_path = os.path.join(resolved_base, url)
        else:
            return match.group(0)
        base64_data = get_image_base64(abs_path)
        if base64_data:
            return f'url("{base64_data}")'
        return match.group(0)

    html = re.sub(r'url\(["\']?([^"\')\s]+)["\']?\)', replace_bg_url, html)

    # 한글 폰트 보장을 위한 CSS (Google Fonts 로드 실패 시에도 시스템 폰트 fallback)
    korean_font_css = """
<style>
@import url('https://fonts.googleapis.com/css2?family=Noto+Sans+KR:wght@300;400;500;700;900&display=swap');
* { font-family: 'Noto Sans KR', 'Apple SD Gothic Neo', 'Malgun Gothic', '맑은 고딕', sans-serif !important; }
</style>"""

    # 이미 완전한 HTML 문서인 경우 한글 폰트 CSS만 주입
    if html.strip().lower().startswith("<!doctype") or html.strip().lower().startswith("<html"):
        # </head> 앞에 폰트 CSS 삽입
        if '</head>' in html:
            html = html.replace('</head>', korean_font_css + '\n</head>', 1)
        elif '<body' in html:
            html = html.replace('<body', korean_font_css + '\n<body', 1)
        return html

    # 그렇지 않으면 래핑 (Google Fonts + Lottie CDN 포함)
    return f"""<!DOCTYPE html>
<html><head><meta charset="UTF-8">
<script src="https://cdn.tailwindcss.com"></script>
<link rel="preconnect" href="https://fonts.googleapis.com">
<link rel="preconnect" href="https://fonts.gstatic.com" crossorigin>
<link href="https://fonts.googleapis.com/css2?family=Black+Han+Sans&family=Do+Hyeon&family=Gothic+A1:wght@400;700;900&family=Noto+Sans+KR:wght@300;400;500;700;900&family=Sunflower:wght@300;500;700&family=Jua&family=Inter:wght@300;400;500;600;700;800;900&family=Montserrat:wght@400;600;700;800;900&family=Playfair+Display:wght@400;600;700;900&family=Poppins:wght@300;400;500;600;700;800&display=swap" rel="stylesheet">
<link rel="stylesheet" href="https://cdnjs.cloudflare.com/ajax/libs/animate.css/4.1.1/animate.min.css"/>
<script src="https://unpkg.com/lucide@latest"></script>
<script src="https://cdnjs.cloudflare.com/ajax/libs/gsap/3.12.2/gsap.min.js"></script>
<script src="https://unpkg.com/@lottiefiles/lottie-player@2/dist/lottie-player.js"></script>
<style>body,html{{margin:0;padding:0;width:{video_width}px;height:{video_height}px;overflow:hidden;font-family:'Noto Sans KR',sans-serif;}}</style>
</head><body>{html}</body></html>"""


def _auto_split_scenes(scenes, narration_texts, default_duration):
    """단일 HTML에 여러 씬이 들어있는 경우 자동 분리.

    AI 에이전트가 모든 씬을 하나의 HTML에 <div id="scene1">, <div id="scene2">... 형태로
    넣는 경우를 감지하여 각 씬을 독립 HTML 문서로 분리합니다.
    """
    import re
    from bs4 import BeautifulSoup

    if len(scenes) != 1:
        return scenes, narration_texts

    html = scenes[0].get("html", "")
    duration = scenes[0].get("duration", default_duration)
    base_path = scenes[0].get("base_path")

    # scene/씬 패턴의 id를 가진 div 요소가 2개 이상인지 확인
    scene_id_pattern = re.compile(r'id=["\'](?:scene|씬)\s*(\d+)["\']', re.IGNORECASE)
    scene_matches = scene_id_pattern.findall(html)

    if len(scene_matches) < 2:
        return scenes, narration_texts

    print(f"[create_html_video] 단일 HTML에 {len(scene_matches)}개의 씬 감지 → 자동 분리")

    try:
        soup = BeautifulSoup(html, 'html.parser')

        # <head> 내용 추출 (공통 스타일/스크립트)
        head_tag = soup.find('head')
        head_content = str(head_tag) if head_tag else "<head><meta charset='UTF-8'></head>"

        # 씬 div 찾기 (id="scene1", id="scene2", ...)
        scene_divs = []
        for match in scene_matches:
            div = soup.find(id=re.compile(f'^(?:scene|씬)\\s*{match}$', re.IGNORECASE))
            if div:
                scene_divs.append(div)

        if len(scene_divs) < 2:
            return scenes, narration_texts

        # body의 style 추출 (배경색, 폰트 등)
        body_tag = soup.find('body')
        body_attrs = ""
        if body_tag and body_tag.attrs:
            attrs_str = " ".join(f'{k}="{v}"' if isinstance(v, str) else f'{k}="{" ".join(v)}"' for k, v in body_tag.attrs.items())
            body_attrs = f" {attrs_str}"

        # 각 씬을 독립 HTML 문서로 구성
        new_scenes = []
        per_scene_duration = duration / len(scene_divs) if duration else default_duration

        for i, div in enumerate(scene_divs):
            # 씬 div에 overflow:hidden과 뷰포트 크기 강제 적용
            div_style = div.get('style', '')
            if 'overflow' not in div_style:
                div['style'] = div_style + '; overflow: hidden;' if div_style else 'overflow: hidden;'

            scene_html = f"""<!DOCTYPE html>
<html lang="ko">
{head_content}
<body{body_attrs}>
{str(div)}
</body>
</html>"""
            new_scenes.append({
                "html": scene_html,
                "duration": per_scene_duration,
                "base_path": base_path
            })

        # narration_texts가 원래 씬 개수와 맞지 않으면 조정
        new_narrations = narration_texts
        if len(narration_texts) == 1 and len(new_scenes) > 1:
            # 하나의 나레이션만 있으면 첫 씬에만 적용
            new_narrations = narration_texts + [""] * (len(new_scenes) - 1)
        elif len(narration_texts) == 0:
            new_narrations = []

        print(f"[create_html_video] 자동 분리 완료: {len(new_scenes)}개 씬, 각 {per_scene_duration:.1f}초")
        return new_scenes, new_narrations

    except Exception as e:
        print(f"[create_html_video] 자동 분리 실패, 원본 유지: {e}")
        return scenes, narration_texts


def create_html_video(tool_input, output_base):
    """HTML 씬들을 프레임별 스크린샷으로 캡처하여 MP4 동영상을 생성합니다.

    각 씬은 독립된 HTML 문서와 재생 시간(duration)으로 구성됩니다.
    CSS 애니메이션, GSAP, Animate.css 등 브라우저에서 렌더링 가능한 모든 애니메이션이 지원됩니다.

    파이프라인:
      1) 나레이션 TTS 먼저 생성
      2) 나레이션 길이에 맞게 씬 duration 자동 조정 (겹침 방지)
      3) HTML → Playwright 프레임 캡처(PNG)
      4) FFmpeg MP4 인코딩
      5) 나레이션/BGM 합성 → 최종 MP4
    """
    import subprocess
    import shutil
    from playwright.sync_api import sync_playwright

    scenes = tool_input.get("scenes", [])
    narration_texts = tool_input.get("narration_texts", [])
    bgm_path = tool_input.get("bgm_path")
    voice = tool_input.get("voice", "ko-KR-SunHiNeural")
    rate = tool_input.get("rate", "+0%")
    pitch = tool_input.get("pitch", "+0Hz")
    output_filename = tool_input.get("output_filename", "html_video.mp4")
    video_width = tool_input.get("width", 1280)
    video_height = tool_input.get("height", 720)
    fps = tool_input.get("fps", 24)
    default_duration = tool_input.get("duration_per_scene", 5)
    # 씬 전환 효과 설정
    transition_type = tool_input.get("transition", "fade")  # fade, wipeleft, wiperight, slidedown, slideup, circleopen, dissolve, none
    transition_duration = tool_input.get("transition_duration", 0.5)  # 전환 시간(초)

    if not scenes:
        if tool_input.get("topic"):
            return (
                "오류: topic만으로는 영상을 만들 수 없습니다. 이 액션은 완성된 HTML 씬을 합성합니다 — "
                "scenes(각 {html, duration} 배열) 또는 scene_dir(저장된 씬 디렉토리)을 주세요. "
                "주제→슬라이드 자동 생성은 [engines:slide] 또는 [engines:lecture_*] 파이프라인을 사용하세요."
            )
        return "오류: scenes 리스트가 비어 있습니다 (scenes 배열 또는 scene_dir 필요)."

    # scenes 검증: html 키 필수
    missing_html = [i for i, s in enumerate(scenes) if isinstance(s, dict) and "html" not in s]
    if missing_html:
        wrong_keys = list(scenes[missing_html[0]].keys()) if missing_html else []
        return (
            f"오류: scenes[{missing_html[0]}]에 'html' 키가 없습니다. "
            f"(현재 키: {wrong_keys})\n"
            f"각 씬은 완전한 HTML 문서가 필요합니다.\n"
            f"올바른 형식:\n"
            f'  scenes: [{{"html": "<!DOCTYPE html><html><head><meta charset=\\"utf-8\\"><script src=\\"https://cdn.tailwindcss.com\\"></script></head>'
            f'<body><div style=\\"width:{video_width}px;height:{video_height}px\\" class=\\"flex items-center justify-center bg-gradient-to-br from-slate-900 to-purple-900\\">'
            f'<h1 class=\\"text-6xl font-bold text-white\\">제목</h1></div></body></html>", "duration": 5}}]\n'
            f"narration_texts: [\"나레이션 텍스트\"]  (씬별 나레이션은 별도 배열)"
        )

    # 자동 씬 분리: 단일 HTML에 여러 씬이 들어있는 경우 감지 및 분리
    scenes, narration_texts = _auto_split_scenes(scenes, narration_texts, default_duration)

    output_path = os.path.join(output_base, output_filename)
    temp_dir = os.path.join(output_base, f"temp_htmlvideo_{uuid.uuid4().hex[:8]}")
    frames_dir = os.path.join(temp_dir, "frames")
    os.makedirs(frames_dir, exist_ok=True)

    try:
        # ============================================================
        # 1단계: 나레이션 TTS를 먼저 생성 (씬 duration 조정을 위해)
        # ============================================================
        narration_audio_paths = []
        narration_durations = []  # 각 나레이션의 실제 재생 시간
        has_narration = False

        for i, scene in enumerate(scenes):
            if i < len(narration_texts) and narration_texts[i]:
                tts_path = os.path.join(temp_dir, f"narration_{i}.mp3")
                asyncio.run(generate_tts(narration_texts[i], tts_path, voice, rate, pitch))
                narration_audio_paths.append(tts_path)
                has_narration = True

                # 나레이션 길이 측정
                from moviepy import AudioFileClip as MpAudioClip
                narr_clip = MpAudioClip(tts_path)
                narration_durations.append(narr_clip.duration)
                narr_clip.close()
            else:
                narration_audio_paths.append(None)
                narration_durations.append(0)

        # ============================================================
        # 2단계: 나레이션 길이에 맞게 씬 duration 자동 조정
        #   나레이션이 씬보다 길면 → 씬을 나레이션 + 여유시간으로 늘림
        #   이렇게 해야 나레이션끼리 겹치지 않음
        # ============================================================
        NARRATION_PADDING = 0.5  # 나레이션 끝나고 다음 씬 시작까지 여유 시간(초)

        for i, scene in enumerate(scenes):
            original_dur = scene.get("duration", default_duration)
            narr_dur = narration_durations[i] if i < len(narration_durations) else 0

            if narr_dur > 0:
                needed_dur = narr_dur + NARRATION_PADDING
                if needed_dur > original_dur:
                    print(f"[create_html_video] 씬 {i+1}: duration {original_dur:.1f}s → {needed_dur:.1f}s (나레이션 {narr_dur:.1f}s에 맞춤)")
                    scene["duration"] = needed_dur

        # ============================================================
        # 3단계: 씬별 프레임 캡처 및 개별 MP4 인코딩
        # ============================================================
        scene_videos = []  # 씬별 MP4 경로 리스트
        scene_durations = []  # 씬별 실제 duration 리스트

        with sync_playwright() as p:
            browser = p.chromium.launch()
            page = browser.new_page(viewport={"width": video_width, "height": video_height})

            for i, scene in enumerate(scenes):
                html = scene.get("html", "")
                duration = scene.get("duration", default_duration)
                base_path_opt = scene.get("base_path") or output_base

                if not html:
                    continue

                html_ready = _prepare_scene_html(html, base_path_opt, video_width, video_height)

                scene_html_dir = base_path_opt if base_path_opt and os.path.isdir(base_path_opt) else temp_dir
                scene_html_path = os.path.join(scene_html_dir, f"_scene_{i}_{uuid.uuid4().hex[:6]}.html")
                with open(scene_html_path, "w", encoding="utf-8") as sf:
                    sf.write(html_ready)
                page.goto(f"file://{scene_html_path}")

                # 모든 애니메이션을 즉시 일시정지
                page.evaluate("""() => {
                    const style = document.createElement('style');
                    style.id = '__anim_pause__';
                    style.textContent = '*, *::before, *::after { animation-play-state: paused !important; }';
                    document.head.appendChild(style);
                }""")

                page.wait_for_load_state("networkidle")
                page.wait_for_timeout(1500 if i == 0 else 500)

                # 리소스 로딩 완료 후 애니메이션 재개
                page.evaluate("""() => {
                    const pauseStyle = document.getElementById('__anim_pause__');
                    if (pauseStyle) pauseStyle.remove();
                }""")

                # 씬별 프레임 디렉토리
                scene_frames_dir = os.path.join(temp_dir, f"frames_scene_{i}")
                os.makedirs(scene_frames_dir, exist_ok=True)

                total_frames = int(duration * fps)
                frame_interval_ms = 1000.0 / fps

                for f in range(total_frames):
                    frame_path = os.path.join(scene_frames_dir, f"frame_{f:06d}.png")
                    page.screenshot(path=frame_path)
                    page.wait_for_timeout(int(frame_interval_ms))

                if total_frames == 0:
                    continue

                # 씬별 MP4 인코딩
                scene_mp4 = os.path.join(temp_dir, f"scene_{i}.mp4")
                enc_result = subprocess.run([
                    "ffmpeg", "-y",
                    "-framerate", str(fps),
                    "-i", os.path.join(scene_frames_dir, "frame_%06d.png"),
                    "-c:v", "libx264", "-preset", "fast", "-crf", "20",
                    "-pix_fmt", "yuv420p",
                    scene_mp4
                ], capture_output=True)

                if enc_result.returncode != 0:
                    print(f"[create_html_video] 씬 {i} 인코딩 실패: {enc_result.stderr.decode()}")
                    continue

                scene_videos.append(scene_mp4)
                scene_durations.append(duration)

            browser.close()

        if not scene_videos:
            shutil.rmtree(temp_dir)
            return "오류: 캡처된 프레임이 없습니다."

        # ============================================================
        # 4단계: 씬 전환 효과 적용 및 병합
        # ============================================================
        merged_video = os.path.join(temp_dir, "merged.mp4")
        use_transition = transition_type != "none" and len(scene_videos) > 1 and transition_duration > 0

        if use_transition:
            # FFmpeg xfade 필터를 사용하여 씬 간 전환 효과 적용
            # 유효한 xfade 전환 효과 목록
            valid_transitions = {
                "fade", "wipeleft", "wiperight", "wipeup", "wipedown",
                "slidedown", "slideup", "slideleft", "slideright",
                "circleopen", "circleclose", "dissolve", "pixelize",
                "radial", "hblur", "fadegrays", "squeezeh", "squeezev",
            }
            if transition_type not in valid_transitions:
                transition_type = "fade"

            td = min(transition_duration, min(scene_durations) * 0.4)  # 전환이 씬 길이의 40%를 넘지 않도록

            if len(scene_videos) == 2:
                # 2개 씬: 단일 xfade
                offset = scene_durations[0] - td
                ffmpeg_cmd = [
                    "ffmpeg", "-y",
                    "-i", scene_videos[0], "-i", scene_videos[1],
                    "-filter_complex",
                    f"[0:v][1:v]xfade=transition={transition_type}:duration={td}:offset={offset},format=yuv420p[v]",
                    "-map", "[v]", "-c:v", "libx264", "-preset", "fast", "-crf", "20",
                    merged_video
                ]
            else:
                # 3개 이상 씬: xfade 체이닝
                inputs = []
                for sv in scene_videos:
                    inputs.extend(["-i", sv])

                filter_parts = []
                cumulative_offset = 0.0

                # 첫 번째 xfade
                cumulative_offset = scene_durations[0] - td
                filter_parts.append(
                    f"[0:v][1:v]xfade=transition={transition_type}:duration={td}:offset={cumulative_offset}[v1]"
                )

                # 나머지 xfade 체이닝
                for j in range(2, len(scene_videos)):
                    prev_label = f"v{j-1}"
                    out_label = f"v{j}"
                    # 이전까지의 총 길이 = 이전 누적 + 현재 씬 duration - 전환 겹침
                    cumulative_offset += scene_durations[j-1] - td
                    filter_parts.append(
                        f"[{prev_label}][{j}:v]xfade=transition={transition_type}:duration={td}:offset={cumulative_offset}[{out_label}]"
                    )

                final_label = f"v{len(scene_videos)-1}"
                filter_complex = ";".join(filter_parts) + f";[{final_label}]format=yuv420p[vout]"

                ffmpeg_cmd = [
                    "ffmpeg", "-y",
                    *inputs,
                    "-filter_complex", filter_complex,
                    "-map", "[vout]", "-c:v", "libx264", "-preset", "fast", "-crf", "20",
                    merged_video
                ]

            print(f"[create_html_video] 씬 전환 적용: {transition_type}, {td:.2f}초")
            xfade_result = subprocess.run(ffmpeg_cmd, capture_output=True)

            if xfade_result.returncode != 0:
                # xfade 실패 시 fallback: 전환 없이 단순 concat
                print(f"[create_html_video] xfade 실패, concat으로 대체: {xfade_result.stderr.decode()[:300]}")
                use_transition = False

        if not use_transition:
            # 전환 효과 없이 단순 연결 (1개 씬이거나 transition=none이거나 xfade 실패 시)
            if len(scene_videos) == 1:
                shutil.copy2(scene_videos[0], merged_video)
            else:
                concat_list = os.path.join(temp_dir, "concat.txt")
                with open(concat_list, "w") as cl:
                    for sv in scene_videos:
                        cl.write(f"file '{sv}'\n")
                concat_result = subprocess.run([
                    "ffmpeg", "-y", "-f", "concat", "-safe", "0",
                    "-i", concat_list,
                    "-c:v", "libx264", "-preset", "fast", "-crf", "20",
                    "-pix_fmt", "yuv420p",
                    merged_video
                ], capture_output=True)
                if concat_result.returncode != 0:
                    shutil.rmtree(temp_dir)
                    return f"FFmpeg concat 오류: {concat_result.stderr.decode()}"

        # ============================================================
        # 5단계: 나레이션과 BGM 합성
        # ============================================================
        if bgm_path:
            bgm_path = os.path.abspath(bgm_path)
        has_bgm = bgm_path and os.path.exists(bgm_path)

        if has_narration or has_bgm:
            from moviepy import AudioFileClip as MpAudioClip, CompositeAudioClip as MpCompositeAudioClip, VideoFileClip

            video_clip = VideoFileClip(merged_video)
            audio_layers = []

            if has_narration:
                # 씬 전환 시 xfade 겹침만큼 시작 타이밍을 앞당김
                actual_td = 0.0
                if use_transition and len(scene_videos) > 1:
                    actual_td = min(transition_duration, min(scene_durations) * 0.4)

                scene_start = 0.0
                for i, scene in enumerate(scenes):
                    dur = scene.get("duration", default_duration)
                    if i < len(narration_audio_paths) and narration_audio_paths[i]:
                        narr = MpAudioClip(narration_audio_paths[i])
                        narr = narr.with_start(scene_start)
                        audio_layers.append(narr)
                    # xfade 적용 시 씬 간 겹침을 반영
                    if i < len(scenes) - 1 and actual_td > 0:
                        scene_start += dur - actual_td
                    else:
                        scene_start += dur

            if has_bgm:
                bgm_clip = MpAudioClip(bgm_path).with_duration(video_clip.duration)
                if has_narration:
                    bgm_clip = bgm_clip.multiply_volume(0.2)
                else:
                    bgm_clip = bgm_clip.multiply_volume(0.5)
                audio_layers.append(bgm_clip)

            if audio_layers:
                final_audio = MpCompositeAudioClip(audio_layers)
                video_clip = video_clip.with_audio(final_audio)

            video_clip.write_videofile(output_path, fps=fps, codec="libx264", audio_codec="aac")
            video_clip.close()
        else:
            shutil.copy2(merged_video, output_path)

        # 임시 파일 정리
        shutil.rmtree(temp_dir)
        # output_base에 생성된 임시 씬 HTML 파일도 정리
        import glob
        for tmp_html in glob.glob(os.path.join(output_base, "_scene_*_*.html")):
            try:
                os.remove(tmp_html)
            except OSError:
                pass

        transition_info = ""
        if use_transition and len(scene_videos) > 1:
            actual_td = min(transition_duration, min(scene_durations) * 0.4)
            transition_info = f" | 씬 전환: {transition_type} ({actual_td:.1f}초)"

        return f"HTML 동영상 제작 완료: {os.path.abspath(output_path)}{transition_info}"
    except subprocess.CalledProcessError as e:
        return f"FFmpeg 오류: {e.stderr.decode() if e.stderr else str(e)}"
    except Exception as e:
        return f"HTML 동영상 제작 중 오류 발생: {str(e)}"


# ── 문서 IR(공유 문서 모델) → 산출물 emitter ───────────────────────────
# 문서 IR: {title?, blocks:[{type, ...}]}. 블록 타입:
#   heading{level,text} · paragraph{text} · list{ordered?,items[]} · image{src,caption?}
#   · table{columns,rows}(=데이터 통화 재사용) · quote{text,cite?} · code{text,lang?} · divider
# 포맷 중립 IR을 여러 emitter가 렌더(현재 html). slide/newspaper와 달리 단일 IR이 단일 진실 소스.
def _doc_blocks_to_html(blocks: list) -> str:
    import html as _html

    def esc(s):
        return _html.escape(str(s if s is not None else ""))

    parts = []
    for b in blocks:
        if not isinstance(b, dict):
            continue
        t = (b.get("type") or "paragraph").lower()
        if t == "heading":
            lvl = max(1, min(6, int(b.get("level") or 2)))
            txt = b.get("text")
            anchor = b.get("anchor") or (str(txt) if txt else "")  # 목차 점프용 id
            id_attr = f' id="{esc(anchor)}"' if anchor else ""
            parts.append(f"<h{lvl}{id_attr}>{esc(txt)}</h{lvl}>")
        elif t == "list":
            tag = "ol" if b.get("ordered") else "ul"
            li = []
            for i in (b.get("items") or []):
                # 항목은 문자열 또는 {text, url}(링크 — 목차·북마크 등)
                if isinstance(i, dict):
                    if i.get("url"):
                        li.append(f'<li><a href="{esc(i.get("url"))}">{esc(i.get("text"))}</a></li>')
                    else:
                        li.append(f"<li>{esc(i.get('text'))}</li>")
                else:
                    li.append(f"<li>{esc(i)}</li>")
            parts.append(f"<{tag}>{''.join(li)}</{tag}>")
        elif t == "image":
            src = b.get("src") or b.get("path") or ""
            cap = b.get("caption")
            cap_html = f"<figcaption>{esc(cap)}</figcaption>" if cap else ""
            parts.append(f'<figure><img src="{esc(src)}" alt="{esc(cap)}">{cap_html}</figure>')
        elif t == "table":
            # 데이터 통화 재사용: {columns, rows}
            cols = b.get("columns") or []
            rows = b.get("rows") or []
            thead = "".join(f"<th>{esc(c)}</th>" for c in cols)
            tbody = "".join("<tr>" + "".join(f"<td>{esc(c)}</td>" for c in r) + "</tr>"
                            for r in rows if isinstance(r, (list, tuple)))
            parts.append(f"<table><thead><tr>{thead}</tr></thead><tbody>{tbody}</tbody></table>")
        elif t == "quote":
            cite = b.get("cite")
            cite_html = f"<cite>— {esc(cite)}</cite>" if cite else ""
            parts.append(f"<blockquote>{esc(b.get('text'))}{cite_html}</blockquote>")
        elif t == "code":
            parts.append(f"<pre><code>{esc(b.get('text'))}</code></pre>")
        elif t == "cards":
            # 링크 달린 카드 그리드 — 뉴스/검색결과/북마크/상품목록 공용 문서 원시.
            # 각 item: {title, meta?, summary?, url?, link_label?}. columns(기본 2).
            ncol = max(1, min(4, int(b.get("columns") or 2)))
            cell = []
            for it in (b.get("items") or []):
                if not isinstance(it, dict):
                    continue
                img = it.get("image")
                img_h = f'<img class="card-img" src="{esc(img)}" alt="" loading="lazy">' if img else ""
                meta_h = f'<p class="card-meta">{esc(it.get("meta"))}</p>' if it.get("meta") else ""
                sum_h = f'<p class="card-sum">{esc(it.get("summary"))}</p>' if it.get("summary") else ""
                url = it.get("url")
                link_h = (f'<a href="{esc(url)}" target="_blank" rel="noopener">'
                          f'{esc(it.get("link_label") or "열기")}</a>') if url else ""
                cell.append(f'<div class="card">{img_h}<h3>{esc(it.get("title"))}</h3>{meta_h}{sum_h}{link_h}</div>')
            parts.append(f'<div class="cards" style="--cols:{ncol}">{"".join(cell)}</div>')
        elif t == "divider":
            parts.append("<hr>")
        else:  # paragraph (기본)
            parts.append(f"<p>{esc(b.get('text'))}</p>")
    return "\n".join(parts)


def _doc_css(theme: str) -> str:
    """문서 emitter의 <style> 본문을 테마별로. default(기사형) / newspaper(제호+카드 그리드)."""
    # 공통: 카드 그리드 골격(테마 무관 동일 구조, 색·여백만 테마가 덧칠)
    base_cards = """
.cards{display:grid;grid-template-columns:repeat(var(--cols,2),1fr);gap:18px;margin:1.2em 0}
@media(max-width:680px){.cards{grid-template-columns:1fr}}
.card{border:1px solid #e1e4e8;border-radius:10px;padding:16px 18px;display:flex;flex-direction:column;background:#fff}
.card img.card-img{width:100%;max-height:200px;object-fit:contain;border-radius:6px;margin-bottom:10px;background:#f5f5f5}
.card h3{margin:0 0 8px;font-size:1.08em;line-height:1.4;color:#22223b}
.card .card-meta{color:#888;font-size:0.82em;margin:0 0 8px}
.card .card-sum{color:#444;font-size:0.92em;margin:0 0 10px;flex:1}
.card a{margin-top:auto;color:#3d5a80;font-weight:bold;font-size:0.9em;text-decoration:none}
.card a:hover{text-decoration:underline}
"""
    if theme == "newspaper":
        return """
body{max-width:1100px;margin:30px auto;padding:0 16px;font-family:'Noto Sans KR','Pretendard',-apple-system,sans-serif;line-height:1.6;color:#333;background:#f0f2f5}
.docwrap{background:#fff;padding:40px;border-radius:12px;box-shadow:0 4px 20px rgba(0,0,0,0.08)}
h1{color:#1a1a2e;font-size:2.5em;margin:0 0 15px;border-bottom:4px solid #1a1a2e;padding-bottom:15px;text-align:center}
.doc-meta{text-align:center;color:#666;font-size:0.95em;margin-bottom:30px;background:#f8f9fa;padding:15px;border-radius:8px}
h2{color:#1a1a2e;font-size:1.8em;margin:40px 0 20px;padding-bottom:8px;border-bottom:2px solid #eee}
h3{color:#22223b}
img{max-width:100%;border-radius:8px} figure{margin:1.2em 0} figcaption{color:#666;font-size:0.9em;text-align:center}
table{border-collapse:collapse;width:100%;margin:1.2em 0} th,td{border:1px solid #ddd;padding:8px 12px;text-align:left} th{background:#f5f5f5}
blockquote{border-left:4px solid #ccc;margin:1.2em 0;padding:0.5em 1em;color:#555}
hr{border:none;border-top:1px solid #eee;margin:30px 0}
""" + base_cards
    # default(기사형 단일단)
    return """
body{max-width:760px;margin:40px auto;padding:0 20px;font-family:'Pretendard','Noto Sans KR',sans-serif;line-height:1.7;color:#1a1a1a}
.docwrap{}
h1,h2,h3,h4{line-height:1.3;margin:1.4em 0 0.5em} h1{font-size:2em}
.doc-meta{color:#666;font-size:0.95em;margin:-0.5em 0 1.5em}
img{max-width:100%;border-radius:8px} figure{margin:1.5em 0} figcaption{color:#666;font-size:0.9em;text-align:center;margin-top:0.5em}
table{border-collapse:collapse;width:100%;margin:1.5em 0} th,td{border:1px solid #ddd;padding:8px 12px;text-align:left} th{background:#f5f5f5}
blockquote{border-left:4px solid #ccc;margin:1.5em 0;padding:0.5em 1em;color:#555} blockquote cite{display:block;margin-top:0.5em;font-size:0.9em}
pre{background:#f6f8fa;padding:1em;border-radius:6px;overflow:auto} hr{border:none;border-top:1px solid #e0e0e0;margin:2em 0}
""" + base_cards


def _resolve_image_bytes(src: str):
    """이미지 src(로컬 경로/data URI/http URL)를 file-like(BytesIO)로 해소. 실패 시 None.
    docx·pptx emitter 공용 — 둘 다 파일/스트림만 받음(HTML <img>와 달리)."""
    import io
    import os
    import base64
    if not src:
        return None
    try:
        s = str(src).strip()
        if s.startswith("data:"):  # data:image/png;base64,....
            b64 = s.split(",", 1)[1]
            return io.BytesIO(base64.b64decode(b64))
        if s.startswith("file://"):
            s = s[7:]
        if s.startswith(("http://", "https://")):
            import urllib.request
            with urllib.request.urlopen(s, timeout=15) as r:
                return io.BytesIO(r.read())
        if os.path.isfile(s):
            with open(s, "rb") as f:
                return io.BytesIO(f.read())
    except Exception:
        return None
    return None


def _typ_esc(s) -> str:
    """typst 마크업 특수문자 이스케이프."""
    s = str(s if s is not None else "")
    for ch in ("\\", "#", "$", "*", "_", "`", "<", ">", "@", "[", "]"):
        s = s.replace(ch, "\\" + ch)
    return s


def _doc_blocks_to_typst(blocks: list, title: str, meta: str, out_path: str):
    """문서 IR → typst 컴파일 PDF (책 품질 조판). 산문·보고서에 최적 — 정렬·페이지·타이포가 강점.
    HTML theme/cards 그리드는 무시(조판 모델이 다름). 한글 = Apple SD Gothic Neo."""
    import os
    import subprocess
    import tempfile

    lines = [
        '#set text(font: "Apple SD Gothic Neo", size: 11pt, lang: "ko")',
        '#set page(paper: "a4", margin: (x: 2.2cm, y: 2.4cm), numbering: "1")',
        '#set par(justify: true, leading: 0.8em)',
        '#show heading: set block(above: 1.2em, below: 0.6em)',
        '#set heading(numbering: none)',
        "",
    ]
    if title:
        lines.append(f'#align(center)[#text(size: 22pt, weight: "bold")[{_typ_esc(title)}]]')
    if meta:
        lines.append(f'#align(center)[#text(size: 9pt, fill: gray)[{_typ_esc(meta)}]]')
    if title or meta:
        lines.append("#v(0.5em)")
        lines.append("")

    for b in blocks:
        if not isinstance(b, dict):
            continue
        t = (b.get("type") or "paragraph").lower()
        if t == "heading":
            lvl = max(1, min(5, int(b.get("level") or 2)))
            lines.append("=" * lvl + " " + _typ_esc(b.get("text")))
        elif t == "list":
            for it in (b.get("items") or []):
                txt = it.get("text") if isinstance(it, dict) else it
                lines.append("- " + _typ_esc(txt))
        elif t == "table":
            cols = b.get("columns") or []
            rows = [r for r in (b.get("rows") or []) if isinstance(r, (list, tuple))]
            ncol = max([len(cols)] + [len(r) for r in rows] or [0])
            if ncol:
                cells = []
                if cols:
                    cells += [f"[*{_typ_esc(c)}*]" for c in cols[:ncol]] + ["[]"] * (ncol - len(cols[:ncol]))
                for r in rows:
                    cells += [f"[{_typ_esc(v)}]" for v in r[:ncol]] + ["[]"] * (ncol - len(r[:ncol]))
                lines.append(f"#table(columns: {ncol}, " + ", ".join(cells) + ")")
        elif t == "cards":
            for it in (b.get("items") or []):
                if not isinstance(it, dict):
                    continue
                lines.append("=== " + _typ_esc(it.get("title")))
                if it.get("meta"):
                    lines.append(f'#text(size: 9pt, fill: gray)[{_typ_esc(it.get("meta"))}]')
                if it.get("summary"):
                    lines.append(_typ_esc(it.get("summary")))
                if it.get("url"):
                    lines.append(f'#link("{it.get("url")}")[{_typ_esc(it.get("link_label") or "열기")}]')
        elif t == "quote":
            cite = b.get("cite")
            q = f'#quote(block: true)[{_typ_esc(b.get("text"))}]'
            lines.append(q + (f" #text(size: 9pt, fill: gray)[— {_typ_esc(cite)}]" if cite else ""))
        elif t == "code":
            lines.append("```\n" + str(b.get("text") or "") + "\n```")
        elif t == "image":
            src = b.get("src") or b.get("path") or ""
            if src and os.path.isfile(str(src)):
                lines.append(f'#figure(image("{src}", width: 80%)' +
                             (f', caption: [{_typ_esc(b.get("caption"))}]' if b.get("caption") else "") + ")")
        elif t == "divider":
            lines.append("#line(length: 100%, stroke: 0.5pt + gray)")
        else:
            lines.append(_typ_esc(b.get("text")))
        lines.append("")

    typ_src = "\n".join(lines)
    with tempfile.NamedTemporaryFile("w", suffix=".typ", delete=False, encoding="utf-8") as tf:
        tf.write(typ_src)
        typ_file = tf.name
    try:
        proc = subprocess.run(["typst", "compile", typ_file, out_path],
                              capture_output=True, text=True, timeout=60)
        if proc.returncode != 0:
            raise RuntimeError(f"typst compile 실패: {proc.stderr[:300]}")
    finally:
        try:
            os.unlink(typ_file)
        except Exception:
            pass


def _add_hyperlink(paragraph, url: str, text: str):
    """python-docx 문단에 클릭 가능한 하이퍼링크 추가(네이티브 미지원이라 관계+XML 수작업)."""
    from docx.oxml.ns import qn
    from docx.oxml import OxmlElement
    try:
        part = paragraph.part
        r_id = part.relate_to(url, "http://schemas.openxmlformats.org/officeDocument/2006/relationships/hyperlink",
                              is_external=True)
        hyperlink = OxmlElement("w:hyperlink")
        hyperlink.set(qn("r:id"), r_id)
        new_run = OxmlElement("w:r")
        rPr = OxmlElement("w:rPr")
        color = OxmlElement("w:color"); color.set(qn("w:val"), "3D5A80"); rPr.append(color)
        u = OxmlElement("w:u"); u.set(qn("w:val"), "single"); rPr.append(u)
        new_run.append(rPr)
        t = OxmlElement("w:t"); t.text = text or url; new_run.append(t)
        hyperlink.append(new_run)
        paragraph._p.append(hyperlink)
    except Exception:
        paragraph.add_run(f"{text}: {url}" if text else url)


def _doc_blocks_to_docx(blocks: list, title: str, out_path: str):
    """문서 IR → .docx (python-docx). html emitter와 같은 IR을 소비.
    table 블록 = 데이터 통화 {columns,rows} 그대로 재사용."""
    from docx import Document
    from docx.shared import Pt, Inches
    from docx.oxml.ns import qn
    from docx.oxml import OxmlElement

    doc = Document()
    if title:
        doc.add_heading(str(title), level=0)

    for b in blocks:
        if not isinstance(b, dict):
            continue
        t = (b.get("type") or "paragraph").lower()
        if t == "heading":
            lvl = max(1, min(6, int(b.get("level") or 2)))
            doc.add_heading(str(b.get("text") or ""), level=lvl)
        elif t == "list":
            ordered = bool(b.get("ordered"))
            style = "List Number" if ordered else "List Bullet"
            for it in (b.get("items") or []):
                txt = it.get("text") if isinstance(it, dict) else it
                doc.add_paragraph(str(txt), style=style)
        elif t == "image":
            stream = _resolve_image_bytes(b.get("src") or b.get("path") or "")
            if stream is not None:
                try:
                    doc.add_picture(stream, width=Inches(6.0))
                except Exception:
                    pass
            cap = b.get("caption")
            if cap:
                p = doc.add_paragraph(str(cap))
                p.italic = True
        elif t == "table":
            cols = b.get("columns") or []
            rows = [r for r in (b.get("rows") or []) if isinstance(r, (list, tuple))]
            ncol = max([len(cols)] + [len(r) for r in rows] or [0])
            if ncol:
                tbl = doc.add_table(rows=0, cols=ncol)
                try:
                    tbl.style = "Light Grid Accent 1"
                except Exception:
                    pass
                if cols:
                    hdr = tbl.add_row().cells
                    for i, c in enumerate(cols[:ncol]):
                        hdr[i].text = str(c)
                    for cell in hdr:
                        for para in cell.paragraphs:
                            for run in para.runs:
                                run.bold = True
                for r in rows:
                    cells = tbl.add_row().cells
                    for i, v in enumerate(r[:ncol]):
                        cells[i].text = str(v)
        elif t == "quote":
            p = doc.add_paragraph(str(b.get("text") or ""))
            try:
                p.style = "Intense Quote"
            except Exception:
                pass
            cite = b.get("cite")
            if cite:
                doc.add_paragraph(f"— {cite}")
        elif t == "code":
            p = doc.add_paragraph()
            run = p.add_run(str(b.get("text") or ""))
            run.font.name = "Courier New"
            run.font.size = Pt(9.5)
        elif t == "cards":
            for it in (b.get("items") or []):
                if not isinstance(it, dict):
                    continue
                if it.get("image"):  # 썸네일(표지 등) — 다운로드 실패는 graceful
                    stream = _resolve_image_bytes(it.get("image"))
                    if stream is not None:
                        try:
                            doc.add_picture(stream, width=Inches(1.6))
                        except Exception:
                            pass
                doc.add_heading(str(it.get("title") or ""), level=3)
                if it.get("meta"):
                    mp = doc.add_paragraph()
                    mr = mp.add_run(str(it.get("meta")))
                    mr.italic = True
                    mr.font.size = Pt(9)
                if it.get("summary"):
                    doc.add_paragraph(str(it.get("summary")))
                if it.get("url"):
                    _add_hyperlink(doc.add_paragraph(), str(it.get("url")),
                                   str(it.get("link_label") or "열기"))
        elif t == "divider":
            p = doc.add_paragraph()
            pPr = p._p.get_or_add_pPr()
            pbdr = OxmlElement("w:pBdr")
            bottom = OxmlElement("w:bottom")
            bottom.set(qn("w:val"), "single")
            bottom.set(qn("w:sz"), "6")
            bottom.set(qn("w:space"), "1")
            bottom.set(qn("w:color"), "auto")
            pbdr.append(bottom)
            pPr.append(pbdr)
        else:  # paragraph
            doc.add_paragraph(str(b.get("text") or ""))

    doc.save(out_path)


def _doc_blocks_to_pptx(blocks: list, title: str, out_path: str):
    """문서 IR → .pptx (python-pptx). ★종류 경계 주의: 슬라이드 IR이 아니라 *문서 IR을 슬라이드로 투영*.
    문서 IR이 정본, pptx는 emitter일 뿐 — heading(level≤2)이 새 슬라이드, 그 아래 내용이 글머리표.
    슬라이드 전용 시각 레이아웃이 필요하면 engines:slide(슬라이드 IR)를 써야지 이걸 쓰면 안 됨."""
    from pptx import Presentation
    from pptx.util import Inches, Pt

    prs = Presentation()
    SW, SH = prs.slide_width, prs.slide_height
    blank = prs.slide_layouts[6]
    title_layout = prs.slide_layouts[0]

    # 표지 슬라이드(title 있으면)
    if title:
        s = prs.slides.add_slide(title_layout)
        s.shapes.title.text = str(title)

    state = {"slide": None, "body": None}

    def new_content_slide(heading_text=""):
        # "제목+내용" 레이아웃(1) — 제목 placeholder + 본문 placeholder
        s = prs.slides.add_slide(prs.slide_layouts[1])
        s.shapes.title.text = str(heading_text or "")
        body_tf = None
        for ph in s.placeholders:
            if ph.placeholder_format.idx == 1:
                body_tf = ph.text_frame
                break
        if body_tf is not None:
            body_tf.clear()
            body_tf.word_wrap = True
        state["slide"], state["body"] = s, body_tf

    def add_bullet(text, level=0, italic=False, mono=False):
        if state["body"] is None:
            new_content_slide("")
        tf = state["body"]
        # clear()가 남긴 빈 첫 문단 재사용, 이후엔 add
        if len(tf.paragraphs) == 1 and not tf.paragraphs[0].runs:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = str(text)
        p.level = min(4, max(0, level))
        for run in p.runs:
            if italic:
                run.font.italic = True
            if mono:
                run.font.name = "Courier New"
                run.font.size = Pt(14)

    for b in blocks:
        if not isinstance(b, dict):
            continue
        t = (b.get("type") or "paragraph").lower()
        if t == "heading":
            lvl = int(b.get("level") or 2)
            if lvl <= 2 or state["body"] is None:
                # 큰 섹션, 또는 표·구분선 뒤 첫 제목이면 새 슬라이드 제목으로
                new_content_slide(b.get("text") or "")
            else:  # 진행 중 슬라이드의 하위 섹션 = 글머리표
                add_bullet(b.get("text") or "", level=0)
        elif t == "paragraph":
            add_bullet(b.get("text") or "", level=0)
        elif t == "list":
            for it in (b.get("items") or []):
                add_bullet(it.get("text") if isinstance(it, dict) else it, level=1)
        elif t == "quote":
            txt = str(b.get("text") or "")
            cite = b.get("cite")
            add_bullet(f"“{txt}”" + (f" — {cite}" if cite else ""), level=1, italic=True)
        elif t == "code":
            add_bullet(b.get("text") or "", level=1, mono=True)
        elif t == "cards":
            for it in (b.get("items") or []):
                if not isinstance(it, dict):
                    continue
                add_bullet(it.get("title") or "", level=0)
                sub = " / ".join(x for x in [it.get("meta"), it.get("summary")] if x)
                if sub:
                    add_bullet(sub, level=1)
        elif t == "image":
            stream = _resolve_image_bytes(b.get("src") or b.get("path") or "")
            if stream is not None:
                s = prs.slides.add_slide(blank)
                try:
                    s.shapes.add_picture(stream, Inches(0.6), Inches(0.6), width=SW - Inches(1.2))
                except Exception:
                    pass
                cap = b.get("caption")
                if cap:
                    tb = s.shapes.add_textbox(Inches(0.6), SH - Inches(0.9), SW - Inches(1.2), Inches(0.7))
                    tb.text_frame.text = str(cap)
                state["slide"], state["body"] = None, None  # 이미지 후 새 내용 슬라이드 강제
        elif t == "table":
            cols = b.get("columns") or []
            rows = [r for r in (b.get("rows") or []) if isinstance(r, (list, tuple))]
            ncol = max([len(cols)] + [len(r) for r in rows] or [0])
            nrow = len(rows) + (1 if cols else 0)
            if ncol and nrow:
                s = prs.slides.add_slide(blank)
                gt = s.shapes.add_table(nrow, ncol, Inches(0.5), Inches(0.6),
                                        SW - Inches(1.0), Inches(0.4) * nrow).table
                ri = 0
                if cols:
                    for ci, c in enumerate(cols[:ncol]):
                        gt.cell(0, ci).text = str(c)
                    ri = 1
                for r in rows:
                    for ci, v in enumerate(r[:ncol]):
                        gt.cell(ri, ci).text = str(v)
                    ri += 1
                state["slide"], state["body"] = None, None
        elif t == "divider":
            state["slide"], state["body"] = None, None  # 새 슬라이드 경계

    if len(prs.slides) == 0:  # 표지도 내용도 없으면 빈 슬라이드 하나
        prs.slides.add_slide(blank)
    prs.save(out_path)


# ── B: 구조화 원자 — 콘텐츠 → 문서 IR (A획득→B구조화→IR→emit 파이프라인) ──
_STRUCTURE_PROMPT = """당신은 콘텐츠를 깔끔한 문서 구조로 정리하는 편집자입니다. 주어진 내용을 문서 IR(JSON)로 변환합니다.

출력은 JSON 한 객체만: {"title": "...", "blocks": [ ... ]}
블록 타입:
- {"type":"heading","level":2,"text":"..."}   (level 1~4)
- {"type":"paragraph","text":"..."}
- {"type":"list","ordered":false,"items":["...","..."]}
- {"type":"table","columns":["...","..."],"rows":[["...","..."]]}
- {"type":"quote","text":"...","cite":"..."}
- {"type":"code","text":"...","lang":"..."}
- {"type":"divider"}

원칙: 내용을 지어내지 말고 주어진 것에서만. title=핵심을 담은 명제. 긴 글은 heading으로 섹션화, 나열은 list, 비교·수치는 table. JSON 외 텍스트 금지."""


def structure_document(tool_input, output_base="."):
    """[table:structure] — 원본 콘텐츠를 문서 IR(blocks)로 구조화 (LLM 편집자).

    파라미터: content(필수, 원본 텍스트) · instruction(선택, 정리 방향).
    반환: {success, title, blocks, block_count}. render_document로 이어 렌더(>> 파이프 지원).
    """
    import json as _json

    content = (tool_input.get("content") or "").strip()
    # >> 파이프: 이전 액션의 텍스트 결과를 content로 받음
    if not content:
        pr = tool_input.get("_prev_result")
        if isinstance(pr, str):
            content = pr.strip()
        elif isinstance(pr, dict):
            content = str(pr.get("summary") or pr.get("content") or pr.get("text") or "").strip()
    if not content:
        return _json.dumps({"success": False, "message": "content(구조화할 원본 내용)가 필요합니다."},
                           ensure_ascii=False)

    instruction = (tool_input.get("instruction") or "").strip()
    user = f"# 정리할 내용\n{content[:16000]}"
    if instruction:
        user += f"\n\n# 정리 방향\n{instruction}"
    user += "\n\n위 내용을 문서 IR(JSON 한 객체)로 출력하라."

    try:
        from consciousness_agent import lightweight_ai_call
        resp = lightweight_ai_call(user, system_prompt=_STRUCTURE_PROMPT)
    except Exception as e:
        return _json.dumps({"success": False, "message": f"구조화 AI 호출 실패: {e}"}, ensure_ascii=False)
    if not resp or not resp.strip():
        return _json.dumps({"success": False, "message": "구조화 AI 응답 없음"}, ensure_ascii=False)

    txt = resp.strip()
    if txt.startswith("```"):
        txt = txt.strip("`")
        if txt.lower().startswith("json"):
            txt = txt[4:]
        txt = txt.split("```")[0].strip()
    try:
        a, b = txt.find("{"), txt.rfind("}")
        ir = _json.loads(txt[a:b + 1])
    except Exception as e:
        return _json.dumps({"success": False, "message": f"IR JSON 파싱 실패: {e}", "raw": resp[:300]},
                           ensure_ascii=False)
    blocks = ir.get("blocks")
    if not isinstance(blocks, list) or not blocks:
        return _json.dumps({"success": False, "message": "blocks가 없습니다.", "raw": resp[:300]},
                           ensure_ascii=False)
    return _json.dumps({"success": True, "title": ir.get("title", ""), "blocks": blocks,
                        "block_count": len(blocks),
                        "message": f"{len(blocks)}블록 문서 IR로 구조화."}, ensure_ascii=False)


def render_document(tool_input, output_base="."):
    """문서 IR → 산출물. 현재 emitter: html (단일 IR, 향후 pdf/docx/pptx emitter 추가).

    파라미터: blocks(필수, IR 블록 배열) · title(선택) · format(기본 html) · filename(선택).
    반환: {success, path, format, blocks}.
    """
    import os
    import html as _html
    import json as _json

    blocks = tool_input.get("blocks")
    if not blocks:
        # >> 파이프: 이전 생산자 결과(_prev_result)의 blocks·title·meta·theme 자동 수용
        pr = tool_input.get("_prev_result")
        if pr:
            try:
                po = _json.loads(pr) if isinstance(pr, str) else pr
                if isinstance(po, dict):
                    _rows = po.get("items")
                    if po.get("blocks"):
                        blocks = po["blocks"]
                    elif isinstance(_rows, list) and _rows:
                        if isinstance(_rows[0], dict) and "type" in _rows[0] and "text" in _rows[0]:
                            # 문서 IR items(type+text — crawl·read 등) = blocks 그 자체(산문).
                            blocks = _rows
                        else:
                            # 단일 통화 items([{title,meta,summary,url,image}]) → cards 블록으로 래핑.
                            blocks = [{"type": "cards", "columns": 2, "items": _rows}]
                    for k in ("title", "meta", "theme"):
                        if not tool_input.get(k) and po.get(k):
                            tool_input[k] = po[k]
            except Exception:
                pass
    if isinstance(blocks, str):
        try:
            blocks = _json.loads(blocks)
        except Exception:
            blocks = None
    if not isinstance(blocks, list) or not blocks:
        return _json.dumps({"success": False, "message": "blocks(문서 IR 블록 배열)가 필요합니다."},
                           ensure_ascii=False)

    title = tool_input.get("title") or ""
    meta = tool_input.get("meta") or ""
    theme = (tool_input.get("theme") or "default").strip().lower()
    fmt = (tool_input.get("format") or "html").strip().lower()
    note = ""
    if fmt not in ("html", "pdf", "png", "docx", "pptx", "typst"):
        note = f" (format '{fmt}' 미지원 — html로 산출)"
        fmt = "html"

    os.makedirs(output_base, exist_ok=True)
    base = tool_input.get("filename") or "document"
    base = os.path.splitext(os.path.basename(str(base)))[0] or "document"

    # typst emitter — 책 품질 조판 PDF(산문·보고서). HTML theme/cards 그리드는 무시(조판 모델 상이).
    if fmt == "typst":
        try:
            out_path = os.path.join(output_base, f"{base}.pdf")
            _doc_blocks_to_typst(blocks, title, meta, out_path)
            return _json.dumps({"success": True, "path": out_path, "file": out_path,
                                "title": title, "format": "typst_pdf", "blocks": len(blocks),
                                "message": f"문서 {len(blocks)}블록을 typst 책 품질 PDF로 조판했습니다."},
                               ensure_ascii=False)
        except Exception as e:
            note = f" (typst 조판 실패 → 브라우저 PDF 폴백: {e})"
            fmt = "pdf"

    # docx/pptx emitter — 같은 문서 IR을 사무 포맷으로. pptx는 문서 IR을 슬라이드로 *투영*(종류 경계 주의).
    if fmt in ("docx", "pptx"):
        try:
            out_path = os.path.join(output_base, f"{base}.{fmt}")
            if fmt == "docx":
                _doc_blocks_to_docx(blocks, title, out_path)
            else:
                _doc_blocks_to_pptx(blocks, title, out_path)
            return _json.dumps({"success": True, "path": out_path, "file": out_path,
                                "title": title, "format": fmt, "blocks": len(blocks),
                                "message": f"문서 {len(blocks)}블록을 {fmt.upper()}로 렌더했습니다."},
                               ensure_ascii=False)
        except Exception as e:
            note = f" ({fmt} 렌더 실패 → HTML 폴백: {e})"
            fmt = "html"

    body = _doc_blocks_to_html(blocks)
    title_h = f"<h1>{_html.escape(str(title))}</h1>" if title else ""
    meta_h = f'<div class="doc-meta">{_html.escape(str(meta))}</div>' if meta else ""
    doc = f"""<!DOCTYPE html><html lang="ko"><head><meta charset="utf-8">
<title>{_html.escape(str(title))}</title><style>{_doc_css(theme)}</style></head><body>
<div class="docwrap">
{title_h}
{meta_h}
{body}
</div>
</body></html>"""

    # 같은 문서 IR을 여러 emitter로 — html/pdf/png. pdf·png는 Playwright로 동일 HTML 렌더.
    if fmt in ("pdf", "png"):
        try:
            from playwright.sync_api import sync_playwright
            out_path = os.path.join(output_base, f"{base}.{fmt}")
            with sync_playwright() as pw:
                br = pw.chromium.launch()
                pg = br.new_page(viewport={"width": 900, "height": 1200})
                pg.set_content(doc, wait_until="networkidle")
                pg.wait_for_timeout(300)
                if fmt == "pdf":
                    pg.pdf(path=out_path, format="A4", print_background=True,
                           margin={"top": "20mm", "bottom": "20mm", "left": "16mm", "right": "16mm"})
                else:
                    pg.screenshot(path=out_path, full_page=True)
                br.close()
            return _json.dumps({"success": True, "path": out_path, "file": out_path,
                                "title": title, "format": fmt, "blocks": len(blocks),
                                "message": f"문서 {len(blocks)}블록을 {fmt.upper()}로 렌더했습니다."},
                               ensure_ascii=False)
        except Exception as e:
            # emitter 실패 시 HTML로 폴백(산출 보존)
            note = f" ({fmt} 렌더 실패 → HTML 폴백: {e})"

    out_path = os.path.join(output_base, f"{base}.html")
    with open(out_path, "w", encoding="utf-8") as f:
        f.write(doc)
    return _json.dumps({"success": True, "path": out_path, "file": out_path,
                        "title": title, "format": "html", "blocks": len(blocks),
                        # 렌더된 HTML을 결과에 동봉 — 액션이 다른 몸(맥)으로 포워드돼 파일이
                        # 거기 생겨도, 호출한 몸(폰)이 파일 위치 의존 없이 콘텐츠로 바로 띄운다.
                        "html": doc,
                        "message": f"문서 {len(blocks)}블록을 HTML로 렌더했습니다.{note}"},
                       ensure_ascii=False)


def render_html_to_image(tool_input, output_base="."):
    """HTML을 이미지로 렌더링

    로컬 이미지 사용 방법:
    1. base_path 옵션: HTML 내 상대 경로의 기준 디렉토리 지정
       - base_path="/Users/.../project" 설정 시 <img src="images/photo.jpg">가 작동
    2. 절대 경로: <img src="file:///Users/.../image.jpg">
    3. Base64: <img src="data:image/png;base64,...">
    """
    from playwright.sync_api import sync_playwright
    import re
    import tempfile

    html = tool_input.get("html")
    output_path = tool_input.get("output_path")
    width = tool_input.get("width", 1280)
    height = tool_input.get("height", 720)
    selector = tool_input.get("selector")
    base_path = tool_input.get("base_path")  # 로컬 파일 기준 경로

    if not html:
        return "오류: html은 필수입니다."

    # output_path가 지정되면 파일명만 추출하여 output_base에 저장
    if output_path:
        filename = os.path.basename(output_path)
        output_path = os.path.join(output_base, filename)
    else:
        output_path = os.path.join(output_base, f"rendered_{uuid.uuid4().hex[:8]}.png")

    os.makedirs(os.path.dirname(os.path.abspath(output_path)), exist_ok=True)

    try:
        with sync_playwright() as p:
            browser = p.chromium.launch()
            page = browser.new_page(viewport={"width": width, "height": height})

            # base_path가 있으면 임시 HTML 파일 생성 후 file:// 프로토콜로 로드
            if base_path:
                base_path = os.path.abspath(base_path)

                # HTML 내 상대 경로 이미지를 Base64로 변환
                def replace_img_src(match):
                    src = match.group(1)
                    # 이미 data: 또는 http로 시작하면 스킵
                    if src.startswith(('data:', 'http://', 'https://', 'file://')):
                        return match.group(0)
                    # 상대 경로를 절대 경로로 변환
                    if not os.path.isabs(src):
                        abs_path = os.path.join(base_path, src)
                    else:
                        abs_path = src
                    # Base64로 변환
                    base64_data = get_image_base64(abs_path)
                    if base64_data:
                        return f'src="{base64_data}"'
                    return match.group(0)  # 변환 실패시 원본 유지

                # src="..." 패턴 찾아서 변환
                html = re.sub(r'src=["\']([^"\']+)["\']', replace_img_src, html)

                # CSS background-image의 url()도 처리
                def replace_bg_url(match):
                    url = match.group(1)
                    if url.startswith(('data:', 'http://', 'https://', 'file://')):
                        return match.group(0)
                    if not os.path.isabs(url):
                        abs_path = os.path.join(base_path, url)
                    else:
                        abs_path = url
                    base64_data = get_image_base64(abs_path)
                    if base64_data:
                        return f'url("{base64_data}")'
                    return match.group(0)

                html = re.sub(r'url\(["\']?([^"\')\s]+)["\']?\)', replace_bg_url, html)

            # 스타일 추가: 여백 제거 및 크기 강제
            html_with_style = f"""
            <style>
                body, html {{ margin: 0; padding: 0; width: {width}px; height: {height}px; overflow: hidden; }}
            </style>
            {html}
            """
            page.set_content(html_with_style)

            # 이미지 로딩 대기
            page.wait_for_load_state("networkidle")
            page.wait_for_timeout(500)  # Tailwind/폰트 로딩 대기

            if selector:
                element = page.query_selector(selector)
                if element:
                    element.screenshot(path=output_path)
                else:
                    browser.close()
                    return f"오류: 셀렉터 '{selector}'를 찾을 수 없습니다."
            else:
                page.screenshot(path=output_path)
            browser.close()
        # 절대 경로로 변환하여 반환 (에이전트 간 경로 혼동 방지)
        return f"렌더링 완료: {os.path.abspath(output_path)}"
    except Exception as e:
        return f"렌더링 중 오류 발생: {str(e)}"

def generate_ai_image(tool_input, output_base):
    import urllib.parse
    import httpx
    prompt = tool_input.get("prompt")
    if not prompt:
        return "오류: prompt는 필수입니다."
    output_path = tool_input.get("output_path")
    width = min(tool_input.get("width", 1024), 2048)
    height = min(tool_input.get("height", 1024), 2048)
    model = tool_input.get("model", "flux")
    seed = tool_input.get("seed")
    # output_path가 지정되면 파일명만 추출하여 output_base에 저장
    if output_path:
        filename = os.path.basename(output_path)
        output_path = os.path.join(output_base, filename)
    else:
        output_path = os.path.join(output_base, f"ai_image_{uuid.uuid4().hex[:8]}.png")
    os.makedirs(os.path.dirname(os.path.abspath(output_path)), exist_ok=True)
    encoded_prompt = urllib.parse.quote(prompt)
    url = f"https://image.pollinations.ai/prompt/{encoded_prompt}?width={width}&height={height}&model={model}&nologo=true"
    if seed is not None:
        url += f"&seed={seed}"
    try:
        with httpx.Client(timeout=120.0, follow_redirects=True) as client:
            response = client.get(url)
            response.raise_for_status()
            with open(output_path, "wb") as f:
                f.write(response.content)
        # 절대 경로로 변환하여 반환 (에이전트 간 경로 혼동 방지)
        return f"AI 이미지 생성 완료: {os.path.abspath(output_path)}\n프롬프트: {prompt}"
    except Exception as e:
        return f"이미지 생성 중 오류 발생: {str(e)}"


STYLE_PRESETS = {
    "vintage_book": (
        "Hand-drawn pen and ink illustration on aged beige parchment paper, "
        "two-tone palette of deep navy blue (#2c3e6f) and rust brown (#a55a3e), "
        "fine cross-hatching, geometric grid background lines, subtle paper texture and grain, "
        "vintage scientific manuscript aesthetic, balanced empty space around the central subject, "
        "centered composition with breathing room. "
        "Do NOT include any Korean or Hangul characters. Do NOT add decorative Latin text, ciphers, or unreadable script. "
        "Only the subject illustration — minimal or no text inside the image."
    ),
    "academic_paper": (
        "Clean academic diagram on bright white paper, monochrome with deep navy (#161c2a) line art and crimson (#8b1a1a) accent, "
        "thin precise pen lines, minimal shading, scholarly figure illustration style, generous white space. "
        "Do NOT include Korean or Hangul characters. Only English labels if any text is needed."
    ),
    "tech_minimal": (
        "Minimal vector-style illustration on dark navy background (#0d0f17), "
        "neon cyan (#1ce0ff) accent lines, thin geometric strokes, subtle glow, "
        "Linear/Vercel design aesthetic, isometric or flat composition, generous negative space. "
        "Do NOT include Korean or Hangul characters."
    ),
    "magazine_modern": (
        "Bold editorial illustration in New Yorker / Wired magazine style, "
        "high contrast black ink on pure white background with selective vivid red (#e6182b) accent, "
        "confident brush strokes, modernist composition, ample white space. "
        "Do NOT include Korean or Hangul characters."
    ),
    "sf_blueprint": (
        "Sci-fi HUD blueprint infographic on a deep navy (#050d1a) background, "
        "glowing cyan (#1ad3ff) and pale ice-blue (#6ee0ff) line art, thin precise pen strokes with subtle outer glow, "
        "wireframe / x-ray rendering of the subject, surrounded by faint technical schematic grid lines and HUD frame brackets in the corners, "
        "small English technical labels and arrow annotations allowed (e.g. 'EYE', 'MASS', 'F=ma', 'FAILURE STATE'), "
        "NotebookLM diagram aesthetic, holographic / hologram feel, "
        "composition deliberately leaves clear empty regions on the sides or bottom so Korean caption boxes can be overlaid later without occluding the subject. "
        "Strictly NO Korean / Hangul characters anywhere in the image. "
        "Avoid decorative gibberish Latin text — only meaningful English labels if any."
    ),
}


def _build_image_prompt(user_prompt: str, style_preset: str = None) -> str:
    """스타일 프리셋을 사용자 프롬프트와 결합. 디자인 시스템과 어울리는 일러스트 생성용."""
    if not style_preset or style_preset == "default":
        return user_prompt
    style = STYLE_PRESETS.get(style_preset)
    if not style:
        return user_prompt
    return f"{user_prompt}\n\n--- STYLE GUIDELINES ---\n{style}"


def critique_gemini_image(tool_input, output_base):
    """Gemini Vision으로 이미지를 평가한다 — 일러스트가 의도와 정합하는지 검증.

    파라미터:
      - image_path (필수): 평가할 이미지 절대 경로 (또는 base64 data URI)
      - intent (필수): "이 일러스트가 무엇을 표현해야 하는가" — 자연어 설명
      - checks (선택): 추가 체크 리스트 (예: ["다이어그램형인가", "한글이 들어갔는가", "오른쪽 1/3이 비어있는가"])
      - style_preset (선택): 디자인 시스템 톤 일관성 평가 기준 (sf_blueprint 등)
      - api_key (선택)

    반환:
      JSON 문자열 — {"passed": bool, "score": 0-10, "issues": [...], "notes": "..."}
      그리고 사람이 읽을 수 있는 요약 텍스트가 앞에 옴.
    """
    import httpx
    import base64
    import json as _json

    # path 는 IBL 표준 파라미터(self:read/grep/edit 모두 path) — image_path 미지정 시 폴백 수용.
    image_path = tool_input.get("image_path") or tool_input.get("path")
    intent = tool_input.get("intent", "")
    if not image_path:
        return "오류: image_path(또는 path)가 필요합니다."
    if not intent:
        return "오류: intent(이 일러스트가 무엇을 표현해야 하는지)가 필요합니다."

    api_key = tool_input.get("api_key") or os.environ.get("GEMINI_API_KEY")
    if not api_key:
        return "오류: GEMINI_API_KEY가 설정되지 않았습니다."

    # 이미지 로드 (절대 경로 또는 data URI)
    if image_path.startswith("data:"):
        # data URI → base64 페이로드만 추출
        try:
            _, b64data = image_path.split(",", 1)
            mime = image_path.split(";", 1)[0].split(":", 1)[1]
        except Exception:
            return "오류: 잘못된 data URI."
    else:
        if not os.path.exists(image_path):
            return f"오류: 파일이 없습니다: {image_path}"
        with open(image_path, "rb") as f:
            b64data = base64.b64encode(f.read()).decode("utf-8")
        ext = os.path.splitext(image_path)[1].lower()
        mime = {".png": "image/png", ".jpg": "image/jpeg", ".jpeg": "image/jpeg", ".webp": "image/webp"}.get(ext, "image/png")

    checks = tool_input.get("checks") or []
    style_preset = tool_input.get("style_preset", "")
    # preset: "slide_illustration"(기본, 현행 슬라이드 일러스트 체크) | "general"(임의 산출물 범용)
    preset = (tool_input.get("preset") or "slide_illustration").strip().lower()

    if preset == "general":
        default_checks = [
            "이미지가 의도(intent)를 정확하고 충분히 충족하는가?",
            "시각적 결함(텍스트 잘림·겹침, 레이아웃 불균형, 깨짐, 저해상도, 빈 공간 과다)이 없는가?",
        ]
        if style_preset:
            default_checks.append(f"스타일/톤이 '{style_preset}'와 일관되는가?")
        intro = "당신은 산출물 품질 평가자입니다. 아래 이미지가 의도를 잘 충족하는지 엄격하게 평가하세요."
        hard_rule = ""
    else:
        default_checks = [
            "이 일러스트는 회화적 '씬(scene)'이 아니라 정보를 전달하는 '다이어그램/인포그래픽'인가? (NotebookLM 양식)",
            "한글(Hangul) 문자가 일러스트 안에 들어가 있는가? (있으면 실패 — 한글은 텍스트 레이어에서 처리)",
            "라벨 박스가 들어갈 빈 공간(여백)이 의도된 위치에 정말 비어 있는가? (intent에 명시된 빈 공간 위치 확인)",
            "주요 객체가 일러스트의 핵심 영역(중앙/지정 위치)에 명확하게 배치되어 있는가?",
        ]
        if style_preset:
            default_checks.append(f"디자인 시스템 톤이 '{style_preset}'와 일관되는가? (색·선·분위기)")
        intro = "당신은 강의 슬라이드 일러스트 평가자입니다. 다음 일러스트가 의도를 잘 표현하는지 엄격하게 평가하세요."
        hard_rule = " 한글이 일러스트에 들어가 있으면 무조건 passed=false."

    all_checks = default_checks + checks

    instruction = (
        intro + "\n\n"
        f"**의도 (이 결과물이 충족해야 할 것)**:\n{intent}\n\n"
        f"**체크 항목** ({len(all_checks)}개):\n"
        + "\n".join(f"{i+1}. {c}" for i, c in enumerate(all_checks))
        + "\n\n반드시 다음 JSON 형식 한 개만 출력하세요. 다른 텍스트 금지.\n"
        "```json\n"
        "{\n"
        '  "passed": true|false,\n'
        '  "score": 0-10,\n'
        '  "issues": ["체크 N번 실패 — 구체적 이유", ...],\n'
        '  "notes": "전반적 평가 (1~2문장)"\n'
        "}\n"
        "```\n"
        "passed는 issues가 없거나 score>=7일 때 true." + hard_rule
    )

    payload = {
        "contents": [{
            "parts": [
                {"inlineData": {"mimeType": mime, "data": b64data}},
                {"text": instruction},
            ]
        }],
        "generationConfig": {"temperature": 0.2, "responseModalities": ["TEXT"]},
    }

    # VLM 모델 — pro 우선, fallback flash
    model = tool_input.get("model") or "gemini-3-pro-preview"
    url = f"https://generativelanguage.googleapis.com/v1beta/models/{model}:generateContent"

    try:
        with httpx.Client(timeout=120.0) as client:
            r = client.post(url, params={"key": api_key}, json=payload,
                            headers={"Content-Type": "application/json"})
            if r.status_code == 404:
                # 폴백 모델
                model = "gemini-2.5-pro"
                url = f"https://generativelanguage.googleapis.com/v1beta/models/{model}:generateContent"
                r = client.post(url, params={"key": api_key}, json=payload,
                                headers={"Content-Type": "application/json"})
            r.raise_for_status()
            data = r.json()
    except Exception as e:
        return f"오류: VLM 호출 실패: {e}"

    parts = (data.get("candidates", [{}])[0].get("content", {}) or {}).get("parts", [])
    text = "".join(p.get("text", "") for p in parts).strip()
    # ```json ... ``` 코드 펜스 제거
    if text.startswith("```"):
        text = text.strip("`")
        if text.lower().startswith("json"):
            text = text[4:].strip()
        text = text.split("```")[0].strip()

    try:
        verdict = _json.loads(text)
    except Exception:
        return f"VLM 응답 파싱 실패 — 원문:\n{text[:500]}"

    summary_lines = [
        f"이미지 평가: {image_path}",
        f"의도: {intent[:80]}{'...' if len(intent) > 80 else ''}",
        f"평가 결과: {'✓ 통과' if verdict.get('passed') else '✗ 실패'} (score={verdict.get('score', '?')}/10)",
    ]
    issues = verdict.get("issues") or []
    if issues:
        summary_lines.append("문제점:")
        summary_lines.extend(f"  - {i}" for i in issues)
    if verdict.get("notes"):
        summary_lines.append(f"메모: {verdict['notes']}")
    summary_lines.append("")
    summary_lines.append(f"verdict_json: {_json.dumps(verdict, ensure_ascii=False)}")
    return "\n".join(summary_lines)


def read_gemini_image(tool_input, output_base):
    """Gemini Vision으로 이미지를 *읽어* 질문에 자유서술로 답한다 (시각 QA / OCR / 검증).

    critique_gemini_image 와 다른 점: 합격/점수 채점이 아니라, 주어진 질문에 대한 자유 텍스트
    답을 돌려준다. "스크린샷의 숫자를 읽어줘", "이 그림에 무엇이 보이나", "이 통계 수치가
    332인가 136인가" 같은 시각 읽기·검증에 쓴다. 생성 일러스트 품질 평가는 image_critic 을 쓸 것.

    파라미터:
      - image_path (또는 path): 읽을 이미지 절대 경로 또는 base64 data URI (필수)
      - question (또는 query/intent/prompt): 무엇을 읽거나 답할지 (없으면 전체 묘사)
      - api_key (선택), model (선택)
    """
    import httpx
    import base64

    image_path = tool_input.get("image_path") or tool_input.get("path")
    if not image_path:
        return "오류: image_path(또는 path)가 필요합니다."
    question = (tool_input.get("question") or tool_input.get("query")
               or tool_input.get("intent") or tool_input.get("prompt") or "").strip()

    api_key = tool_input.get("api_key") or os.environ.get("GEMINI_API_KEY")
    if not api_key:
        return "오류: GEMINI_API_KEY가 설정되지 않았습니다."

    # 이미지 로드 (절대 경로 또는 data URI) — critique_gemini_image 와 동일.
    if image_path.startswith("data:"):
        try:
            _, b64data = image_path.split(",", 1)
            mime = image_path.split(";", 1)[0].split(":", 1)[1]
        except Exception:
            return "오류: 잘못된 data URI."
    else:
        if not os.path.exists(image_path):
            return f"오류: 파일이 없습니다: {image_path}"
        with open(image_path, "rb") as f:
            b64data = base64.b64encode(f.read()).decode("utf-8")
        ext = os.path.splitext(image_path)[1].lower()
        mime = {".png": "image/png", ".jpg": "image/jpeg", ".jpeg": "image/jpeg",
                ".webp": "image/webp", ".gif": "image/gif"}.get(ext, "image/png")

    if question:
        instruction = (
            "당신은 이미지를 정확히 읽는 시각 분석가입니다. 아래 이미지를 보고 질문에 "
            "사실에 근거해 답하세요. 이미지에 적힌 텍스트·숫자는 보이는 그대로 정확히 옮기고, "
            "보이지 않거나 불확실하면 추측하지 말고 그렇다고 밝히세요.\n\n"
            f"**질문**: {question}"
        )
    else:
        instruction = ("아래 이미지를 보고 무엇이 있는지 상세히 묘사하세요. 이미지에 적힌 "
                       "텍스트·숫자가 있으면 보이는 그대로 정확히 옮기세요.")

    payload = {
        "contents": [{
            "parts": [
                {"inlineData": {"mimeType": mime, "data": b64data}},
                {"text": instruction},
            ]
        }],
        "generationConfig": {"temperature": 0.0, "responseModalities": ["TEXT"]},
    }
    model = tool_input.get("model") or "gemini-3-pro-preview"
    url = f"https://generativelanguage.googleapis.com/v1beta/models/{model}:generateContent"

    try:
        with httpx.Client(timeout=120.0) as client:
            r = client.post(url, params={"key": api_key}, json=payload,
                            headers={"Content-Type": "application/json"})
            if r.status_code == 404:
                model = "gemini-2.5-pro"
                url = f"https://generativelanguage.googleapis.com/v1beta/models/{model}:generateContent"
                r = client.post(url, params={"key": api_key}, json=payload,
                                headers={"Content-Type": "application/json"})
            r.raise_for_status()
            data = r.json()
    except Exception as e:
        return f"오류: VLM 호출 실패: {e}"

    parts = (data.get("candidates", [{}])[0].get("content", {}) or {}).get("parts", [])
    text = "".join(p.get("text", "") for p in parts).strip()
    if not text:
        return "오류: VLM 빈 응답 (이미지를 읽지 못했습니다)."
    return text


def generate_gemini_image(tool_input, output_base):
    """Gemini API를 사용하여 이미지를 생성합니다."""
    import httpx
    import base64

    prompt = tool_input.get("prompt")
    if not prompt:
        return "오류: prompt는 필수입니다."

    api_key = tool_input.get("api_key") or os.environ.get("GEMINI_API_KEY")
    if not api_key:
        return "오류: GEMINI_API_KEY 환경변수가 설정되지 않았거나 api_key 파라미터가 필요합니다."

    output_path = tool_input.get("output_path")
    aspect_ratio = tool_input.get("aspect_ratio", "1:1")
    image_size = tool_input.get("image_size", "1K")  # 512/1K/2K/4K (3.x), 2.5는 무시
    style_preset = tool_input.get("style_preset")
    final_prompt = _build_image_prompt(prompt, style_preset)

    # 모델 선택 — 기본은 Nano Banana 2 (Gemini 3.1 Flash, 2026-02 출시)
    # quality 별칭으로 간편 선택, 또는 model로 직접 지정 가능
    quality = tool_input.get("quality")  # "fast" | "pro" | "legacy"
    quality_map = {
        "fast": "gemini-3.1-flash-image-preview",  # Nano Banana 2 (기본)
        "pro": "gemini-3-pro-image-preview",       # Nano Banana Pro — 4K, 더 정밀
        "legacy": "gemini-2.5-flash-image",        # 구버전 폴백
    }
    model = tool_input.get("model") or quality_map.get(quality, "gemini-3.1-flash-image-preview")

    # output_path가 지정되면 파일명만 추출하여 output_base에 저장
    if output_path:
        filename = os.path.basename(output_path)
        output_path = os.path.join(output_base, filename)
    else:
        output_path = os.path.join(output_base, f"gemini_image_{uuid.uuid4().hex[:8]}.png")
    os.makedirs(os.path.dirname(os.path.abspath(output_path)), exist_ok=True)

    url = f"https://generativelanguage.googleapis.com/v1beta/models/{model}:generateContent"

    # 페이로드 구조: 양쪽 모델 모두 generationConfig.imageConfig 사용.
    #   - 2.5: aspectRatio만 지원
    #   - 3.x: aspectRatio + imageSize (1K/2K/4K/512) 지원
    is_legacy = model.startswith("gemini-2.5")
    image_config = {"aspectRatio": aspect_ratio}
    if not is_legacy:
        image_config["imageSize"] = image_size
    payload = {
        "contents": [{"parts": [{"text": final_prompt}]}],
        "generationConfig": {
            "responseModalities": ["IMAGE", "TEXT"],
            "imageConfig": image_config,
        },
    }

    try:
        with httpx.Client(timeout=180.0) as client:
            response = client.post(
                url,
                params={"key": api_key},
                json=payload,
                headers={"Content-Type": "application/json"}
            )
            response.raise_for_status()
            data = response.json()

        # 응답에서 이미지 데이터 추출
        candidates = data.get("candidates", [])
        if not candidates:
            return f"오류: Gemini API 응답에 결과가 없습니다. 응답: {data}"

        parts = candidates[0].get("content", {}).get("parts", [])
        image_saved = False
        description = ""

        for part in parts:
            if "inlineData" in part:
                img_data = base64.b64decode(part["inlineData"]["data"])
                with open(output_path, "wb") as f:
                    f.write(img_data)
                image_saved = True
            elif "text" in part:
                description = part["text"]

        if not image_saved:
            return f"오류: 응답에 이미지 데이터가 없습니다. 텍스트 응답: {description or data}"

        size_note = "" if is_legacy else f" / {image_size}"
        result = (
            f"Gemini 이미지 생성 완료: {os.path.abspath(output_path)}\n"
            f"모델: {model}{size_note} (aspect {aspect_ratio})\n"
            f"프롬프트: {prompt}"
        )
        if description:
            result += f"\n설명: {description}"
        return result

    except httpx.HTTPStatusError as e:
        return f"Gemini API 오류 ({e.response.status_code}): {e.response.text}"
    except Exception as e:
        return f"Gemini 이미지 생성 중 오류 발생: {str(e)}"


def create_tts(tool_input, output_base):
    """텍스트를 음성 파일(MP3)로 변환합니다. (Edge TTS)"""
    text = tool_input.get("text")
    if not text:
        return "오류: text는 필수입니다."

    voice = tool_input.get("voice", "ko-KR-SunHiNeural")
    rate = tool_input.get("rate", "+0%")
    pitch = tool_input.get("pitch", "+0Hz")
    output_filename = tool_input.get("output_filename")

    if output_filename:
        output_path = os.path.join(output_base, os.path.basename(output_filename))
    else:
        output_path = os.path.join(output_base, f"tts_{uuid.uuid4().hex[:8]}.mp3")

    os.makedirs(os.path.dirname(os.path.abspath(output_path)), exist_ok=True)

    try:
        asyncio.run(generate_tts(text, output_path, voice, rate, pitch))

        # 길이 측정
        from moviepy import AudioFileClip as MpAudioClip
        clip = MpAudioClip(output_path)
        duration = clip.duration
        clip.close()

        abs_path = os.path.abspath(output_path)
        return f"TTS 생성 완료: {abs_path}\n길이: {duration:.1f}초\n음성: {voice}"
    except Exception as e:
        return f"TTS 생성 중 오류 발생: {str(e)}"


def render_html_video(tool_input, output_base):
    """HTML 파일 경로 목록을 받아서 MP4 동영상으로 렌더링합니다.

    코드 작성은 에이전트가 system:write로 미리 해두고,
    이 함수는 순수 렌더링만 담당합니다.

    입력:
      scene_files: [{path: "scene1.html", duration: 5}, ...]
      또는
      scene_dir: "씬 HTML 파일들이 있는 폴더" (알파벳 순 정렬)

      narration_files: ["narr1.mp3", "narr2.mp3", ...]  (선택)
      bgm_path: "bgm.mp3"  (선택)
    """
    import subprocess
    import shutil
    import glob as glob_module
    from playwright.sync_api import sync_playwright

    scene_files = tool_input.get("scene_files", [])
    scene_dir = tool_input.get("scene_dir")
    narration_files = tool_input.get("narration_files", [])
    bgm_path = tool_input.get("bgm_path")
    output_filename = tool_input.get("output_filename", "rendered_video.mp4")
    video_width = tool_input.get("width", 1280)
    video_height = tool_input.get("height", 720)
    fps = tool_input.get("fps", 24)
    default_duration = tool_input.get("duration_per_scene", 5)

    # scene_dir이 지정되면 폴더에서 HTML 파일들을 자동 수집
    if scene_dir and not scene_files:
        scene_dir = os.path.abspath(scene_dir)
        if not os.path.isdir(scene_dir):
            return f"오류: scene_dir '{scene_dir}'이 존재하지 않습니다."
        html_paths = sorted(glob_module.glob(os.path.join(scene_dir, "*.html")))
        if not html_paths:
            return f"오류: '{scene_dir}'에 HTML 파일이 없습니다."
        scene_files = [{"path": p, "duration": default_duration} for p in html_paths]

    if not scene_files:
        return "오류: scene_files 또는 scene_dir이 필요합니다."

    # 경로 검증
    for i, sf in enumerate(scene_files):
        path = sf.get("path", "")
        if not os.path.isabs(path):
            path = os.path.join(output_base, path)
            sf["path"] = path
        if not os.path.exists(path):
            return f"오류: scene_files[{i}]의 파일이 없습니다: {path}"

    output_path = os.path.join(output_base, output_filename)
    temp_dir = os.path.join(output_base, f"temp_render_{uuid.uuid4().hex[:8]}")
    frames_dir = os.path.join(temp_dir, "frames")
    os.makedirs(frames_dir, exist_ok=True)

    try:
        # 프레임 캡처
        global_frame = 0

        with sync_playwright() as p:
            browser = p.chromium.launch()
            page = browser.new_page(viewport={"width": video_width, "height": video_height})

            for i, sf in enumerate(scene_files):
                html_path = sf["path"]
                duration = sf.get("duration", default_duration)

                # HTML 파일을 읽어서 전처리 (이미지 Base64 변환 등)
                with open(html_path, "r", encoding="utf-8") as f:
                    html = f.read()

                base_path = os.path.dirname(html_path)
                html_ready = _prepare_scene_html(html, base_path, video_width, video_height)

                # 전처리된 HTML을 임시 파일로 저장 후 로드
                temp_html = os.path.join(temp_dir, f"_render_scene_{i}.html")
                with open(temp_html, "w", encoding="utf-8") as f:
                    f.write(html_ready)
                page.goto(f"file://{temp_html}")

                # 애니메이션 일시정지 → 리소스 로드 → 재개
                page.evaluate("""() => {
                    const style = document.createElement('style');
                    style.id = '__anim_pause__';
                    style.textContent = '*, *::before, *::after { animation-play-state: paused !important; }';
                    document.head.appendChild(style);
                }""")
                page.wait_for_load_state("networkidle")
                page.wait_for_timeout(1500 if i == 0 else 500)
                page.evaluate("""() => {
                    const s = document.getElementById('__anim_pause__');
                    if (s) s.remove();
                }""")

                total_frames = int(duration * fps)
                frame_interval_ms = 1000.0 / fps

                for f_idx in range(total_frames):
                    frame_path = os.path.join(frames_dir, f"frame_{global_frame:06d}.png")
                    page.screenshot(path=frame_path)
                    global_frame += 1
                    page.wait_for_timeout(int(frame_interval_ms))

            browser.close()

        if global_frame == 0:
            shutil.rmtree(temp_dir)
            return "오류: 캡처된 프레임이 없습니다."

        # FFmpeg 인코딩
        merged_video = os.path.join(temp_dir, "merged.mp4")
        ffmpeg_result = subprocess.run([
            "ffmpeg", "-y",
            "-framerate", str(fps),
            "-i", os.path.join(frames_dir, "frame_%06d.png"),
            "-c:v", "libx264", "-preset", "fast", "-crf", "20",
            "-pix_fmt", "yuv420p",
            merged_video
        ], capture_output=True)

        if ffmpeg_result.returncode != 0:
            shutil.rmtree(temp_dir)
            return f"FFmpeg 인코딩 오류: {ffmpeg_result.stderr.decode()}"

        # 오디오 합성 (나레이션 + BGM)
        if bgm_path:
            bgm_path = os.path.abspath(bgm_path)
        has_bgm = bgm_path and os.path.exists(bgm_path)
        has_narration = any(narration_files)

        if has_narration or has_bgm:
            from moviepy import AudioFileClip as MpAudioClip, CompositeAudioClip as MpCompositeAudioClip, VideoFileClip

            video_clip = VideoFileClip(merged_video)
            audio_layers = []

            if has_narration:
                scene_start = 0.0
                for i, sf in enumerate(scene_files):
                    dur = sf.get("duration", default_duration)
                    if i < len(narration_files) and narration_files[i]:
                        narr_path = narration_files[i]
                        if not os.path.isabs(narr_path):
                            narr_path = os.path.join(output_base, narr_path)
                        if os.path.exists(narr_path):
                            narr = MpAudioClip(narr_path)
                            narr = narr.with_start(scene_start)
                            audio_layers.append(narr)
                    scene_start += dur

            if has_bgm:
                bgm_clip = MpAudioClip(bgm_path).with_duration(video_clip.duration)
                if has_narration:
                    bgm_clip = bgm_clip.multiply_volume(0.2)
                else:
                    bgm_clip = bgm_clip.multiply_volume(0.5)
                audio_layers.append(bgm_clip)

            if audio_layers:
                final_audio = MpCompositeAudioClip(audio_layers)
                video_clip = video_clip.with_audio(final_audio)

            video_clip.write_videofile(output_path, fps=fps, codec="libx264", audio_codec="aac")
            video_clip.close()
        else:
            shutil.copy2(merged_video, output_path)

        shutil.rmtree(temp_dir)
        abs_output = os.path.abspath(output_path)
        scene_count = len(scene_files)
        total_dur = sum(sf.get("duration", default_duration) for sf in scene_files)
        return f"HTML 동영상 렌더링 완료: {abs_output}\n씬 수: {scene_count}개, 총 길이: {total_dur:.1f}초"

    except Exception as e:
        if os.path.exists(temp_dir):
            shutil.rmtree(temp_dir)
        return f"렌더링 중 오류 발생: {str(e)}"
