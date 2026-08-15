"""
shadcn_slides.py
shadcn/ui 컴포넌트와 web-builder 테마를 활용한 슬라이드 생성
"""

import os
import json
import uuid
import base64
import urllib.request
from jinja2 import Template

# 상수 모듈 분리(2026-08-06, 1500줄 규칙): 디자인 정체성=shadcn_design.py / 템플릿=shadcn_layouts.py.
# ★이 모듈 자체가 spec_from_file_location 으로 로드되므로(lecture_workspace/slide_ai.py)
# 상대 import 불가 — 하우스 패턴(sys.modules 미등록 spec-load: 싱글턴 임포트 레이스 없음).
import importlib.util as _ilu

def _load_sibling(fname: str, modname: str):
    spec = _ilu.spec_from_file_location(modname, os.path.join(os.path.dirname(__file__), fname))
    mod = _ilu.module_from_spec(spec)
    spec.loader.exec_module(mod)
    return mod

_design = _load_sibling("shadcn_design.py", "mp_shadcn_design")
_layouts = _load_sibling("shadcn_layouts.py", "mp_shadcn_layouts")
THEMES = _design.THEMES
DESIGN_SYSTEMS = _design.DESIGN_SYSTEMS
SLIDE_BASE_TEMPLATE = _layouts.SLIDE_BASE_TEMPLATE
SLIDE_LAYOUTS = _layouts.SLIDE_LAYOUTS


def get_image_base64(image_path: str) -> str:
    """이미지를 Base64로 변환"""
    if not image_path or not os.path.exists(image_path):
        return None
    try:
        with open(image_path, "rb") as f:
            encoded = base64.b64encode(f.read()).decode('utf-8')
            ext = os.path.splitext(image_path)[1].lower()
            mime = {
                ".jpg": "image/jpeg",
                ".jpeg": "image/jpeg",
                ".png": "image/png",
                ".webp": "image/webp",
                ".gif": "image/gif"
            }.get(ext, "image/png")
            return f"data:{mime};base64,{encoded}"
    except Exception as e:
        print(f"이미지 로드 실패: {e}")
        return None


def render_slide(slide_data: dict, theme_name: str = "default", width: int = 1280, height: int = 720, design_system_name: str = "default") -> str:
    """슬라이드 HTML 생성

    Args:
        slide_data: 슬라이드 dict (layout + 레이아웃별 키)
        theme_name: 색 테마 (THEMES dict 키). design_system이 theme_override를 가지면 무시됨.
        design_system_name: 디자인 시스템 (DESIGN_SYSTEMS dict 키). 색+폰트+텍스처+장식 일관 묶음.
    """
    from jinja2 import Environment, BaseLoader, Undefined

    # Undefined 변수에 대해 빈 문자열 반환하는 환경 설정
    class SilentUndefined(Undefined):
        def _fail_with_undefined_error(self, *args, **kwargs):
            return ''
        def __str__(self):
            return ''
        def __iter__(self):
            return iter([])
        def __bool__(self):
            return False
        def __getitem__(self, key):
            return ''
        def __getattr__(self, name):
            return SilentUndefined()

    env = Environment(loader=BaseLoader(), undefined=SilentUndefined)

    # 디자인 시스템 가져오기 (먼저 — theme_override 적용 위해)
    ds = DESIGN_SYSTEMS.get(design_system_name, DESIGN_SYSTEMS["default"])

    # 테마 결정: 디자인 시스템이 override하면 그쪽, 아니면 외부 인자
    effective_theme_name = ds.get("theme_override") or theme_name
    theme = THEMES.get(effective_theme_name, THEMES["default"])

    # 레이아웃 타입
    layout_type = slide_data.get("layout", "hero")
    layout_template = SLIDE_LAYOUTS.get(layout_type, SLIDE_LAYOUTS["hero"])

    # 이미지 처리
    if slide_data.get("image_path"):
        slide_data["image_data"] = get_image_base64(slide_data["image_path"])
    if slide_data.get("avatar_path"):
        slide_data["avatar_data"] = get_image_base64(slide_data["avatar_path"])
    if slide_data.get("left_image_path"):
        slide_data["left_image_data"] = get_image_base64(slide_data["left_image_path"])
    if slide_data.get("right_image_path"):
        slide_data["right_image_data"] = get_image_base64(slide_data["right_image_path"])

    # 레이아웃 렌더링 (SilentUndefined 사용)
    layout_tpl = env.from_string(layout_template)
    content_html = layout_tpl.render(**slide_data)

    # style_overrides — spec에 사용자가 지정한 미세 조정값 (없으면 빈 CSS)
    style_overrides_css = _build_style_overrides_css(slide_data.get("style_overrides"))

    # 베이스 템플릿에 삽입 (디자인 시스템 주입 포함)
    base_tpl = env.from_string(SLIDE_BASE_TEMPLATE)
    full_html = base_tpl.render(
        theme=theme,
        width=width,
        height=height,
        content=content_html,
        design_system_head=ds.get("extra_head", ""),
        design_system_css=ds.get("extra_css", ""),
        design_system_html=ds.get("extra_html", ""),
        style_overrides_css=style_overrides_css,
    )

    return full_html


def _build_style_overrides_css(overrides) -> str:
    """spec.style_overrides → 슬라이드별 추가 CSS.

    지원 키:
      - font_scale: float (0.7 ~ 1.4). html root font-size 조정 → 모든 rem 단위 비례.
      - text_align: 'left' | 'center' | 'right'. 본문 텍스트 강제 정렬.
      - accent_color: hex 문자열 (예: '#a55a3e'). primary 색 override.

    1280×720 viewport에서 font_scale은 텍스트만 비례 확대/축소 (이미지·여백은 그대로).
    너무 큰 값은 슬라이드가 잘릴 수 있음 — 0.85~1.25 권장.
    """
    if not overrides or not isinstance(overrides, dict):
        return ""

    parts = []

    # font_scale — html root font-size 조정 (tailwind rem 단위 비례)
    fs = overrides.get("font_scale")
    if isinstance(fs, (int, float)) and 0.5 <= float(fs) <= 2.0 and float(fs) != 1.0:
        parts.append(f"html {{ font-size: {16 * float(fs):.2f}px !important; }}")

    # text_align — 본문 텍스트 강제 정렬
    ta = overrides.get("text_align")
    if ta in ("left", "center", "right"):
        parts.append(
            f".slide-container h1, .slide-container h2, .slide-container h3, "
            f".slide-container p, .slide-container li, .slide-container blockquote "
            f"{{ text-align: {ta} !important; }}"
        )

    # accent_color — primary 색 override (text/bg/border-primary 클래스)
    ac = overrides.get("accent_color")
    if isinstance(ac, str) and ac.startswith("#") and len(ac) in (4, 7):
        parts.append(
            f".slide-container [class*='text-primary'] {{ color: {ac} !important; }}\n"
            f".slide-container [class*='bg-primary'] {{ background-color: {ac} !important; }}\n"
            f".slide-container [class*='border-primary'] {{ border-color: {ac} !important; }}"
        )

    return "\n".join(parts)


_LEGACY_SLIDE_THEMES = {
    "modern", "tech", "business", "title_bold", "dark_tech",
    "glassmorphism", "gradient_modern", "split_asymmetric",
    "minimal_white", "image_fullscreen", "data_card", "tailwind",
}


def _adapt_legacy_slide_input(tool_input: dict) -> dict:
    """옛 slide 스키마 입력을 [engines:slide_shadcn] 입력으로 호환 변환.

    옛 slide dict: {title, body, theme(modern/tech/...), image_path, bg_color, ...}
    새 slide dict: {layout, title, body, ...}

    layout이 없는 슬라이드에 한해 image_path/body 유무 기반으로 기본 layout 주입.
    옛 theme(slide_shadcn의 theme와 의미 다름)과 색상 키들은 제거.
    """
    adapted = dict(tool_input)
    slides = adapted.get("slides")
    if not isinstance(slides, list):
        return adapted
    new_slides = []
    for s in slides:
        if not isinstance(s, dict):
            new_slides.append(s)
            continue
        s = dict(s)
        if "layout" not in s:
            if s.get("image_path") or s.get("image_data"):
                s["layout"] = "content_image"
            elif s.get("body"):
                s["layout"] = "lecture_body"
            else:
                s["layout"] = "hero"
        for legacy_key in ("theme", "bg_color", "text_color", "accent_color"):
            s.pop(legacy_key, None)
        new_slides.append(s)
    adapted["slides"] = new_slides
    if adapted.get("theme") in _LEGACY_SLIDE_THEMES:
        adapted.pop("theme")
    return adapted


def _bundle_slides(png_paths: list, output_dir: str, fmt: str, width: int, height: int) -> str:
    """렌더된 슬라이드 PNG들을 단일 공유 파일로 묶음 — 디자인 보존(이미지 그대로).
    pdf=슬라이드당 1페이지 / pptx=슬라이드당 풀블리드 이미지. 반환: 산출 파일 절대경로."""
    if fmt == "pdf":
        from PIL import Image
        imgs = [Image.open(p).convert("RGB") for p in png_paths]
        out = os.path.join(output_dir, "slides.pdf")
        imgs[0].save(out, save_all=True, append_images=imgs[1:])
        return out
    if fmt == "pptx":
        from pptx import Presentation
        from pptx.util import Emu
        prs = Presentation()
        prs.slide_width = Emu(int(width / 96 * 914400))   # px→EMU (96dpi 가정, 슬라이드 비율 유지)
        prs.slide_height = Emu(int(height / 96 * 914400))
        blank = prs.slide_layouts[6]
        for p in png_paths:
            s = prs.slides.add_slide(blank)
            s.shapes.add_picture(p, 0, 0, width=prs.slide_width, height=prs.slide_height)
        out = os.path.join(output_dir, "slides.pptx")
        prs.save(out)
        return out
    raise ValueError(f"지원하지 않는 format: {fmt}")


def create_shadcn_slides(tool_input: dict, output_base: str) -> str:
    """
    shadcn 스타일 슬라이드 생성"""
    tool_input = _adapt_legacy_slide_input(tool_input)
    return _create_shadcn_slides_impl(tool_input, output_base)


def _create_shadcn_slides_impl(tool_input: dict, output_base: str) -> str:
    """
    shadcn 스타일 슬라이드 생성 (실제 구현)

    Args:
        tool_input: {
            "slides": [
                {
                    # layout 옵션:
                    #   마케팅: hero / hero_image / features / stats / testimonial / pricing / cta / content_image / steps / custom
                    #   강의:   lecture_body / metaphor_story / comparison_table / factbox / quote
                    #   통합형 일러스트(NotebookLM 스타일):
                    #           hero_illustration / illustration_anchor / split_concept /
                    #           illustration_background / comparison_iconic
                    "layout": "hero",
                    "title": "제목",
                    "subtitle": "부제목",
                    # 일러스트가 필요한 layout에 image_path 절대경로(자동 base64 변환):
                    #   - image_path  → image_data  (hero_illustration / illustration_anchor / illustration_background / hero_image / content_image)
                    #   - left_image_path  → left_image_data  (split_concept)
                    #   - right_image_path → right_image_data (split_concept)
                    #   - avatar_path → avatar_data           (testimonial)
                    ...
                }
            ],
            "theme": "blue",          # default / blue / green / purple / orange / dark
            "design_system": "vintage_book",  # default / vintage_book / blueprint / architect / ink_orange
            "output_dir": "경로",     # 선택. 미지정 시 output_base/shadcn_slides_<8자hex>/
            "width": 1280,            # 선택. 슬라이드 가로 픽셀 (기본 1280)
            "height": 720             # 선택. 슬라이드 세로 픽셀 (기본 720)
        }

    Returns:
        JSON: {success, message, output_dir, images[], html_files[], theme}
    """
    slides_data = tool_input.get("slides", [])

    # 가드: 빈 slides + 옛 호출 패턴 감지
    if not slides_data:
        legacy_keys = [k for k in ("topic", "file") if k in tool_input]
        if legacy_keys:
            return (
                "오류: [engines:slide_shadcn]는 'slides' 배열 인라인 호출만 받습니다. "
                f"'{', '.join(legacy_keys)}'는 지원하지 않습니다. "
                "올바른 호출: [engines:slide_shadcn]{slides: [{layout: \"hero\", title: \"제목\", subtitle: \"부제\"}, ...]}. "
                "레이아웃 목록은 read_guide(query=\"슬라이드\")로 가이드 확인."
            )
        return (
            "오류: [engines:slide_shadcn] 호출에 slides 배열이 비어있습니다. "
            "[engines:slide_shadcn]{slides: [{layout, title, ...}, ...]} 형태로 인라인 배열을 전달하세요."
        )

    theme_name = tool_input.get("theme", "default")
    design_system_name = tool_input.get("design_system", "default")
    if design_system_name not in DESIGN_SYSTEMS:
        return (
            f"오류: 알 수 없는 design_system: '{design_system_name}'. "
            f"사용 가능: {', '.join(DESIGN_SYSTEMS.keys())}"
        )
    custom_output_dir = tool_input.get("output_dir")
    width = tool_input.get("width", 1280)
    height = tool_input.get("height", 720)

    output_dir = custom_output_dir if custom_output_dir else os.path.join(output_base, f"shadcn_slides_{uuid.uuid4().hex[:8]}")
    output_dir = os.path.abspath(output_dir)  # 외부와 소통하는 모든 경로는 절대경로 (AI 안내가 cwd에 의존하지 않도록)
    os.makedirs(output_dir, exist_ok=True)

    generated_paths = []
    html_paths = []

    try:
        from playwright.sync_api import sync_playwright

        with sync_playwright() as p:
            browser = p.chromium.launch()
            page = browser.new_page(viewport={"width": width, "height": height})

            for i, slide in enumerate(slides_data):
                # HTML 생성
                html_content = render_slide(slide, theme_name, width, height, design_system_name)

                # HTML 파일 저장 (디버깅용)
                html_path = os.path.join(output_dir, f"slide_{i+1:02d}.html")
                with open(html_path, "w", encoding="utf-8") as f:
                    f.write(html_content)
                html_paths.append(html_path)

                # 렌더링
                page.set_content(html_content)
                page.wait_for_load_state("networkidle")
                page.wait_for_timeout(800)  # Tailwind/폰트/아이콘 로딩 대기

                # 스크린샷
                png_path = os.path.join(output_dir, f"slide_{i+1:02d}.png")
                page.screenshot(path=png_path)
                generated_paths.append(png_path)

            browser.close()

        # emitter 정리: 슬라이드 IR(slides[]) → png(기본, 슬라이드별 이미지) / pdf / pptx.
        # ★슬라이드 layout은 디자인된 HTML이라, pdf·pptx는 렌더된 PNG를 그대로 보존(네이티브 도형 재구성=디자인 파괴).
        fmt = (tool_input.get("format") or "png").strip().lower()
        result = {
            "success": True,
            "message": f"{len(generated_paths)}개의 슬라이드가 생성되었습니다",
            "output_dir": output_dir,
            "images": generated_paths,
            "html_files": html_paths,
            "theme": theme_name,
            "format": "png",
        }
        if fmt in ("pdf", "pptx") and generated_paths:
            try:
                bundle = _bundle_slides(generated_paths, output_dir, fmt, width, height)
                result["format"] = fmt
                result["path"] = bundle
                result["file"] = bundle
                result["message"] = f"{len(generated_paths)}개 슬라이드를 {fmt.upper()}로 묶었습니다."
            except Exception as e:
                result["message"] += f" (단, {fmt} 묶기 실패 → PNG 유지: {e})"
        return json.dumps(result, ensure_ascii=False)

    except ImportError:
        return json.dumps({
            "success": False,
            "error": "Playwright가 설치되어 있지 않습니다. 'pip install playwright' 후 "
                     "'python scripts/check_playwright_browsers.py --install' 실행 필요 "
                     "(맨손 playwright install 은 기본 캐시로 받아 백엔드가 보는 주소와 어긋난다)"
        }, ensure_ascii=False)
    except Exception as e:
        return json.dumps({
            "success": False,
            "error": str(e)
        }, ensure_ascii=False)


if __name__ == "__main__":
    # 테스트
    test_input = {
        "theme": "blue",
        "slides": [
            {
                "layout": "hero",
                "badge": "New Release",
                "title": "IndieBiz OS",
                "subtitle": "AI 기반 통합 비즈니스 관리 시스템",
                "cta_text": "시작하기"
            },
            {
                "layout": "features",
                "title": "주요 기능",
                "features": [
                    {"icon": "🤖", "title": "AI 에이전트", "description": "맞춤형 AI 비서가 업무를 도와줍니다"},
                    {"icon": "📊", "title": "데이터 분석", "description": "실시간 인사이트를 제공합니다"},
                    {"icon": "🔗", "title": "통합 연동", "description": "다양한 서비스와 연결됩니다"}
                ]
            },
            {
                "layout": "stats",
                "title": "성과",
                "stats": [
                    {"value": "10K+", "label": "활성 사용자"},
                    {"value": "99.9%", "label": "가동률"},
                    {"value": "24/7", "label": "지원"},
                    {"value": "50+", "label": "통합 서비스"}
                ]
            },
            {
                "layout": "cta",
                "title": "지금 시작하세요",
                "subtitle": "무료로 체험해보고 비즈니스를 성장시키세요",
                "cta_text": "무료 체험 시작"
            }
        ]
    }

    result = create_shadcn_slides(test_input, "/tmp/test_slides")
    print(result)

